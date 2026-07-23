# -*- coding: utf-8 -*-
"""完成版:タイトル→企画内容→見積(木曜SB打ち合わせ提出用・円谷フォーマット)"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=11.5, row_h=0.5):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(headers), x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, font_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            for k, line in enumerate(str(text).split("\n")):
                p = c.text_frame.paragraphs[0] if k == 0 else c.text_frame.add_paragraph()
                r = p.add_run(); r.text = line
                set_font(r, font_size, color, bold)
    return tbl



COVER_IMG = os.path.join(ASSETS, "band_cover.png")
TSUB_NAVY = RGBColor(0x1F, 0x33, 0x8C)

# =============================================================================
# 1. 表紙
# =============================================================================
s = add_slide()
s.shapes.add_picture(COVER_IMG, 0, 0, width=Inches(4.09), height=SH)
textbox(s, Inches(4.5), Inches(1.25), Inches(8.4), Inches(0.5),
        [("ソフトバンク株式会社 御中", 14, SUB, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(4.5), Inches(2.25), Inches(8.4), Inches(2.0), [
    ("ウルトラマンIP ご活用のご提案", 30, INK, True),
    ("企画内容とお見積り", 26, INK, True),
], align=PP_ALIGN.CENTER, gap=10)
textbox(s, Inches(4.5), Inches(4.2), Inches(8.4), Inches(0.9), [
    ("店頭イベント・イベントキット・店頭訴求ツールへの", 13.5, INK, False),
    ("ウルトラマンIP導入(年間ライセンス)と今期スポットテストのご提案", 13.5, INK, False),
], align=PP_ALIGN.CENTER, gap=4)
textbox(s, Inches(4.5), Inches(5.7), Inches(8.4), Inches(1.0), [
    ("2026年7月", 15, TSUB_NAVY, True),
    ("(株)円谷プロダクション", 15, TSUB_NAVY, True),
    ("金額はすべて税別・概算", 10.5, SILVER, False),
], align=PP_ALIGN.CENTER, gap=3)
textbox(s, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
        [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(10.25), Inches(7.1), Inches(2.5), Inches(0.3),
        [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)

# =============================================================================
# 2. 企画内容:ご活用の全体像(3階建て)
# =============================================================================
s = add_slide()
header(s, "企画内容:ウルトラマンIPご活用の全体像",
       "貴社の現行のご運用(イベント+イベント外)と同じ枠組み・同じ運用フローのまま、ウルトラマンをご活用いただけます")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "内容", "費用(税別)"],
          [
              [("① 年間IP使用料(固定)", True), "イベント・イベントキット・店頭訴求ツール・LP/Web・SNSでのIP利用と作成物の監修対応を包括。着ぐるみイベントの開催権利を含みます", ("初年度特別 2,000万円/年\n通常 2,400万円/年", True, RED)],
              [("② 着ぐるみイベント(従量)", True), "店頭写真撮影会・グリーティング。ディレクター+スーツアクター2名の専属体制・演技対応込み。実施した開催数の分だけのお支払いです", ("20万円/開催(土日2日間)", True, INK)],
              [("③ 作成許諾ロイヤリティ(従量)", True), "イベントキット・ノベルティ・店頭訴求ツールが対象(作成主体は貴社)。作成実績(実額)ベースで精算します", ("作成費用の20%", True, INK)],
          ],
          col_widths=[1.3, 2.6, 1.1], row_h=0.78, font_size=12)
b = rect(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.5), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(5.03), Inches(11.3), Inches(1.3), [
    ("この構造のポイント", 13, RED, True),
    ("・固定でお支払いいただくのは①のみ。②③は実施・作成した分だけの完全従量型で、規模のご判断に費用が連動します", 12, INK, False),
    ("・着ぐるみは1回30分×1日最大5回、1回あたり約50組にご対応可能。キャラクター登場時のBGM音素材もご用意します", 12, INK, False),
], gap=5)
textbox(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.55), [
    ("※着ぐるみ運搬・交通・宿泊費、スタッフ昼食費等の実費を別途申し受けます。※控室(非公開の着替え場所)のご用意など、実施ルールの詳細は別紙レギュレーションをご参照ください。", 10, SUB, False),
], gap=3)

# =============================================================================
# 3. 企画内容:ウルトラマンならではの価値
# =============================================================================
s = add_slide()
header(s, "企画内容:ウルトラマンならではの価値",
       "シリーズ60周年。ほかのIPにはない4つの強みで、貴社店頭の集客を年間を通じて支えます")

def value_block(x, y, title, lines):
    bb = rect(s, x, y, Inches(5.75), Inches(2.1), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bb.adjustments[0] = 0.07
    textbox(s, x + Inches(0.28), y + Inches(0.16), Inches(5.2), Inches(1.8),
            [(title, 13.5, RED, True)] + [(t, 11.5, INK, False) for t in lines], gap=4)

value_block(Inches(0.7), Inches(1.6), "1. 選べるヒーロー(8キャラクター)", [
    "昭和・平成・令和の8ヒーローから店舗・企画ごとに選択可能",
    "(初代ウルトラマン/タロウ/ティガ/ダイナ/ゼロ/アーク/オメガ/テオ)",
    "季節・エリアごとの「キャラクター切替」で来店動機が持続します",
])
value_block(Inches(6.85), Inches(1.6), "2. 繁忙期の夏休み商戦も全力稼働", [
    "7・8月を含むフル12か月の稼働が可能です",
    "夏商戦・夏休みの店頭集客に切れ目がありません",
    "フル稼働なら年間最大720開催まで拡張できます",
])
value_block(Inches(0.7), Inches(3.9), "3. 映画・新TVシリーズ・新商品との連動", [
    "劇場映画・新TVシリーズ・玩具新商品と連動した季節企画をご提供",
    "キット・ノベルティのデザイン切替で「変化」を作り、",
    "リピート来店・継続的な話題化につなげます",
])
value_block(Inches(6.85), Inches(3.9), "4. 高いイベント運用品質", [
    "専属ディレクター+スーツアクター2名体制・演技対応込み",
    "1回30分×1日最大5回・1回約50組",
])
textbox(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.6), [
    ("※年間の最低開催数・最低作成数をお約束いただける場合は、着ぐるみ単価の優遇(例:400開催以上のコミットで17万円/開催)をご用意しています(次ページ)。", 10.5, SUB, False),
], gap=3)

# =============================================================================
# 4. 見積:規模別お見積り3パターン
# =============================================================================
s = add_slide()
header(s, "規模別お見積り:3パターン(年間・税別概算)",
       "7/17お打ち合わせを受け、全量・半分・3IP分割(約1/3)の3規模でご提示します(単価は共通:IP使用料は初年度特別2,000万円(通常2,400万円)/着ぐるみ20万円/ロイヤリティ20%)")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["パターン", "着ぐるみ", "キット・ノベルティ想定", "①IP使用料(固定)", "②着ぐるみ費用", "③ロイヤリティ", "年間合計"],
          [
              [("最大(貴社実績同等・全量)", True), "600開催", "キット800式・ノベルティ2.1万セット(実績)・ツール全店2回", "2,000万円", "1億2,000万円", "約6,480万円", ("約2.05億円", True, RED)],
              [("半分(実績の50%)", True), "300開催", "キット400式・ノベルティ1.05万セット・ツール全店1回", "2,000万円", "6,000万円", "約3,240万円", ("約1.12億円", True, RED)],
              [("3IP分割(約1/3配分)", True), "200開催", "キット270式・ノベルティ7,000セット・ツール1/3", "2,000万円", "4,000万円", "約2,170万円", ("約8,200万円", True, RED)],
          ],
          col_widths=[1.15, 0.65, 1.85, 0.95, 0.95, 0.85, 0.95], row_h=0.72, font_size=10.5)
b = rect(s, Inches(0.7), Inches(4.65), Inches(11.9), Inches(1.45), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.25), [
    ("お打ち合わせを受けたご提案", 13, RED, True),
    ("・年間の最低開催数・最低作成数をお約束いただける場合、着ぐるみ単価を優遇します(例:400開催以上のコミットで17万円/開催)", 12, INK, False),
    ("・ウルトラマンは7・8月も稼働可能です(現行IP様は夏季稼働不可)。夏商戦を含むフル12か月稼働なら最大720開催まで拡張できます", 12, INK, False),
    ("・複数IPでのご運用でも、映画・新TVシリーズ・玩具新商品と連動した季節企画で「切替の変化」をご提供できます(キャラクター切替・選べるヒーロー等)", 12, INK, False),
], gap=5)
textbox(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.8), [
    ("※③はキット単価15.5万円/式、ノベルティは貴社FY25ご実績(年約2.1万セット・作成費約1億8,500万円)を配分比で試算。ロイヤリティは作成実績(実額)ベースで精算します。", 10, SUB, False),
    ("※店頭訴求ツールは作成数実績(年間約23.5万部)をご共有いただきました。作成費(金額)のご共有後、③を最終精緻化します。※固定費は①のみ。②③は実施・作成分だけの従量です。", 10, SUB, False),
    ("※②は土日2日間・ディレクター+スーツアクター2名体制。着ぐるみ運搬・交通・宿泊費、スタッフ昼食費等の実費を別途申し受けます。", 10, SUB, False),
], gap=3)

# =============================================================================
# 5. 見積:今期スポットテストのご提案
# =============================================================================
s = add_slide()
header(s, "今期スポットテストのご提案(回数ベース・関東エリア)",
       "年間のお取り組みとは切り離した単発テストです。撮影、ステージショーの2形式をご用意しました")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "形式S1:写真撮影会型(グリーティング)", "形式S2:ステージショー型"],
          [
              ["内容", "1回30分×1日最大5回の撮影会・グリーティング。店頭・商業施設の集客と送客を検証", "ステージショー(2回/日)+撮影会のセット。大型集客の瞬発力を検証"],
              ["単位", "1開催=土日2日間", "1開催=土日2日間"],
              [("費用", True, RED), ("100万円/開催(IP利用込み・監修・スーツアクター・運搬込み)", True, INK), ("200万円/開催(ショー演出・音源・追加スタッフ込み)", True, INK)],
              ["貴社側ご負担", "会場手配・集客告知・当日運営(誘導スタッフ)", "同左+ステージ設営"],
              ["想定規模", ("5施設程度(各1開催)を想定 → 総額500万円", True, INK), "うち1〜2施設でのオプション実施"],
              ["地域", "関東エリア", "同左"],
          ],
          col_widths=[0.6, 1.7, 1.7], row_h=0.56, font_size=11)
textbox(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.3), [
    [("想定プラン:", 12.5, RED, True),
     ("エリア・施設タイプの異なる5施設で実施し、集客・送客効果を比較検証(S1×5施設=総額500万円・IP使用料込み)。", 12.5, INK, False)],
    [("継続パック(S1):", 12.5, RED, True),
     ("隔週(月2開催)190万円/月、毎週(月4開催)360万円/月(単発比の優遇設定)。", 12.5, INK, False)],
    [("実施条件:", 12.5, RED, True),
     ("当社の既存ライセンス条件の調整を経て、実施時期・会場を確定します(調整は当社にて速やかに進めます)。", 12, INK, False)],
    ("※通常イベントキットを用いるテストはキット製作を伴うため来年度からのご提案。", 10.5, SUB, False),
], gap=3)


# =============================================================================
# 6. ご留意事項:実施可否とキャンセル料
# =============================================================================
s = add_slide()
header(s, "ご留意事項:実施可否とキャンセル料",
       "イベントがお受けできない場合と、実施決定後のキャンセル料についてご確認ください")
b = rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.55), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2), [
    ("イベントがお受けできない場合について", 14, RED, True),
    ("・エリアバッティングの場合", 12.5, INK, True),
    ("   イベント日前8日間・後5日間、半径8キロ圏内に他イベントがある場合は、イベントの実施が不可となります。", 12, INK, False),
    ("   ※バッティング先が競合施設でない場合や、キャラクター違いの場合は、調整したのち実施できる場合がございます。", 11, SUB, False),
    ("・スーツアクターやスタッフが手配できない場合", 12.5, INK, True),
], gap=6)
b = rect(s, Inches(0.7), Inches(4.45), Inches(11.9), Inches(1.95), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(4.65), Inches(11.3), Inches(1.6), [
    ("キャンセル料について", 14, RED, True),
    ("イベント実施が決定した後のキャンセルは、キャンセル料が発生いたします。", 12, INK, False),
    ("・1か月前〜6日前:お見積り金額の30%   ・6日以内:お見積り金額の50%   ・前日:お見積り金額の100%", 12.5, INK, True),
], gap=6)

# =============================================================================
# 7. Appendix:お申し込みの流れ
# =============================================================================
s = add_slide()
header(s, "Appendix:お申し込みの流れ",
       "お申し込みからイベント開催までの4ステップです。イベント会期の2か月前を目安に申請書をご提出ください")
steps = [
    ("STEP 1  イベント開催申請書の提出", "イベント会期の2か月前が目安(申請内容の詳細は別紙をご参照ください)"),
    ("STEP 2  仮押さえ", "エリアバッティング確認OKの場合、エリア・キャラクター・スタッフを仮押さえで手配"),
    ("STEP 3  イベント決定・キャンセルの連絡", "イベント会期の1か月前までに、決定・キャンセルのご連絡をお願いします"),
    ("STEP 4  正式決定・開催", "実施に向け正式に手配。これ以降のキャンセルは規定のキャンセル料が発生します"),
]
yy = 1.6
for t, d in steps:
    bb = rect(s, Inches(0.7), Inches(yy), Inches(7.3), Inches(1.0), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bb.adjustments[0] = 0.12
    textbox(s, Inches(1.0), Inches(yy + 0.13), Inches(6.9), Inches(0.8),
            [(t, 13, NAVY, True), (d, 11, SUB, False)], gap=3)
    yy += 1.22
b = rect(s, Inches(8.3), Inches(1.6), Inches(4.3), Inches(4.86), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.05
textbox(s, Inches(8.6), Inches(1.85), Inches(3.8), Inches(4.4), [
    ("エリアバッティングとは", 12.5, RED, True),
    ("近隣でウルトラマンイベントの実施予定がある場合、出演NGとなります。", 11, INK, False),
    ("キャンセル料", 12.5, RED, True),
    ("決定以降のキャンセルはキャンセル料が発生いたします(料率はP6をご参照ください)。", 11, INK, False),
    ("イベント実施に向けて", 12.5, RED, True),
    ("・告知物の監修", 11, INK, False),
    ("・実施確認書の提出", 11, INK, False),
], gap=6)

out = "ウルトラマンIPご活用のご提案とお見積り.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
