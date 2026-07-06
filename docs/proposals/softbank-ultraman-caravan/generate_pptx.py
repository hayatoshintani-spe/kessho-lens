# -*- coding: utf-8 -*-
"""
SoftBank × ウルトラマン「未来防衛隊キャラバン」外部説明用プレゼン資料の生成スクリプト。
python-pptx で 16:9 の PPTX を出力する。

usage: python3 generate_pptx.py [output.pptx]
"""
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----------------------------------------------------------------
INK = RGBColor(0x21, 0x25, 0x2B)      # 本文ダークグレー
SUB = RGBColor(0x5A, 0x62, 0x6B)      # 補足グレー
RED = RGBColor(0xD6, 0x00, 0x12)      # アクセント(ヒーローレッド)
NAVY = RGBColor(0x12, 0x1A, 0x2E)     # 濃紺(表紙・中扉)
SILVER = RGBColor(0x8E, 0x94, 0x9B)   # シルバー
PALE = RGBColor(0xF2, 0xF4, 0xF6)     # 薄グレー背景
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB) # 薄赤背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD5, 0xD9, 0xDE)

FONT = "Yu Gothic"

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

page_no = [0]


def _set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    # 日本語フォントを明示
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6):
    """lines: list of (text, size, color, bold) or list of runs [(t,s,c,b),...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        runs = item if isinstance(item, list) else [item]
        for (t, s, c, b) in runs:
            r = p.add_run()
            r.text = t
            _set_font(r, s, c, b)
    return tb


def header(slide, title, sub=None):
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.13), Inches(0.62), fill=RED)
    textbox(slide, Inches(0.85), Inches(0.38), Inches(11.9), Inches(0.7),
            [(title, 25, INK, True)])
    if sub:
        textbox(slide, Inches(0.85), Inches(1.02), Inches(11.9), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    rect(slide, Inches(0.55), Inches(1.42), Inches(12.23), Pt(1.4), fill=LINE)
    footer(slide)


def footer(slide):
    textbox(slide, Inches(0.55), Inches(7.06), Inches(9.0), Inches(0.3),
            [("SoftBank × ウルトラマン 未来防衛隊キャラバン(ご提案)", 9, SILVER, False)])
    textbox(slide, Inches(12.2), Inches(7.06), Inches(0.6), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=11.5,
              header_size=11.5, row_h=0.34, header_fill=NAVY):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_h * n_rows))
    table = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(w * cw / total)
    # header
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htxt
        _set_font(r, header_size, WHITE, True)
    # body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            bold = False
            color = INK
            text = val
            if isinstance(val, tuple):
                text, bold, color = val[0], val[1], (val[2] if len(val) > 2 else INK)
            r = p.add_run()
            r.text = text
            _set_font(r, font_size, color, bold)
    return table


def bullets(slide, x, y, w, h, items, size=14, gap=8):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            t, opts = it
            lines.append([("● ", size - 3, RED, True), (t, opts.get("size", size), opts.get("color", INK), opts.get("bold", False))])
        else:
            lines.append([("● ", size - 3, RED, True), (it, size, INK, False)])
    textbox(slide, x, y, w, h, lines, space_after=gap)


def flow_down(slide, x, y, w, steps, box_h=0.52, gap=0.16, size=13,
              fill=PALE, first_fill=None, last_fill=None):
    """縦フロー(ボックス+下矢印)"""
    cy = y
    n = len(steps)
    for i, s in enumerate(steps):
        f = fill
        tc = INK
        bold = False
        if i == 0 and first_fill:
            f, tc, bold = first_fill, WHITE, True
        if i == n - 1 and last_fill:
            f, tc, bold = last_fill, WHITE, True
        bx = rect(slide, x, cy, w, Inches(box_h), fill=f, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        bx.adjustments[0] = 0.16
        tf = bx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        _set_font(r, size, tc, bold)
        cy += Inches(box_h)
        if i < n - 1:
            ar = rect(slide, x + w / 2 - Inches(0.09), cy + Inches(0.015),
                      Inches(0.18), Inches(gap - 0.03), fill=RED,
                      shape=MSO_SHAPE.DOWN_ARROW)
            cy += Inches(gap)


def divider(num, title, sub=None):
    s = add_slide()
    rect(s, 0, 0, SW, SH, fill=NAVY)
    rect(s, Inches(0.9), Inches(3.05), Inches(0.16), Inches(1.15), fill=RED)
    textbox(s, Inches(1.35), Inches(2.55), Inches(11), Inches(0.6),
            [(num, 18, SILVER, True)])
    textbox(s, Inches(1.35), Inches(3.05), Inches(11.2), Inches(1.2),
            [(title, 34, WHITE, True)])
    if sub:
        textbox(s, Inches(1.35), Inches(4.35), Inches(10.8), Inches(1.2),
                [(sub, 15, RGBColor(0xC3, 0xC9, 0xD2), False)])
    textbox(s, Inches(12.2), Inches(7.06), Inches(0.6), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)
    return s


# =============================================================================
# 1. 表紙
# =============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(4.9), SW, Inches(0.06), fill=RED)
textbox(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5),
        [("ソフトバンク株式会社 御中", 15, RGBColor(0xC3, 0xC9, 0xD2), False)])
textbox(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(2.6), [
    ("SoftBank × ウルトラマン", 30, SILVER, True),
    ("未来防衛隊キャラバン", 54, WHITE, True),
], space_after=10)
textbox(s, Inches(0.9), Inches(4.05), Inches(11.6), Inches(0.9), [
    ("商業施設に親子・ファミリーを呼び込み、SoftBankショップ・BB・でんき・モバイル相談へ送客する", 16, WHITE, False),
    ("店頭イベント戦略のご提案", 16, WHITE, False),
], space_after=4)
textbox(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.4), [
    ("ウルトラマンが、家族を商業施設へ呼び込む。SoftBankが、暮らしの通信・電気・安心を守る。", 13.5, RGBColor(0xE8, 0xB4, 0xB8), False),
    ("2026年7月|Confidential|金額はすべて税別・概算", 12, SILVER, False),
], space_after=8)

# =============================================================================
# 2. エグゼクティブサマリー
# =============================================================================
s = add_slide()
header(s, "エグゼクティブサマリー", "本企画はキャラクター露出ではなく、商業施設への「目的来館」を作り、店頭相談へ送客する施策です")
box = rect(s, Inches(0.55), Inches(1.7), Inches(12.23), Inches(1.15), fill=PALE_RED)
textbox(s, Inches(0.85), Inches(1.86), Inches(11.7), Inches(0.9), [
    [("結論:", 15, RED, True),
     ("ウルトラマンを“商業施設に顧客を呼び込む店頭イベントの核”として活用し、来館した親子・ファミリー層を", 15, INK, False)],
    [("SoftBankショップ、BB、でんき、モバイル相談へ送客する。", 15, INK, True),
     ("推奨は、FY27年度の年間活用を前提とした1.5億円プラン(B案)です。", 15, RED, True)],
], space_after=4)
bullets(s, Inches(0.85), Inches(3.15), Inches(11.8), Inches(3.6), [
    ("背景:現行のキャラクターIP施策は、FY27以降もコンテンツ活用が継続される一方、イベントでの着ぐるみ活用は終了予定と伺っております。", {}),
    ("課題:FY27以降、商業施設イベントにおける「会いに行く理由」=リアル集客機能に空白が生まれます。", {}),
    ("解決:現行IPの置き換えではなく、リアルイベント集客機能をウルトラマンで補完します。2026年7月10日にシリーズ60周年を迎える追い風もあります。", {}),
    ("仕組み:ヒーロー撮影会・館内ミッションラリー・デジタル隊員証で商業施設に呼び込み、SoftBankショップの相談(スマホ・BB・でんき)へ自然に送客します。", {}),
    ("ご提案:A案3,500万円(短期検証)/B案1億5,000万円(FY27年間・推奨)/C案2億5,000万円(60周年大型展開)。", {"bold": True}),
], size=14, gap=12)

# =============================================================================
# 3. 提案ポジション(貴社コンシューマ事業戦略との接続)
# =============================================================================
s = add_slide()
header(s, "本提案の位置づけ:コンシューマ事業のリアル接点拡張施策",
       "単なるキャラクターコラボではなく、貴社コンシューマ事業の成長に直結する施策としてご提案します")
bullets(s, Inches(0.85), Inches(1.75), Inches(11.8), Inches(1.7), [
    "貴社コンシューマ事業は、付加価値サービスを備えた料金プランの浸透、グループ経済圏の活用、長期利用が期待できるユーザー層の獲得・定着を通じ、モバイルサービス売上とセグメント利益の継続成長を掲げていらっしゃいます。",
    "SoftBank/ワイモバイル/LINEMOのマルチブランド展開も、大きな競争優位と拝察しております。",
    "本企画は、この戦略に対して「商業施設でのリアル顧客接点の創出」という形で貢献するものです。",
], size=13.5, gap=10)
box = rect(s, Inches(0.85), Inches(4.15), Inches(11.6), Inches(1.7), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box.adjustments[0] = 0.08
textbox(s, Inches(1.25), Inches(4.42), Inches(10.8), Inches(1.3), [
    ("本企画の主語", 13, RGBColor(0xE8, 0xB4, 0xB8), True),
    ("ウルトラマン60周年を活用し、商業施設に親子・ファミリーを呼び込み、", 17, WHITE, True),
    ("SoftBankショップ・BB・でんき・モバイル相談へ送客する、コンシューマ事業成長施策。", 17, WHITE, True),
], space_after=6)

# =============================================================================
# 4. 現行キャラクターIP施策の整理
# =============================================================================
s = add_slide()
header(s, "現行キャラクターIP施策の整理(ご共有いただいた前提)",
       "FY27以降もコンテンツ活用は継続予定。ただし着ぐるみ活用が終了し、リアル集客機能に課題が残ります")
add_table(s, Inches(0.7), Inches(1.75), Inches(8.1),
          ["項目", "現状"],
          [
              ["契約形態", "年度単位(4月〜3月)での契約"],
              ["利用可能領域", "イベント、BB、でんき"],
              ["ノベルティ", "作成数制限なし(数量に応じた許諾対価が別途必要)"],
              ["他コンテンツ制限", "なし(他IPとの併用が可能)"],
              ["着ぐるみ活用", ("FY26 3月まで", True, INK)],
              ["FY27以降", ("コンテンツ活用は継続予定。ただし着ぐるみ活用は終了", True, RED)],
              ["その他", "複数のキャラクターIP活用も並行してご検討中"],
          ],
          col_widths=[1, 2.4], row_h=0.42)
box = rect(s, Inches(9.05), Inches(1.75), Inches(3.7), Inches(3.8), fill=PALE)
textbox(s, Inches(9.3), Inches(2.0), Inches(3.2), Inches(3.4), [
    ("ここからの論点", 14, RED, True),
    ("① 現行IPはBB・でんき・既存コンテンツで継続活用が有効", 12.5, INK, False),
    ("② ただしFY27以降、リアルイベントの「目的来館」機能に空白が生まれる", 12.5, INK, False),
    ("③ その空白を補完できるIPが必要になる", 12.5, INK, True),
], space_after=10)

# =============================================================================
# 5. FY27の課題
# =============================================================================
s = add_slide()
header(s, "FY27以降の課題:リアル集客機能の空白",
       "着ぐるみ利用不可により、商業施設イベントの「会いに行く理由」が弱まります")
flow_down(s, Inches(1.2), Inches(1.8), Inches(6.4), [
    "現行IPのコンテンツ活用はFY27以降も継続",
    "ただし、着ぐるみは利用不可",
    "「会いに行く理由」「写真を撮る理由」「親子で参加する理由」が弱まる",
    "商業施設イベントでのリアル集客機能に空白",
], box_h=0.62, gap=0.28, size=13.5, last_fill=RED)
box = rect(s, Inches(8.2), Inches(1.8), Inches(4.55), Inches(3.9), fill=PALE)
textbox(s, Inches(8.5), Inches(2.05), Inches(4.0), Inches(3.5), [
    ("FY27の店頭イベントに必要なもの", 14, RED, True),
    ("親子・ファミリーを商業施設へ呼び込める「リアル集客コンテンツ」", 15, INK, True),
    ("・ヒーローの登場/撮影会", 12.5, INK, False),
    ("・ステージ/館内ラリーとの親和性", 12.5, INK, False),
    ("・親子二世代での認知", 12.5, INK, False),
    ("・通信・家族・暮らし・未来との文脈接続", 12.5, INK, False),
], space_after=9)
textbox(s, Inches(0.85), Inches(6.05), Inches(11.8), Inches(0.7), [
    [("答え:", 15, RED, True),
     ("現行IPを置き換えるのではなく、FY27以降担いにくくなる“リアルイベント集客機能”をウルトラマンで補完する。", 15, INK, True)],
])

# =============================================================================
# 6. IPポートフォリオ
# =============================================================================
s = add_slide()
header(s, "IPポートフォリオにおけるウルトラマンの役割",
       "既存IPと競合させず、役割分担で全体最適を作ります(独占は求めない前提)")
add_table(s, Inches(0.7), Inches(1.85), Inches(11.9),
          ["IP", "役割"],
          [
              ["現行キャラクターIP", "BB・でんき・既存コンテンツで継続活用(FY27以降も活用継続)"],
              ["かわいい系キャラクターIP(並行ご検討中)", "SNS拡散、女性・低年齢層向け"],
              [("ウルトラマン(本提案)", True, RED), ("商業施設イベント、親子撮影会、ヒーロー登場、通信・未来・防衛文脈", True, INK)],
          ],
          col_widths=[1.1, 2.3], row_h=0.62, font_size=14)
box = rect(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.6), fill=PALE_RED)
textbox(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(1.2), [
    ("ウルトラマンの位置づけ = 「リアルイベント集客IP」", 16, RED, True),
    ("ヒーロー登場・撮影会・ステージ・館内ミッションラリーと相性が高く、商業施設に人を呼ぶイベントの核になれる。", 14, INK, False),
], space_after=8)

# =============================================================================
# 7. なぜ今ウルトラマンか
# =============================================================================
s = add_slide()
header(s, "なぜ今、ウルトラマンなのか — 60周年「with U」との接続",
       "2026年7月10日にシリーズ60周年。公式60周年プロジェクトのコアメッセージ「with U」はSoftBankコンシューマ事業と好相性です")
add_table(s, Inches(0.7), Inches(1.8), Inches(11.9),
          ["ウルトラマンの文脈", "SoftBank側の接続"],
          [
              ["守る", "家族のスマホ安全、通信の安心、暮らしの電気"],
              ["つながる", "モバイル、BB、Wi-Fi、家族契約"],
              ["未来", "5G、AI、デジタルサービス、次世代の生活"],
              ["親子二世代", "商業施設のファミリー集客"],
              ["ヒーロー登場", "目的来館・撮影会・イベント集客"],
          ],
          col_widths=[1, 2.2], row_h=0.5, font_size=13.5)
textbox(s, Inches(0.85), Inches(5.05), Inches(11.8), Inches(1.4), [
    [("重要な視点:", 15, RED, True),
     ("ウルトラマンを「懐かしいIP」として使うのではなく、", 15, INK, False)],
    ("通信・家族・暮らし・未来をわかりやすく伝える“体験装置”として使う。", 17, INK, True),
], space_after=8)

# =============================================================================
# 8. 企画コンセプト(中扉)
# =============================================================================
divider("CONCEPT", "SoftBank × ウルトラマン 未来防衛隊キャラバン",
        "ウルトラヒーローの登場、親子撮影会、館内ミッションラリー、デジタル隊員証、SoftBankショップでのスマホ・BB・でんき相談を一体化し、\n商業施設全体の来館動機を創出する「送客型キャンペーン」")

# =============================================================================
# 9. 企画コンセプト詳細
# =============================================================================
s = add_slide()
header(s, "企画コンセプト", "商業施設とSoftBank、双方に価値が生まれる一体型イベントパッケージ")
textbox(s, Inches(0.85), Inches(1.75), Inches(11.8), Inches(1.2), [
    ("ウルトラマンを店頭イベントの核として活用し、商業施設に親子・ファミリー層を呼び込む送客型キャンペーン。", 15.5, INK, True),
    ("ヒーロー登場/親子撮影会/館内ミッションラリー/デジタル隊員証/ショップでのスマホ・BB・でんき相談を一体化。", 14, SUB, False),
], space_after=8)
b1 = rect(s, Inches(0.85), Inches(3.15), Inches(5.75), Inches(3.1), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b1.adjustments[0] = 0.05
textbox(s, Inches(1.15), Inches(3.4), Inches(5.15), Inches(2.6), [
    ("商業施設にとって", 15, RED, True),
    ("・週末・祝日のファミリー集客", 13.5, INK, False),
    ("・館内回遊の増加", 13.5, INK, False),
    ("・滞在時間の向上", 13.5, INK, False),
    ("・テナント全体を盛り上げる参加型イベント", 13.5, INK, False),
], space_after=9)
b2 = rect(s, Inches(6.95), Inches(3.15), Inches(5.75), Inches(3.1), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b2.adjustments[0] = 0.05
textbox(s, Inches(7.25), Inches(3.4), Inches(5.15), Inches(2.6), [
    ("SoftBankにとって", 15, RED, True),
    ("・店舗前通行量の増加", 13.5, INK, False),
    ("・相談件数の創出", 13.5, INK, False),
    ("・家族契約・機種変更への送客", 13.5, INK, False),
    ("・BB(SoftBank 光/Air)・でんきへの送客", 13.5, INK, False),
    [("・最終KGI:新規契約(モバイル・BB・でんき)の獲得", 13.5, RED, True)],
], space_after=9)

# =============================================================================
# 10. 中核戦略
# =============================================================================
s = add_slide()
header(s, "中核戦略:「店頭で待つ」から「商業施設へ呼び込む」へ",
       "すでに来店意思のある人への販促ではなく、来店意向のない層を呼び込み、相談接点へ転換します")
b1 = rect(s, Inches(0.85), Inches(1.85), Inches(4.6), Inches(4.3), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b1.adjustments[0] = 0.04
textbox(s, Inches(1.15), Inches(2.1), Inches(4.0), Inches(3.7), [
    ("旧来型の店頭施策", 15, SUB, True),
    ("SoftBankショップに来た人へ販促する", 14, INK, False),
    ("=すでに来店意思のある人への施策", 13, SUB, False),
    ("新規の来館動機は作れない", 13, SUB, False),
], space_after=12)
flow_down(s, Inches(6.1), Inches(1.85), Inches(6.6), [
    "ウルトラマンを目的に商業施設へ来てもらう",
    "館内イベントに参加してもらう",
    "SoftBankショップへ送客する",
    "スマホ・BB・でんき相談へつなげる",
    "見積・予約・申込・後日フォローへ展開する",
], box_h=0.56, gap=0.2, size=13, first_fill=NAVY, last_fill=RED)
textbox(s, Inches(0.85), Inches(6.4), Inches(11.8), Inches(0.6), [
    [("本企画の最重要メッセージ:", 14, RED, True),
     ("店頭で待つ施策ではなく、商業施設に顧客を呼び込む施策。", 15, INK, True)],
])

# =============================================================================
# 11. 商業施設との相性
# =============================================================================
s = add_slide()
header(s, "商業施設とのWin-Win", "商業施設側の集客課題を解決するイベントパッケージのため、施設側も受け入れやすい企画です")
add_table(s, Inches(0.7), Inches(1.8), Inches(11.9),
          ["商業施設側の課題", "ウルトラマン施策での解決"],
          [
              ["週末・祝日の集客を増やしたい", "ウルトラマン撮影会で「目的来館」を作る"],
              ["親子・ファミリー層を呼びたい", "親子二世代に認知があるIP"],
              ["館内回遊を増やしたい", "ミッションラリーで複数フロアを回遊"],
              ["滞在時間を伸ばしたい", "ステージ・撮影会・ラリー・相談を組み合わせる"],
              ["テナント全体を盛り上げたい", "館内参加型イベントとして展開できる"],
              ["天候に左右されたくない", "屋内商業施設イベントと相性が高い"],
          ],
          col_widths=[1, 1.4], row_h=0.47, font_size=13)
textbox(s, Inches(0.85), Inches(5.35), Inches(11.8), Inches(1.2), [
    ("SoftBankショップ単体でのイベントよりも、商業施設全体を巻き込むことで、", 14, INK, False),
    ("「来館動機」「館内回遊」「店舗送客」を同時に作れる。", 16, RED, True),
], space_after=6)

# =============================================================================
# 12. ターゲット
# =============================================================================
s = add_slide()
header(s, "ターゲット", "来店意向のない親子・ファミリー層を中心に、既存ユーザー・他キャリア利用者までカバーします")
add_table(s, Inches(0.7), Inches(1.75), Inches(7.3),
          ["メインターゲット", "狙い"],
          [
              ["親子・ファミリー層", "目的来館、家族契約、子どもスマホ相談"],
              ["30〜50代の保護者", "家庭の通信費、BB、でんき、機種変更相談"],
              ["既存SoftBankユーザー", "料金見直し、機種変更、BB・でんき追加"],
              ["他キャリア利用者", "料金診断、MNP、ワイモバイル・LINEMO送客"],
              ["商業施設来館者", "偶発接触から相談・予約へ誘導"],
          ],
          col_widths=[1, 1.5], row_h=0.5, font_size=12.5)
box = rect(s, Inches(8.3), Inches(1.75), Inches(4.45), Inches(4.55), fill=PALE)
textbox(s, Inches(8.55), Inches(1.98), Inches(3.95), Inches(4.2), [
    ("特に狙うべき生活者", 14, RED, True),
    ("・子どものスマホを検討している家庭", 12.5, INK, False),
    ("・家族の通信費を見直したい家庭", 12.5, INK, False),
    ("・自宅Wi-Fiに不満がある家庭", 12.5, INK, False),
    ("・電気料金の見直しに関心がある家庭", 12.5, INK, False),
    ("・商業施設に週末来館する親子", 12.5, INK, False),
    ("・料金相談のきっかけがない既存ユーザー", 12.5, INK, False),
], space_after=9)

# =============================================================================
# 13. ユーザー導線
# =============================================================================
s = add_slide()
header(s, "ユーザー体験設計(生活者導線)", "認知から来館・参加・相談・申込まで、8ステップの一気通貫導線")
steps = [
    "① SNS・商業施設告知・SoftBank告知でイベントを知る",
    "② 週末に親子で商業施設へ来館",
    "③ ウルトラマン撮影会・ステージを見る",
    "④ 館内ミッションラリーに参加",
    "⑤ デジタル隊員証を発行",
    "⑥ SoftBankショップで特典・相談導線に接触",
    "⑦ スマホ料金診断・家族契約・BB・でんき相談",
    "⑧ 見積・来店予約・オンライン申込・後日フォロー",
]
# 2列 x 4段のフロー
flow_down(s, Inches(0.85), Inches(1.85), Inches(5.7), steps[:4], box_h=0.62, gap=0.28, size=13, first_fill=NAVY)
flow_down(s, Inches(6.95), Inches(1.85), Inches(5.7), steps[4:], box_h=0.62, gap=0.28, size=13, last_fill=RED)
ar = rect(s, Inches(6.62), Inches(3.35), Inches(0.3), Inches(0.24), fill=RED, shape=MSO_SHAPE.RIGHT_ARROW)
textbox(s, Inches(0.85), Inches(6.15), Inches(11.8), Inches(0.8), [
    [("ポイント:", 14, RED, True),
     ("いきなり契約を狙いすぎない。まず「来館・参加・相談・見積・予約」を作り、その後、店舗またはオンラインで申込につなげる。", 14, INK, True)],
])

# =============================================================================
# 14. 店頭イベントメニュー
# =============================================================================
s = add_slide()
header(s, "店頭イベントの具体内容:「ウルトラマン未来防衛隊 in 商業施設」",
       "集客・回遊・接点化・相談・送客の機能を1パッケージ化した基本メニュー")
add_table(s, Inches(0.7), Inches(1.75), Inches(11.9),
          ["施策", "内容", "SoftBankへの効果"],
          [
              ["ウルトラヒーロー撮影会", "親子で写真撮影", "目的来館を作る"],
              ["未来防衛隊ステージ", "通信・家族・暮らしの安心がテーマのミニステージ", "ブランド好意度向上"],
              ["館内ミッションラリー", "館内を回遊しながらミッション達成", "商業施設回遊、滞在時間向上"],
              ["デジタル隊員証", "QRで参加登録、スマホ上で表示", "顧客接点化"],
              ["家族スマホ安全診断", "子どものスマホ利用、料金、セキュリティ相談", "モバイル相談"],
              ["おうち通信防衛診断", "Wi-Fi、BB、通信環境の相談", "SoftBank 光/Air送客"],
              ["くらしの電気防衛診断", "電気料金、でんき契約の見直し", "でんき送客"],
              ["ショップ限定特典", "相談・見積でノベルティ", "店舗送客"],
          ],
          col_widths=[1, 1.6, 1.1], row_h=0.51, font_size=12.5)

# =============================================================================
# 15. ミッションラリー
# =============================================================================
s = add_slide()
header(s, "未来防衛隊ミッションラリー", "「基本ミッション(SoftBank固定)」+「拡張ミッション(施設側で追加自由)」の2層設計です")
add_table(s, Inches(0.7), Inches(1.75), Inches(7.5),
          ["基本ミッション(SoftBank・固定)", "目的"],
          [
              ["家族スマホ安全ミッション", "子どもスマホ・見守り・セキュリティ相談"],
              ["おうちWi-Fi防衛ミッション", "SoftBank 光/Air相談"],
              ["でんき防衛ミッション", "でんき契約・電気料金見直し"],
              ["料金見直しミッション", "家族通信費の診断"],
              ["ショップ基地ミッション", "SoftBankショップへの送客"],
              ["フォトミッション", "ウルトラマン撮影会参加"],
              ["デジタル隊員証ミッション", "顧客接点の取得"],
          ],
          col_widths=[1, 1.35], row_h=0.44, font_size=12.5)
box = rect(s, Inches(8.5), Inches(1.75), Inches(4.25), Inches(3.35), fill=PALE)
textbox(s, Inches(8.75), Inches(1.98), Inches(3.75), Inches(3.0), [
    ("ラリー完了特典", 14, RED, True),
    ("・未来防衛隊ステッカー", 12.5, INK, False),
    ("・デジタル壁紙", 12.5, INK, False),
    ("・隊員証カード/クリアファイル", 12.5, INK, False),
    ("・店舗相談者限定ノベルティ", 12.5, INK, False),
    ("・抽選応募権", 12.5, INK, False),
], space_after=8)
box = rect(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.92), fill=PALE)
textbox(s, Inches(1.0), Inches(5.52), Inches(11.3), Inches(0.7), [
    [("拡張ミッション(施設・テナント側で追加自由):", 12.5, RED, True),
     ("館内テナントを回遊するミッションを施設側で自由に追加可能。テナント送客協賛も設計でき、", 12.5, INK, False)],
    ("施設全体の回遊・売上に貢献するため、場所代交渉の材料にもなります。", 12.5, INK, False),
], space_after=3)
textbox(s, Inches(0.85), Inches(6.55), Inches(11.8), Inches(0.45), [
    [("ノベルティ設計:", 12, RED, True),
     ("作成数に上限を設けず、数量に応じた許諾対価を別途お支払いする設計を想定しています。", 12, INK, False)],
])

# =============================================================================
# 15b. デジタル隊員証の設計
# =============================================================================
s = add_slide()
header(s, "デジタル隊員証の設計", "アプリ開発を伴わない軽い仕組みで、「接客の質」と「再来館動機」を同時に作ります")
add_table(s, Inches(0.7), Inches(1.75), Inches(11.9),
          ["設計ポイント", "内容"],
          [
              ["発行フロー", "QRからアンケートに回答 → その場でオリジナル隊員証“画像”を発行(アプリ開発不要・画像保存のみの軽い設計)"],
              [("デザインの出し分け", True, RED), ("アンケート結果に応じて隊員証デザインを出し分け → スタッフが隊員証を見た瞬間に、ご案内すべき相談メニューを判別できる(お客様に合ったミッションだけをご案内)", True, INK)],
              ["ランクアップ演出", "来場・ミッション達成で隊員ランクが上がる演出 → 再来館・リピート参加の動機に"],
              ["体験への配慮", "アンケートは遊びの中で自然に答えられる設計とし、営業感を出さない"],
              ["将来拡張", "ポイント連携・アプリ連携・LINE連携は将来オプション(初期スコープ外として切り分け)"],
          ],
          col_widths=[1, 3.2], row_h=0.62, font_size=12)
textbox(s, Inches(0.85), Inches(5.85), Inches(11.8), Inches(0.8), [
    [("狙い:", 14, RED, True),
     ("既にご契約中のお客様に不要なご案内をしない=接客効率と顧客体験を両立し、獲得ターゲットに接客リソースを集中します。", 14, INK, True)],
])

# =============================================================================
# 16. ショップ送客設計
# =============================================================================
s = add_slide()
header(s, "SoftBankショップへの送客設計", "SoftBankショップを「未来防衛隊基地」として位置づけ、ラリーの最終ミッションで店舗へ誘導します")
flow_down(s, Inches(0.85), Inches(1.85), Inches(5.4), [
    "中央イベントスペース(広場)",
    "ミッションラリー参加",
    "SoftBankショップで最終ミッション",
    "料金診断・BB診断・でんき相談",
    "見積・予約・申込",
], box_h=0.56, gap=0.22, size=13, first_fill=NAVY, last_fill=RED)
add_table(s, Inches(6.75), Inches(1.85), Inches(6.0),
          ["店舗の相談メニュー", "目的"],
          [
              ["家族スマホ料金診断", "家族契約、料金見直し、機種変更"],
              ["子どもスマホ相談", "初スマホ、安心設定、フィルタリング"],
              ["おうちWi-Fi診断", "SoftBank 光/Air送客"],
              ["でんき料金診断", "でんき契約相談"],
              ["MNP相談", "他キャリア利用者の獲得"],
              ["ワイモバイル/LINEMO案内", "価格・オンライン志向ユーザーの受け皿"],
          ],
          col_widths=[1, 1.25], row_h=0.44, font_size=12)
textbox(s, Inches(0.85), Inches(6.15), Inches(11.8), Inches(0.8), [
    ("マルチブランド(SoftBank/ワイモバイル/LINEMO)展開により、来館者のニーズに応じた最適な受け皿を用意できる。", 13.5, INK, False),
])

# =============================================================================
# 17. PayPay/LINE/LYP
# =============================================================================
s = add_slide()
header(s, "PayPay/LINE/LYPとの将来連携", "初期はPayPayポイント原資を必須とせず、グループ経済圏への拡張余地を残します")
bullets(s, Inches(0.85), Inches(1.75), Inches(11.8), Inches(0.9), [
    "PayPayとLINEヤフーは、LINE(国内月間利用者数1億)とPayPay(登録ユーザー7,400万人)のアカウント連携を2026年夏より開始と発表。LINEヤフー側とも既に対話を開始しています。",
], size=13.5)
add_table(s, Inches(0.7), Inches(2.5), Inches(11.9),
          ["連携先", "展開案"],
          [
              ["PayPay", "商業施設内の対象店舗利用で抽選、PayPayポイント施策"],
              ["LINE", "イベント告知、デジタル隊員証、来場後フォロー"],
              ["LYP", "SoftBankユーザー向け特典訴求"],
              ["Yahoo! JAPAN", "特集面・キャンペーン告知"],
              ["My SoftBank", "既存ユーザー向けイベント告知"],
          ],
          col_widths=[1, 2.5], row_h=0.46, font_size=13)
textbox(s, Inches(0.85), Inches(5.35), Inches(11.8), Inches(0.8), [
    [("注記:", 13.5, RED, True),
     ("PayPayポイント原資、アプリ本体改修、LINEミニアプリ開発は本見積に含めず、別途費用とする。", 13.5, INK, False)],
])

# =============================================================================
# 18. 見積(中扉)
# =============================================================================
divider("PLAN & BUDGET", "見積もり3パターン",
        "すべて税別・概算。円谷プロ側の正式なIP使用料、監修条件、ヒーロー稼働条件、使用媒体、使用期間により変動します。")

# =============================================================================
# 19. 3プラン比較
# =============================================================================
s = add_slide()
header(s, "3プラン比較", "推奨はB案:FY27標準プラン(1億5,000万円)です")
add_table(s, Inches(0.7), Inches(1.85), Inches(11.9),
          ["プラン", "金額(税別)", "位置づけ", "推奨度"],
          [
              ["A案:短期検証プラン", "3,500万円", "商業施設イベントの小規模検証(FY27期首:2027年4月〜)", "○"],
              [("B案:FY27標準プラン", True, RED), ("1億5,000万円", True, RED), ("年間展開の推奨案。FY27以降のリアル集客を補完", True, INK), ("◎", True, RED)],
              ["C案:60周年大型展開プラン", "2億5,000万円", "全国大型キャンペーン(B案実施後の拡張案)", "○"],
          ],
          col_widths=[1.3, 0.8, 2.2, 0.45], row_h=0.68, font_size=13.5)
box = rect(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.5), fill=PALE_RED)
textbox(s, Inches(1.0), Inches(5.08), Inches(11.3), Inches(1.1), [
    ("推奨の進め方", 14, RED, True),
    ("A案パイロット(FY27期首:2027年4月〜)で反応を検証 → B案で年間展開 → 成果に応じてC案(全国)へ拡張。", 15, INK, True),
], space_after=6)

# =============================================================================
# 19b. 見積の前提と積算の考え方
# =============================================================================
s = add_slide()
header(s, "お見積りの前提と積算の考え方",
       "各金額の積算根拠を明示のうえ、円谷プロ側との条件確認をもって精査版をご提示します")
add_table(s, Inches(0.7), Inches(1.75), Inches(11.9),
          ["項目", "積算の考え方"],
          [
              ["年間IP使用料(B案・C案:8,000万円)",
               "円谷プロとの年間包括契約を想定した上限側の概算。60周年記念年のプレミアム、ヒーロー稼働との一体運用、広範な使用範囲(イベント/BB/でんき/店頭/Web/SNS)を含む水準"],
              ["短期IP使用料(A案:1,200万円)",
               "月額換算約300〜400万円×3〜4か月。短期プロモーションライセンスの一般的な水準(1クール数百万〜数千万円)の範囲内"],
              ["ヒーロー登場・イベント運営費",
               "キャラクターショーの実勢価格(1回あたり数十万円+設営・音響・運営)をもとに、1施設あたり週末1〜2回の開催を想定して積算。B案は首都圏中心、全国展開はC案で想定"],
              ["広告・送客費",
               "集客は商業施設の館内媒体・施設告知との相互送客が主体。広告費はデジタル告知の補助として計上"],
              ["事務局・レポート",
               "店舗・イベント主体の軽量運用を想定。大規模な抽選・賞品発送を伴う場合は別途お見積り"],
              ["変動要因と確定プロセス",
               "正式なIP使用料・監修条件・稼働条件・使用媒体により変動。2026年9〜10月の条件確認をもって精査版をご提示"],
          ],
          col_widths=[1.1, 2.5], row_h=0.66, font_size=11.5)

# =============================================================================
# 20. A案
# =============================================================================
s = add_slide()
header(s, "A案:短期検証プラン|3,500万円(税別)", "FY27期首(2027年4月〜)の3〜4か月で、ウルトラマンを核にした商業施設送客の反応を検証します")
add_table(s, Inches(0.7), Inches(1.75), Inches(6.0),
          ["見積項目", "金額"],
          [
              ["ウルトラマンIP使用料・監修費", "1,200万円"],
              ["企画設計・PM", "250万円"],
              ["キービジュアル・コピー制作", "300万円"],
              ["LP/応募フォーム制作", "350万円"],
              ["館内ラリーツール制作", "250万円"],
              ["ヒーロー登場・イベント運営", "500万円"],
              ["SNS/バナー制作", "250万円"],
              ["広告・送客費", "250万円"],
              ["事務局・レポート", "150万円"],
              [("合計", True), ("3,500万円", True, RED)],
          ],
          col_widths=[1.8, 0.7], row_h=0.4, font_size=11.5)
add_table(s, Inches(7.1), Inches(1.75), Inches(5.6),
          ["KPI目安", "目標"],
          [
              ["実施商業施設数", "2〜3施設"],
              ["総イベント来場者", "5,000〜1.5万人"],
              ["撮影会参加組数", "500〜1,500組"],
              ["館内ラリー参加", "1,500〜4,000人"],
              ["SoftBankショップ送客", "300〜800人"],
              ["相談件数", "100〜300件"],
              ["見積・来店予約", "50〜150件"],
          ],
          col_widths=[1.3, 1], row_h=0.4, font_size=11.5)
textbox(s, Inches(7.1), Inches(5.05), Inches(5.6), Inches(1.6), [
    ("位置づけ", 13, RED, True),
    ("導入の第一歩として有効。ただしFY27以降の年間集客機能を担うには規模が限定的。", 12.5, INK, False),
], space_after=6)
textbox(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5), [
    ("※IP使用料は月額換算約300〜400万円で、短期プロモーションライセンスの一般的な水準の範囲内。ヒーロー登場・運営費は2〜3施設×週末開催(計8〜12回想定)で積算しています。", 11, SUB, False),
])

# =============================================================================
# 21. B案
# =============================================================================
s = add_slide()
header(s, "B案:FY27標準プラン|1億5,000万円(税別)【推奨】",
       "FY27年度(2027年4月〜2028年3月)/商業施設10〜20施設(首都圏中心)/ヒーロー年間稼働枠/BB・でんき本格導線化/他IP併用可")
add_table(s, Inches(0.7), Inches(1.75), Inches(6.0),
          ["見積項目", "金額"],
          [
              ["年間IP使用料・監修費", "8,000万円"],
              ["ヒーロー登場・撮影会稼働枠", "1,500万円"],
              ["商業施設イベント運営費", "1,200万円"],
              ["企画設計・PM", "600万円"],
              ["キービジュアル・店頭イベント素材", "500万円"],
              ["LP/デジタル隊員証/計測導線", "550万円"],
              ["館内ラリー・イベントツール制作", "600万円"],
              ["ノベルティ初期制作・デザイン", "500万円"],
              ["BB/でんき向けコンテンツ", "300万円"],
              ["SNS/告知クリエイティブ", "350万円"],
              ["広告・送客費", "700万円"],
              ["事務局・レポート", "200万円"],
              [("合計", True), ("1億5,000万円", True, RED)],
          ],
          col_widths=[2, 0.7], row_h=0.355, font_size=11)
add_table(s, Inches(7.1), Inches(1.75), Inches(5.6),
          ["KPI目安", "目標"],
          [
              ["実施商業施設数", "10〜20施設"],
              ["総イベント来場者", "3万〜8万人"],
              ["撮影会参加組数", "3,000〜8,000組"],
              ["館内ラリー参加", "1万〜3万人"],
              ["デジタル隊員証発行", "8,000〜2万人"],
              ["SoftBankショップ送客", "3,000〜8,000人"],
              ["相談件数", "1,000〜3,000件"],
              ["見積・来店予約", "500〜1,500件"],
              ["BB/でんき相談", "300〜1,000件"],
              ["モバイル関連相談", "700〜2,000件"],
          ],
          col_widths=[1.3, 1], row_h=0.355, font_size=11)
textbox(s, Inches(7.1), Inches(5.65), Inches(5.6), Inches(1.2), [
    ("※年間IP使用料8,000万円は、60周年プレミアム・ヒーロー稼働との一体運用・広範な使用範囲を含む年間包括契約想定の概算です。運営費・稼働枠は首都圏中心・1施設あたり週末1〜2回の開催を想定して積算(全国展開はC案)。KPIは最終KGI(新規契約数)の先行指標として設計しています。", 10.5, SUB, False),
])

# =============================================================================
# 22. B案を推奨する理由
# =============================================================================
s = add_slide()
header(s, "B案を推奨する理由", "IP使用料単体ではなく、「年間IP契約+リアルイベント実行費+店舗送客導線」の総額としてご評価ください")
textbox(s, Inches(0.85), Inches(1.75), Inches(11.8), Inches(0.8), [
    ("B案の1億5,000万円は、IP使用料だけの金額ではありません。年間の企画・制作・イベント運営・店舗送客までを含む、実行費一式の総額です。", 14, INK, False),
], space_after=4)
items = [
    "年間IP使用料", "ヒーロー登場・撮影会稼働枠", "商業施設イベント運営", "店頭ツール",
    "館内ラリー", "LP/デジタル隊員証", "BB/でんき向けコンテンツ", "SNS告知",
    "広告・送客", "事務局・レポート",
]
x0, y0 = Inches(0.85), Inches(2.6)
bw, bh = Inches(2.28), Inches(0.72)
gapx, gapy = Inches(0.12), Inches(0.16)
for i, t in enumerate(items):
    col, row = i % 5, i // 5
    bx = rect(s, x0 + (bw + gapx) * col, y0 + (bh + gapy) * row, bw, bh,
              fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx.adjustments[0] = 0.14
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    _set_font(r, 12, INK, True)
box = rect(s, Inches(0.85), Inches(4.75), Inches(11.75), Inches(1.5), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box.adjustments[0] = 0.08
textbox(s, Inches(1.2), Inches(5.05), Inches(11.1), Inches(1.0), [
    ("B案 = 年間IP契約 + リアルイベント実行費 + 店舗送客導線 の一体パッケージ", 18, WHITE, True),
], anchor=MSO_ANCHOR.MIDDLE)

# =============================================================================
# 22b. ご参考:場所代削減効果(実質負担の考え方)
# =============================================================================
s = add_slide()
header(s, "ご参考:場所代の削減効果と実質負担の考え方",
       "ファミリー集客企画は商業施設の場所代交渉の材料になり、本企画の実質負担を大きく圧縮し得ます")
bullets(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(1.5), [
    "商業施設のイベントスペース利用料は、1日50万〜100万円が標準的な水準と伺っております。",
    "施設側は「来館者が喜ぶ集客企画」に対して、場所代を減額する交渉余地を持っています。",
    "本企画はファミリーの目的来館・館内回遊・滞在時間向上を直接作るため、場所代交渉の強い材料になります。",
], size=13, gap=7)
add_table(s, Inches(0.7), Inches(3.35), Inches(11.9),
          ["効果の範囲", "考え方(参考試算)"],
          [
              ["① 本企画の実施施設での直接効果", "10〜20施設×年間80〜160日の場所代負担が減額交渉の対象に(数千万円規模)"],
              ["② 全国のイベント網への波及", "全国数百か所で実施される店頭イベントの場所代交渉の標準材料に → 年間数億円規模のコスト構造に波及"],
          ],
          col_widths=[1.1, 2.4], row_h=0.6, font_size=12.5)
box = rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.0), fill=PALE_RED)
textbox(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.8), [
    ("場所代の削減効果が本企画費用(B案1億5,000万円)を上回れば、実質負担ゼロも視野に入ります。", 15.5, RED, True),
], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4), [
    ("※場所代・減額幅は施設・立地・企画内容および施設側との個別交渉により変動します。上記は考え方をお示しする参考値です。", 10.5, SUB, False),
])

# =============================================================================
# 23. C案
# =============================================================================
s = add_slide()
header(s, "C案:60周年大型展開プラン|2億5,000万円(税別)",
       "60周年文脈を活用した全社大型コンシューマキャンペーン(6〜12か月)。B案実施後の拡張案としてご提示します")
add_table(s, Inches(0.7), Inches(1.75), Inches(6.0),
          ["見積項目", "金額"],
          [
              ["年間IP使用料・広範囲監修費", "8,000万円"],
              ["ヒーロー登場・イベント稼働枠", "2,500万円"],
              ["戦略設計・全体PM", "1,200万円"],
              ["キービジュアル・コピー開発", "1,000万円"],
              ["Web動画/ショート動画制作", "2,500万円"],
              ["LP/ミッション基盤制作", "1,200万円"],
              ["AR/デジタルコンテンツ制作", "800万円"],
              ["全国店頭ツール制作", "1,700万円"],
              ["ノベルティ制作・配送管理", "1,500万円"],
              ["BB/でんき/モバイル向けコンテンツ", "500万円"],
              ["Web/SNS広告配信", "3,500万円"],
              ["PR/タイアップ", "400万円"],
              ["事務局・効果測定", "200万円"],
              [("合計", True), ("2億5,000万円", True, RED)],
          ],
          col_widths=[2, 0.7], row_h=0.335, font_size=10.5)
add_table(s, Inches(7.1), Inches(1.75), Inches(5.6),
          ["KPI目安", "目標"],
          [
              ["実施商業施設数", "30〜50施設"],
              ["総イベント来場者", "10万〜25万人"],
              ["撮影会参加組数", "1万〜3万組"],
              ["館内ラリー参加", "5万〜10万人"],
              ["デジタル隊員証発行", "3万〜8万人"],
              ["SoftBankショップ送客", "1万〜3万人"],
              ["相談件数", "3,000〜8,000件"],
              ["見積・来店予約", "1,500〜4,000件"],
              ["BB/でんき相談", "1,000〜3,000件"],
              ["PR露出", "50〜150件"],
          ],
          col_widths=[1.3, 1], row_h=0.35, font_size=11)
textbox(s, Inches(7.1), Inches(5.78), Inches(5.6), Inches(1.2), [
    ("※年間IP使用料はB案と同一の年間包括契約水準(8,000万円)で統一(映像・AR等の広範囲利用の監修費込み)。展開規模の拡大分は稼働枠・制作費・広告費側に計上しています。", 10.5, SUB, False),
])

# =============================================================================
# 24. 別途費用
# =============================================================================
s = add_slide()
header(s, "基本見積に含めない項目(別途費用)", "変動幅が大きい項目・権利体系が異なる項目は、基本見積から切り分けて透明性を担保します")
add_table(s, Inches(0.7), Inches(1.8), Inches(11.9),
          ["別途項目", "理由"],
          [
              ["TVCM媒体費", "出稿量で大きく変動するため"],
              ["交通広告・大型OOH", "掲出場所・期間で変動するため"],
              ["PayPayポイント原資", "付与条件・人数・上限で変動するため"],
              ["全国全店の大型施工", "店舗数・施工仕様で大きく変動するため"],
              ["販売用グッズの商品化ライセンス", "ノベルティとは権利体系が異なるため"],
              ["My SoftBank/PayPayアプリ本体改修", "開発範囲次第で別見積になるため"],
              ["大型タレント起用", "キャスティング費が別途必要なため"],
              ["商業施設の場所代", "施設条件により有償・無償が変わるため"],
          ],
          col_widths=[1.1, 1.4], row_h=0.48, font_size=13)

# =============================================================================
# 25. 契約条件案
# =============================================================================
s = add_slide()
header(s, "契約条件案(円谷プロ側への打診条件)", "独占は求めず、ウルトラマンを「リアルイベント集客IP」として他IPと役割分担する前提です")
add_table(s, Inches(0.7), Inches(1.75), Inches(11.9),
          ["項目", "条件案"],
          [
              ["契約期間", "FY27年度(2027年4月〜2028年3月)"],
              ["使用地域", "日本国内"],
              [("使用業種(プロモーション対象)", True, INK), ("携帯通信、ブロードバンド、でんき(コンシューマ向けプロモーションに限る)", True, INK)],
              ["使用媒体", "イベント、店頭ツール、LP/Web、SNS、メール、ノベルティ、イベント演出"],
              ["ヒーロー登場", "年間稼働枠を設定。追加稼働は別途"],
              ["ノベルティ", "作成数制限なし、数量に応じた許諾対価は別途"],
              ["サンプル納品/他IP併用/独占", "作成物ごとにサンプル提出/他IP併用可/独占は原則なし"],
              ["競合(バッティング)", "既存契約書(ディールメモ)の書式に準拠して業種バッティングを規定"],
              ["範囲外", "法人向け、PayPay等グループ他社サービス、TVCM、交通広告、大型OOH、販売グッズ、アプリ改修、ポイント原資"],
          ],
          col_widths=[1, 2.5], row_h=0.47, font_size=12)

# =============================================================================
# 26. スケジュール
# =============================================================================
s = add_slide()
header(s, "実施スケジュール案(FY26下期=準備、FY27=実働)",
       "FY26下期は設計・準備に充て、ウルトラマンのIP利用・実働は2027年4月(FY27期首)から開始します")
add_table(s, Inches(0.7), Inches(1.7), Inches(11.9),
          ["期間", "内容", "フェーズ"],
          [
              ["2026年7月〜8月", "初回提案・方向性合意", "準備"],
              ["2026年9月〜10月", "円谷プロ条件確認・概算見積精査・契約協議", "準備"],
              ["2026年11月〜2027年3月", "パイロット設計、商業施設候補の選定・調整、制作準備(IP実働前の準備期間)", "準備"],
              [("2027年4月", True, RED), ("FY27契約開始 — ウルトラマンのIP利用・実働はここから", True, RED), ("開始", True, RED)],
              ["2027年4月〜6月", "先行2〜3施設で小規模検証(A案フェーズ)", "検証"],
              ["2027年7月〜8月", "検証を踏まえ10〜20施設へ拡大。屋内商業施設中心の夏休み施策", "本展開"],
              ["2027年9月〜11月", "秋のファミリー送客施策", "本展開"],
              ["2027年12月〜2028年1月", "年末年始商業施設施策", "本展開"],
              ["2028年2月〜3月", "新生活・春商戦連動", "本展開"],
              ["2028年4月〜", "初年度の成果を踏まえ、複数年契約・全国展開(C案)への移行を協議", "拡張"],
          ],
          col_widths=[1.1, 1.9, 0.55], row_h=0.4, font_size=12)
textbox(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.5), [
    ("※現行IPの着ぐるみ活用がFY26年度末(2027年3月)まで可能なため、ウルトラマンの実働はFY27期首からの開始が最も効率的かつ空白なく接続できます。", 11, SUB, False),
])

# =============================================================================
# 27. 実施体制
# =============================================================================
s = add_slide()
header(s, "実施体制", "SoftBank・円谷プロダクション・提案側チームの三者連携で推進します")
add_table(s, Inches(0.7), Inches(1.8), Inches(11.9),
          ["役割", "担当"],
          [
              ["全体統括", "SoftBank コンシューマ事業側"],
              ["IP監修", "円谷プロダクション"],
              ["企画・制作", "提案側チーム"],
              ["商業施設調整", "SoftBank営業/代理店/イベント運営側"],
              ["店舗連携", "SoftBankショップ、代理店運営会社"],
              ["BB/でんき導線", "各サービス担当部門"],
              ["デジタル導線", "LP、計測、デジタル隊員証担当"],
              ["法務・ライセンス", "双方法務・ライセンス担当"],
              ["効果測定", "SoftBankマーケティング/提案側チーム"],
          ],
          col_widths=[1, 2], row_h=0.44, font_size=12.5)

# =============================================================================
# 28. 最終提案
# =============================================================================
s = add_slide()
header(s, "最終ご提案:B案 FY27標準プラン", "FY27以降のリアルイベント集客機能を担う、推奨プランです")
add_table(s, Inches(0.7), Inches(1.8), Inches(11.9),
          ["項目", "内容"],
          [
              ["金額", ("1億5,000万円 税別", True, RED)],
              ["期間", "FY27年度(2027年4月〜2028年3月)"],
              ["目的", "商業施設送客、SoftBankショップ送客、BB・でんき・モバイル相談創出"],
              ["実施規模", "10〜20商業施設"],
              ["中核施策", "ウルトラヒーロー登場、撮影会、館内ラリー、デジタル隊員証、店舗相談"],
              ["KPI", "イベント来場3万〜8万人、店舗送客3,000〜8,000人、相談1,000〜3,000件"],
              ["位置づけ", "FY27以降のリアルイベント集客機能を担う中核施策"],
          ],
          col_widths=[0.9, 2.9], row_h=0.5, font_size=13)
textbox(s, Inches(0.85), Inches(6.05), Inches(11.8), Inches(0.9), [
    [("次アクション:", 14, RED, True),
     ("方向性ご合意のうえ、円谷プロ側への条件打診と概算見積の精査(2026年9月〜10月)、パイロット設計・施設選定(2026年11月〜)に着手し、2027年4月の実働開始に備えます。", 14, INK, False)],
])

# =============================================================================
# 29. クロージング
# =============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, Inches(0.9), Inches(1.5), Inches(0.16), Inches(2.5), fill=RED)
textbox(s, Inches(1.35), Inches(1.45), Inches(11.0), Inches(2.7), [
    ("本企画の最終メッセージ", 15, SILVER, True),
    ("ウルトラマンでSoftBankを目立たせる企画ではない。", 25, WHITE, True),
    ("ウルトラマンを目的来館の核にして、商業施設に顧客を呼び込み、", 25, WHITE, True),
    ("SoftBankの店頭相談・BB・でんき・モバイル獲得へつなげる企画である。", 25, WHITE, True),
], space_after=12)
textbox(s, Inches(1.35), Inches(4.6), Inches(10.8), Inches(1.6), [
    ("現行IP施策との整合/FY27以降の課題/商業施設送客/コンシューマ事業への貢献 — すべてが一本につながります。", 14.5, RGBColor(0xC3, 0xC9, 0xD2), False),
    ("ご清聴ありがとうございました。", 14.5, RGBColor(0xC3, 0xC9, 0xD2), False),
], space_after=10)
rect(s, 0, Inches(6.9), SW, Inches(0.06), fill=RED)

out = sys.argv[1] if len(sys.argv) > 1 else "SoftBank_Ultraman_未来防衛隊キャラバン_提案資料.pptx"
prs.save(out)
print(f"saved: {out} / slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
