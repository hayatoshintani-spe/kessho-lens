# -*- coding: utf-8 -*-
"""SoftBank様ご依頼事項(2026/7/16)への見積回答資料を生成する(円谷フォーマット)。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
TSUB_NAVY = RGBColor(0x1F, 0x33, 0x8C)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")
COVER_IMG = os.path.join(ASSETS, "band_cover.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=12,
              header_size=12, row_h=0.45):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, header_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = text
            set_font(r, font_size, color, bold)
    return tbl


# =============================================================================
# 1. 表紙
# =============================================================================
s = add_slide()
s.shapes.add_picture(COVER_IMG, 0, 0, width=Inches(4.09), height=SH)
textbox(s, Inches(4.5), Inches(1.3), Inches(8.4), Inches(0.5),
        [("ソフトバンク株式会社 御中", 14, SUB, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(4.5), Inches(2.3), Inches(8.4), Inches(2.0), [
    ("ウルトラマンIP 年間ご契約", 30, INK, True),
    ("お見積りのご提示", 30, INK, True),
], align=PP_ALIGN.CENTER, gap=10)
textbox(s, Inches(4.5), Inches(4.25), Inches(8.4), Inches(0.9), [
    ("2026年7月16日付「貴社コンテンツIPご活用について」の", 13.5, INK, False),
    ("ご依頼事項・ご確認ポイントへのご回答", 13.5, INK, False),
], align=PP_ALIGN.CENTER, gap=4)
textbox(s, Inches(4.5), Inches(5.7), Inches(8.4), Inches(1.0), [
    ("2026年7月", 15, TSUB_NAVY, True),
    ("(株)円谷プロダクション", 15, TSUB_NAVY, True),
    ("金額はすべて税別・概算", 10.5, SILVER, False),
], align=PP_ALIGN.CENTER, gap=3)
textbox(s, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
        [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(10.25), Inches(7.1), Inches(2.5), Inches(0.3),
        [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)

# =============================================================================
# 2. ご依頼事項の整理
# =============================================================================
s = add_slide()
header(s, "ご依頼事項の整理", "貴社資料(2026/7/16)のご依頼・ご確認事項を以下の通り理解しております")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["区分", "内容", "貴社想定規模"],
          [
              ["① 着ぐるみイベント", "店頭写真撮影会(グリーティング型)。イベントキットと併用", "600開催/年(全国)"],
              ["② 通常イベントキット", "展示キット+ノベルティ(作成主体は貴社)", "20,000開催/年(全国)"],
              ["③ イベント外", "店頭訴求ツールへのIP活用(ポスター・POP等)", "各種ツールで展開"],
          ],
          col_widths=[1.1, 2.4, 1.1], row_h=0.55, font_size=12.5)
textbox(s, Inches(0.85), Inches(3.7), Inches(11.8), Inches(0.9), [
    [("ご依頼内容:", 14, RED, True),
     ("年間費用および内訳のご提示(IP利用/契約料、着ぐるみ費用、作成許諾費用〔ロイヤリティ〕)", 14, INK, True)],
], gap=4)
b = rect(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(1.75), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(1.45), [
    ("あわせてご回答する確認ポイント(見積もり以外)", 13, RED, True),
    ("1. FY26中のテスト実施可否(単体契約 or 短期契約 等)", 12.5, INK, False),
    ("2. ノベルティ作成時の費用(ロイヤリティ)発生有無", 12.5, INK, False),
    ("3. イベント実施会場に関する制限事項有無(更衣室の準備など)", 12.5, INK, False),
], gap=6)

# =============================================================================
# 3. 年間ご契約費用の全体構造
# =============================================================================
s = add_slide()
header(s, "年間ご契約費用の全体構造", "固定費は年間IP使用料のみ。着ぐるみ費用・ロイヤリティは実施数・作成数に連動する従量型です")
boxes = [
    ("① 年間IP使用料(固定)", "1,650万円/年", "イベント・店頭訴求ツール・Web/SNSでのIP利用と監修対応を包括"),
    ("② 着ぐるみイベント費用(従量)", "15万円/開催(土日2日間)", "撮影会・グリーティング型。実施数に応じたお支払い"),
    ("③ 作成許諾費用(従量)", "作成費用の10%", "イベントキット・ノベルティ・訴求ツール等のロイヤリティ"),
]
x0 = Inches(0.7)
for i, (t, price, desc) in enumerate(boxes):
    y = Inches(1.75 + i * 1.5)
    bx = rect(s, x0, y, Inches(11.9), Inches(1.3), fill=PALE if i != 0 else PALE_RED,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx.adjustments[0] = 0.08
    textbox(s, Inches(1.0), y + Inches(0.14), Inches(4.6), Inches(1.0), [
        (t, 14, RED, True),
        (desc, 11, SUB, False),
    ], gap=4)
    textbox(s, Inches(7.6), y, Inches(4.7), Inches(1.3),
            [(price, 22, INK, True)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
textbox(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5), [
    ("※ポケモン様の現行運用と同一の費用構造(年間契約+着ぐるみ都度+ロイヤリティ)に揃えており、貴社の運用フローを変えずに導入いただけます。", 11, SUB, False),
])

# =============================================================================
# 4. ①年間IP使用料
# =============================================================================
s = add_slide()
header(s, "①年間IP使用料:1,650万円(税別)", "貴社のご予算水準に合わせた年間包括のご提供条件です")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "条件"],
          [
              ["契約期間", "年度単位(12か月)"],
              ["使用業種", "携帯通信、ブロードバンド、でんき(コンシューマ向けプロモーション)"],
              ["使用媒体", "イベント(着ぐるみ・キット)、店頭訴求ツール、LP/Web、SNS、メール"],
              ["含まれるもの", "上記範囲でのIP利用許諾、作成物の監修対応、着ぐるみイベント年600開催の開催権(稼働実費は②)"],
              ["通常キットイベント", "年20,000開催規模での展開可(キット作成費は貴社、ロイヤリティは③)"],
              ["ステージショー(オプション)", "年10回程度まで対応可。1回30万円〜(内容により個別見積)"],
              ["範囲外", "TVCM、交通広告、大型OOH、販売用グッズ、法人向け、グループ他社サービス"],
          ],
          col_widths=[1.15, 2.6], row_h=0.5, font_size=12)
textbox(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.6), [
    ("※60周年イヤーの記念施策・キャンペーン素材のご提供など、初年度限定の特典もあわせてご相談させてください。", 11, SUB, False),
])

# =============================================================================
# 5. ②着ぐるみイベント費用
# =============================================================================
s = add_slide()
header(s, "②着ぐるみイベント費用(撮影会・グリーティング型)",
       "ポケモン様と同一の運用形式・同水準の費用に設定。貴社店舗・量販店・商業施設での週末稼働に対応します")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "内容"],
          [
              ["形式", "写真撮影会・グリーティング(1回30分×1日最大5回)。イベントキットと併用"],
              [("費用", True, RED), ("15万円/開催(土日2日間の概算)。1日のみの場合は8万円", True, INK)],
              ["内訳", "着ぐるみ費用+専属スタッフ人件費(ディレクター+アクターの2名体制)"],
              ["交通・宿泊", "首都圏は費用内。首都圏外は交通・宿泊実費を別途"],
              ["体制", "円谷プロ指定アクター制(貴社・代理店社員による着用は不可)。全国600開催に対応する供給体制を整備"],
              ["会場条件", "外部から見えない更衣室(控室)のご用意/他キャラクターとの同時共存は不可/非公式商品の展示・装飾利用は不可"],
          ],
          col_widths=[1, 3.2], row_h=0.55, font_size=12)
textbox(s, Inches(0.85), Inches(5.6), Inches(11.8), Inches(0.9), [
    [("参考:", 12.5, RED, True),
     ("年600開催(貴社想定)の場合、着ぐるみ費用は年間約9,000万円(15万円×600開催、実施数に連動)。", 12.5, INK, True)],
    ("開催数が変動した場合も費用は実施分のみ。固定のコミットメントはございません。", 11.5, SUB, False),
], gap=5)

# =============================================================================
# 6. ③作成許諾費用(ロイヤリティ)
# =============================================================================
s = add_slide()
header(s, "③作成許諾費用:作成費用の10%",
       "イベントキット・ノベルティ・店頭訴求ツールの作成主体は貴社。作成費用(税別)の10%をロイヤリティとしてお支払いいただきます")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "内容"],
          [
              [("料率", True, RED), ("作成費用(税別)の10%(ポケモン様と同率)", True, INK)],
              ["対象", "イベントキット(ルーレット・POP・パネル等)、ノベルティ、店頭訴求ツール(ポスター・POP・什器シート等)"],
              ["作成数", "上限なし。数量に応じたロイヤリティのお支払いのみ"],
              ["監修", "作成物ごとに事前監修(デザインデータ確認)+完成サンプルのご提出"],
              ["素材提供", "キービジュアル・キャラクター素材は当社より提供(既存素材は使用料内。新規描き起こしは個別見積)"],
              ["お支払い", "請求書発行月の翌月末までにお振込(作成実績ベース)"],
          ],
          col_widths=[1, 3.2], row_h=0.52, font_size=12)
textbox(s, Inches(0.85), Inches(5.5), Inches(11.8), Inches(0.7), [
    [("ご確認ポイント2への回答:", 12.5, RED, True),
     ("ノベルティ作成時のロイヤリティは発生します(作成費用の10%)。", 12.5, INK, True)],
])

# =============================================================================
# 7. 年間費用の概算シミュレーション
# =============================================================================
s = add_slide()
header(s, "年間費用の概算シミュレーション",
       "貴社想定の開催数(着ぐるみ600開催/キット20,000開催)をそのまま当てはめた場合の年間概算です")
add_table(s, Inches(0.7), Inches(1.65), Inches(11.9),
          ["費用項目", "単価・料率", "貴社想定規模", "年間概算"],
          [
              ["① 年間IP使用料(固定)", "—", "—", ("1,650万円", True, RED)],
              ["② 着ぐるみイベント", "15万円/開催(土日)", "600開催/年", "約9,000万円(従量)"],
              ["  ステージショー(任意)", "30万円〜/回", "年10回程度まで", "〜約300万円(任意)"],
              ["③ 作成許諾ロイヤリティ", "作成費用の10%", "キット・ノベルティ・訴求ツール", "作成費用に応じて(従量)"],
          ],
          col_widths=[1.5, 1.1, 1.3, 1.4], row_h=0.55, font_size=12)
b = rect(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(1.15), fill=PALE_RED)
textbox(s, Inches(1.0), Inches(4.78), Inches(11.3), Inches(0.85), [
    ("固定でお支払いいただくのは①1,650万円のみ。②③は実施・作成した分だけの従量型のため、", 13.5, INK, False),
    ("展開規模を貴社側でコントロールしながらご活用いただけます。", 14, RED, True),
], gap=4)
textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.6), [
    ("※金額はすべて税別・概算です。②は実施数、③は作成実績により変動します。訴求ツール作成のロイヤリティはイベント作成物と同水準(10%)です。", 10.5, SUB, False),
])

# =============================================================================
# 8. ご確認ポイントへのご回答
# =============================================================================
s = add_slide()
header(s, "ご確認ポイントへのご回答", "貴社資料の確認事項3点への回答です")
add_table(s, Inches(0.7), Inches(1.65), Inches(11.9),
          ["No.", "ご確認内容", "ご回答"],
          [
              ["1", "FY26中のテスト実施可否(単体契約 or 短期契約)",
               "短期・単体契約での対応は可能です。ただし当社の既存ライセンス条件との調整が必要なため、FY26中の実施は形態・会場・時期を個別協議とさせてください。本格開始は2027年4月(FY27期首)を推奨します"],
              ["2", "ノベルティ作成時の費用(ロイヤリティ)発生有無",
               "発生します。作成費用(税別)の10%をロイヤリティとしてお支払いいただきます(作成数の上限なし)"],
              ["3", "イベント実施会場に関する制限事項有無",
               "①外部から見えない更衣室(控室)のご用意 ②他キャラクターとの同時共存は不可 ③非公式商品の店頭ツール・装飾利用は不可。詳細は運用ガイドラインをご提示します"],
          ],
          col_widths=[0.35, 1.6, 2.9], row_h=0.85, font_size=11.5)
textbox(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.6), [
    ("※着ぐるみの稼働ルール(1回30分×1日最大5回、専属スタッフ2名体制)は、ポケモン様の現行運用と同一の水準で設計しています。", 11, SUB, False),
])

# =============================================================================
# 9. 次アクション
# =============================================================================
s = add_slide()
header(s, "本お見積りのポイントと次アクション", "ご予算・運用フローともに現行IP様の座組と同一の形でご導入いただけます")
b1 = rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(2.3), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b1.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.0), [
    ("本お見積りのポイント", 14, RED, True),
    ("・年間IP使用料1,650万円+従量型(着ぐるみ15万円/開催・ロイヤリティ10%)のシンプルな3階建て構造", 13, INK, False),
    ("・ポケモン様の現行運用(形式・体制・料率)と同一の枠組みのため、貴社・代理店様の運用フロー変更が不要", 13, INK, False),
    ("・60周年イヤーのウルトラマンで「会いに行ける」店頭集客コンテンツを、固定費を抑えて導入可能", 13, INK, False),
], gap=7)
add_table(s, Inches(0.7), Inches(4.3), Inches(11.9),
          ["時期", "アクション"],
          [
              ["2026年7月〜8月", "本お見積りのご確認・条件協議(契約書ドラフトのご提示)"],
              ["2026年9月〜12月", "キット・訴求ツールのデザイン監修、着ぐるみ供給体制の整備、FY26中テストの個別協議"],
              [("2027年4月", True, RED), ("年間契約開始・全国展開スタート(FY27)", True, RED)],
          ],
          col_widths=[1, 2.8], row_h=0.5, font_size=12.5)

out = "SoftBank様向け_お見積りご回答.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
