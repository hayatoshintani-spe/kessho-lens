# -*- coding: utf-8 -*-
"""プランA/B二者択一のポケットスライド(1枚・円谷フォーマット)を生成する。
SBが単価に抵抗した局面でのみ提示する追加資料。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
s = prs.slides.add_slide(prs.slide_layouts[6])


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=4):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


s.shapes.add_picture(os.path.join(ASSETS, "band_header.png"), 0, 0, width=SW, height=Inches(0.86))
tb(Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
   [("ご参考:2つのご契約プラン", 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
tb(Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
   [("単価の柔軟性と、お取り組み規模のご保証はトレードオフの関係です。ご運用に合う形をお選びください", 12.5, SUB, False)])

rows = [
    ("年間IP使用料", "2,000万円", "1,650万円(貴社ご希望水準)"),
    ("着ぐるみイベント(土日2日)", "20万円/開催", "15万円/開催(交通・宿泊実費別途)"),
    ("作成許諾ロイヤリティ", "作成費用の20%", "作成費用の10%"),
    ("開催コミット", "なし(完全従量)", "年間400開催の最低お取り組み保証"),
    ("作成物の最低数量", "なし(完全従量)", "キット600式・ノベルティ3,000セット/年(不足分は精算)"),
    ("契約期間", "単年(FY27)", "2年(FY27〜FY28)"),
]
tbl = s.shapes.add_table(len(rows) + 1, 3, Inches(0.7), Inches(1.6), Inches(11.9),
                         Inches(0.52 * (len(rows) + 1))).table
tbl.columns[0].width = int(Inches(11.9) * 0.28)
tbl.columns[1].width = int(Inches(11.9) * 0.32)
tbl.columns[2].width = int(Inches(11.9) * 0.40)
headers = ["項目", "プランA:フリープラン", "プランB:年間コミットプラン"]
for j, htxt in enumerate(headers):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = NAVY if j != 2 else RED
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    c.margin_left = c.margin_right = Inches(0.08)
    p = c.text_frame.paragraphs[0]
    r = p.add_run(); r.text = htxt
    set_font(r, 13, WHITE, True)
for i, (label, a, b) in enumerate(rows, start=1):
    for j, text in enumerate([label, a, b]):
        c = tbl.cell(i, j)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = text
        set_font(r, 12, INK, bold=(j == 2 and i >= 4))

bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.45), Inches(11.9), Inches(1.05))
bx.adjustments[0] = 0.07
bx.shadow.inherit = False
bx.fill.solid(); bx.fill.fore_color.rgb = PALE_RED
bx.line.fill.background()
tb(Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.8), [
    [("プランBの単価優遇は、", 12.5, INK, False),
     ("全国600開催に対応する専属アクター・スーツ供給体制への先行投資が前提", 12.5, RED, True),
     ("となるため、", 12.5, INK, False)],
    ("最低お取り組み規模のご保証(貴社ご計画600開催の3分の2)と複数年でのご契約をお願いするものです。", 12.5, INK, False),
], gap=4)
tb(Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.4), [
    ("※金額はすべて税別・概算です。中間的な設計(単価・コミット規模の組み合わせ)もご相談可能です。", 10.5, SUB, False),
])
tb(Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
   [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
tb(Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
   [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)

prs.save("プランAB比較_ポケットスライド.pptx")
print("saved")
