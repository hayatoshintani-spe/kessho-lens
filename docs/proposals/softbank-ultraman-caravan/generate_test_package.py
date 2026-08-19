# -*- coding: utf-8 -*-
"""完成版:タイトル→企画内容→見積(木曜SB打ち合わせ提出用・円谷フォーマット)"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=11.5, row_h=0.5):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(headers), x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, font_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            for k, line in enumerate(str(text).split("\n")):
                p = c.text_frame.paragraphs[0] if k == 0 else c.text_frame.add_paragraph()
                r = p.add_run(); r.text = line
                set_font(r, font_size, color, bold)
    return tbl



COVER_IMG = os.path.join(ASSETS, "band_cover.png")
TSUB_NAVY = RGBColor(0x1F, 0x33, 0x8C)

# =============================================================================
# 1. 表紙
# =============================================================================
s = add_slide()
s.shapes.add_picture(COVER_IMG, 0, 0, width=Inches(4.09), height=SH)
textbox(s, Inches(4.5), Inches(1.25), Inches(8.4), Inches(0.5),
        [("ソフトバンク株式会社 御中", 14, SUB, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(4.5), Inches(2.35), Inches(8.4), Inches(1.8), [
    ("スポットテスト実施に向けた", 28, INK, True),
    ("ご準備資料", 28, INK, True),
], align=PP_ALIGN.CENTER, gap=10)
textbox(s, Inches(4.5), Inches(4.15), Inches(8.4), Inches(1.0), [
    ("7/23お打ち合わせの宿題事項へのご回答", 14, INK, False),
    ("(実施可能エリア/ノベルティ/テストイベント形式/検証項目案)", 12.5, SUB, False),
], align=PP_ALIGN.CENTER, gap=4)
textbox(s, Inches(4.5), Inches(5.7), Inches(8.4), Inches(1.0), [
    ("2026年7月31日", 15, TSUB_NAVY, True),
    ("(株)円谷プロダクション", 15, TSUB_NAVY, True),
    ("金額はすべて税別・概算", 10.5, SILVER, False),
], align=PP_ALIGN.CENTER, gap=3)
textbox(s, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
        [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(10.25), Inches(7.1), Inches(2.5), Inches(0.3),
        [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)

# =============================================================================
# 2. 実施可能エリア(宿題①)
# =============================================================================
s = add_slide()
header(s, "① テスト期間の実施可能エリア",
       "当社の既存イベント運用との調整のため、テスト期間は以下のエリアでの実施をお願いいたします")
b = rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.7), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4), [
    ("実施可能エリア", 14, RED, True),
    [("神奈川県・東京都", 20, INK, True), ("(埼玉県隣接地域を除く) を中心にご調整をお願いします", 14, INK, False)],
    ("※上記以外のエリアをご希望の場合は個別にご相談ください。会場候補をいただければ、当社にて開催可否を速やかに確認・回答いたします", 11.5, SUB, False),
], gap=7)
add_table(s, Inches(0.7), Inches(3.65), Inches(11.9),
          ["確認項目", "条件", "備考"],
          [
              ["エリアバッティング", "イベント日の前8日間・後5日間、半径8km圏内に他のウルトラマンイベントがないこと", "競合施設でない場合・キャラクター違いは調整のうえ実施できる場合があります"],
              ["会場条件", "外部から見えない控室(着替え場所)・歩いて移動できる動線", "詳細は別紙レギュレーションご参照"],
              ["スケジュール", "会場・日程の確定は開催の1か月前まで(仮押さえは2か月前目安)", "9月開催をご希望の場合はお早めの候補共有をお願いします"],
          ],
          col_widths=[1.0, 2.6, 1.7], row_h=0.56, font_size=10.5)
textbox(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.7), [
    [("お願い:", 13, RED, True),
     ("商業施設の候補リストをご共有ください。当社にてバッティング・会場条件を確認し、開催可否を即日〜数営業日で回答します。", 13, INK, False)],
], gap=4)

# =============================================================================
# 3. ノベルティのご提案(宿題②)
# =============================================================================
s = add_slide()
header(s, "② ノベルティ・賞品のご提案(既存品より)",
       "参加賞から成約者特典まで、当社既存品でご用意できる4段構成をご提案します(詳細は別添の商品リストをご参照ください)")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["用途", "商品", "税別売価", "ポイント"],
          [
              [("参加賞", True), "ウルトラマンカードゲーム ブースターパック(複数弾)", "180〜200円/パック", "1BOX=24パック入。登場キャラが弾ごとに異なり、お客様が好きな弾を選べます"],
              [("抽選・ルーレット上位賞", True), "ウルトラヒーローソフビシリーズ各種(スペシャルカラー)", "818〜1,100円", "一般玩具売場に出回らない限定カラー。登場ヒーローと合わせた選定が可能です"],
              [("プレミア賞・ご成約特典", True), "フォンタブ&ショルダーストラップ/ネッククールリング", "1,800〜2,000円", "スマートフォン関連・夏の実用アイテム。店頭のご成約特典にも好適です"],
              [("+α(無償ご提供)", True), "アプリゲーム「パズルシュワッチ!!」販促ステッカー", ("無償", True, RED), "数量はご相談ください"],
          ],
          col_widths=[1.1, 2.0, 0.9, 2.2], row_h=0.66, font_size=10.5)
b = rect(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.3), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.1), [
    ("構成例(5会場・各300名想定)", 13, RED, True),
    ("・参加賞:ブースターパック 1,500パック(300パック×5会場) ・上位賞:ソフビ 各会場20体 ・プレミア賞:各会場10個", 12, INK, False),
    ("・数量・構成は会場規模に合わせて調整可能です。賞品は当社にてご用意し、事前に会場へお届けします", 12, INK, False),
], gap=5)
textbox(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.45), [
    ("※在庫状況により一部商品はご用意できない場合があります。※商品リスト(別添Excel)に写真・仕様・在庫目安を記載しています。", 10, SUB, False),
], gap=3)

# =============================================================================
# 4. テストイベント3形式(宿題③・全体)
# =============================================================================
s = add_slide()
header(s, "③ テストイベントの実施形式(3形式)",
       "着ぐるみを用いる2形式に加え、当社既存キットを活用する形式も今期からご提案可能です")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "S1:写真撮影会型", "S2:ステージショー型", "S3:イベントキット型"],
          [
              ["内容", "着ぐるみ撮影会・グリーティング(1回30分×1日最大5回)", "2ヒーローショー 1日2回(各約30分)", "当社保有の体験ゲームキット+景品(着ぐるみなし)"],
              ["単位", "1開催=土日2日間", "1開催=土日2日間", "1日単位"],
              [("費用", True, RED), ("100万円/開催", True, INK), ("200万円/開催", True, INK), ("約29万円/日(2日で約58万円)+交通運搬実費", True, INK)],
              ["含まれるもの", "IP利用・監修・スーツアクター・運搬", "9名編成・音響機材・テント・運営5名・諸経費・交通運搬", "キット一式・運営スタッフ4名・諸経費(テント要否別途3万円/日)"],
              ["貴社側ご負担", "会場手配・集客告知・当日誘導", "同左+ステージ(会場既設 or 施工50〜80万円〜)", "会場手配・集客告知"],
          ],
          col_widths=[0.75, 1.55, 1.55, 1.55], row_h=0.62, font_size=10.5)
textbox(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.3), [
    [("推奨プラン:", 12.5, RED, True),
     ("エリア・施設タイプの異なる5会場でS1を実施(総額500万円・IP利用込み)+1〜2会場でS2またはS3をオプション実施。", 12.5, INK, False)],
    ("※S3は当社の既存保有キットを活用するため今期実施が可能です(貴社仕様での新規キット製作を伴うテストは来年度からのご提案)。", 11, SUB, False),
    ("※S2のショー終了後の撮影会の組み込みは、会場条件により個別にご相談させてください。※いずれも駐車場・スタッフ控室のご用意を会場様とご相談ください。", 10.5, SUB, False),
], gap=4)

# =============================================================================
# 5. S3:イベントキット型の詳細
# =============================================================================
s = add_slide()
header(s, "③-1 イベントキット型(S3)の内容と費用",
       "当社保有の体験ゲームキットを活用した店頭・商業施設イベントです(1日単位・着ぐるみなし)")
textbox(s, Inches(0.7), Inches(1.55), Inches(5.9), Inches(3.3), [
    ("キット構成", 13.5, RED, True),
    ("・体験ゲーム(3種の内1種)+装飾アイテムセット", 12, INK, False),
    ("・ゲーム説明用パネル(1枚)", 12, INK, False),
    ("・テーブルクロス(ビニールカバー付き):3セット", 12, INK, False),
    ("・景品:「修了証」ステッカー 200枚", 12, INK, False),
    ("・会場装飾用バナー:いずれか1種", 12, INK, False),
    ("", 8, INK, False),
    ("装飾バナーについて", 13.5, RED, True),
    ("・バナーは骨組み付きで、当社運営スタッフが設営いたします", 11.5, INK, False),
    ("・テント用横幕は骨組みがないため、固定できる場所がある場合に設置可能です", 11.5, INK, False),
], gap=5)
add_table(s, Inches(6.9), Inches(1.6), Inches(5.7),
          ["項目", "費用(税別)"],
          [
              ["① キット一式", "15万円/1日"],
              ["② 運営スタッフ4名(ディレクター1名+スタッフ3名)", "9万円/1日"],
              ["③ 現場諸経費", "2万円/1日"],
              ["④ 交通運搬費(車両・ガソリン・高速・ツール運搬)", "実費(会場確定後お見積り)"],
              ["⑤ テント(必要な場合)", "3万円/1日"],
              [("合計目安(テント含む)", True), ("約29万円/1日+④実費", True, RED)],
          ],
          col_widths=[2.3, 1.4], row_h=0.6, font_size=11)
textbox(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.5), [
    ("※土日2日間の実施の場合:約58万円+交通運搬実費が目安です。※駐車場・スタッフ控室のご用意を会場様とご相談ください。", 10.5, SUB, False),
], gap=3)

# =============================================================================
# 6. S2:ステージショー型の詳細
# =============================================================================
s = add_slide()
header(s, "③-2 ステージショー型(S2)の内容と費用",
       "2ヒーローショーのフルパッケージです(200万円/開催・土日2日間)")
textbox(s, Inches(0.7), Inches(1.55), Inches(5.9), Inches(4.3), [
    ("200万円に含まれるもの", 13.5, RED, True),
    ("・2ヒーローショーのパッケージ:1日2回(11時〜/14時〜推奨・各約30分)", 12, INK, False),
    ("・出演:ウルトラマン2体/怪獣3〜4体", 12, INK, False),
    ("・スタッフ:チーフ1名+アクター6名+MC+PAオペレーター(計9名)", 12, INK, False),
    ("・音響機材・テント", 12, INK, False),
    ("・運営スタッフ5名(ディレクター1名+スタッフ4名)", 12, INK, False),
    ("・現場諸経費・交通運搬費", 12, INK, False),
], gap=5)
textbox(s, Inches(6.9), Inches(1.55), Inches(5.7), Inches(4.5), [
    ("ステージ・会場について", 13.5, RED, True),
    ("・会場様のお持ちのステージ、または会場ご指定の施工会社様の手配が一般的です。ステージなしでの実施も可能です", 11.5, INK, False),
    ("・ステージを施工する場合:約50〜80万円(税別)〜が目安(設営時間・時間帯により変動)", 11.5, INK, False),
    ("・必要スペース:幅6m×奥行4m程度+左右いずれかに着替えスペース(テント等)+音響機材置き場(長テーブル1本)", 11.5, INK, False),
    ("・お客様との距離は1m以上を確保してください", 11.5, INK, False),
    ("", 6, INK, False),
    ("その他", 13.5, RED, True),
    ("・BGM等のJASRAC申請は主催者様側でのお手続きをお願いします(会場様が包括契約をお持ちの場合もあるためご確認ください)", 11.5, INK, False),
    ("・会場様ご指定のイベント運営会社様がいらっしゃる場合は連携いたします", 11.5, INK, False),
], gap=5)

# =============================================================================
# 7. 検証項目のたたき台(共同策定用)
# =============================================================================
s = add_slide()
header(s, "ご参考:検証項目のたたき台",
       "「検証項目の共同策定」に向けた当社案です。次回お打ち合わせでご一緒に固めさせてください")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["区分", "指標(例)", "取得方法(例)"],
          [
              ["集客力", "イベント来場者数/撮影会参加組数(回別)/待機列の発生状況", "当社スタッフの現場カウント+貴社会場データ"],
              ["送客力", "イベント→ブース・店頭への送客数/クーポン・ノベルティ引換数", "引換券・QR等の設計を共同で検討"],
              ["商談効果", "ブース来訪からの相談・商談件数/ご成約数(可能な範囲で)", "貴社側データ(店舗・出店ブース)"],
              ["話題化", "SNS投稿数・撮影画像の投稿率/会場様の反応", "ハッシュタグ設計+会場ヒアリング"],
              ["運営品質", "オペレーション所要人員・動線・安全面の課題抽出", "双方の現場レビュー(実施後1週間以内)"],
          ],
          col_widths=[0.8, 2.3, 1.9], row_h=0.6, font_size=11)
textbox(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.1), [
    [("ご提案:", 13, RED, True),
     ("5会場をエリア・施設タイプで分散し、上記指標を共通フォーマットで比較検証 → 年間のお取り組み設計(展開エリア・開催数)の根拠データにします。", 13, INK, False)],
    ("※検証項目・取得方法は貴社のご意向に合わせて調整いたします。", 11, SUB, False),
], gap=5)

# =============================================================================
# 8. 次のステップ
# =============================================================================
s = add_slide()
header(s, "次のステップ", "9月の実施に向けて、以下の進行をご提案します")
steps = [
    ("1  会場候補のご共有(貴社)", "商業施設の候補リストをご共有ください。当社にて開催可否を確認します"),
    ("2  会場・日程の確定", "エリア・バッティング確認のうえ、実施会場・日程を確定(9月開催の場合は8月中旬目安)"),
    ("3  検証項目の合意", "本資料P7のたたき台をベースに、次回お打ち合わせで共同策定"),
    ("4  実施準備", "告知物の監修・実施確認書・ノベルティ手配・当日運営の役割分担"),
]
yy = 1.7
for t, d in steps:
    bb = rect(s, Inches(0.7), Inches(yy), Inches(11.9), Inches(1.0), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bb.adjustments[0] = 0.12
    textbox(s, Inches(1.05), Inches(yy + 0.13), Inches(11.2), Inches(0.8),
            [(t, 13.5, NAVY, True), (d, 11.5, SUB, False)], gap=3)
    yy += 1.18
textbox(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5), [
    ("ご不明点は何なりとお申し付けください。会場候補をいただければ、開催可否とあわせて具体的なお見積りをご提示します。", 12, INK, False),
], gap=3)

out = "SB様送付用_テスト実施に向けたご準備資料.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
