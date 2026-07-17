# -*- coding: utf-8 -*-
"""規模別お見積り3パターン+今期スポットテスト提案(7/23宿題・円谷フォーマット)"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=11.5, row_h=0.5):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(headers), x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, font_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = text
            set_font(r, font_size, color, bold)
    return tbl


# =============================================================================
# 1. 規模別お見積り3パターン
# =============================================================================
s = add_slide()
header(s, "規模別お見積り:3パターン(年間・税別概算)",
       "7/17お打ち合わせを受け、全量・半分・3IP分割(約1/3)の3規模でご提示します(単価は共通:IP使用料2,000万円/着ぐるみ20万円/ロイヤリティ20%)")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["パターン", "着ぐるみ", "キット・ノベルティ想定", "①IP使用料(固定)", "②着ぐるみ費用", "③ロイヤリティ", "年間合計"],
          [
              [("最大(貴社実績同等・全量)", True), "600開催", "キット800式・ノベルティ2.1万セット(実績)・ツール全店2回", "2,000万円", "1億2,000万円", "約6,480万円", ("約2.05億円", True, RED)],
              [("半分(実績の50%)", True), "300開催", "キット400式・ノベルティ1.05万セット・ツール全店1回", "2,000万円", "6,000万円", "約3,240万円", ("約1.12億円", True, RED)],
              [("3IP分割(約1/3配分)", True), "200開催", "キット270式・ノベルティ7,000セット・ツール1/3", "2,000万円", "4,000万円", "約2,170万円", ("約8,200万円", True, RED)],
          ],
          col_widths=[1.15, 0.65, 1.85, 0.95, 0.95, 0.85, 0.95], row_h=0.72, font_size=10.5)
b = rect(s, Inches(0.7), Inches(4.65), Inches(11.9), Inches(1.45), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.25), [
    ("お打ち合わせを受けたご提案", 13, RED, True),
    ("・年間の最低開催数・最低作成数をお約束いただける場合、着ぐるみ単価を優遇します(例:400開催以上のコミットで17万円/開催)", 12, INK, False),
    ("・ウルトラマンは7・8月も稼働可能です(現行IP様は夏季稼働不可)。夏商戦を含むフル12か月稼働なら最大720開催まで拡張できます", 12, INK, False),
    ("・複数IPでのご運用でも、映画・新TVシリーズ・玩具新商品と連動した季節企画で「切替の変化」をご提供できます(キャラクター切替・選べるヒーロー等)", 12, INK, False),
], gap=5)
textbox(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.8), [
    ("※③はキット単価15.5万円/式、ノベルティは貴社FY25ご実績(年約2.1万セット・作成費約1億8,500万円)を配分比で試算。ロイヤリティは作成実績(実額)ベースで精算します。", 10, SUB, False),
    ("※店頭訴求ツールの年間作成数(貴社ご確認中)の受領後、③をさらに精緻化します。※固定費は①のみ。②③は実施・作成分だけの従量です。", 10, SUB, False),
], gap=3)

# =============================================================================
# 2. 今期スポットテストのご提案
# =============================================================================
s = add_slide()
header(s, "今期スポットテストのご提案(回数ベース・地域自由)",
       "年間のお取り組みとは切り離した単発テストです。着ぐるみを用いる2形式をご用意しました")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "形式S1:写真撮影会型(グリーティング)", "形式S2:ステージショー型"],
          [
              ["内容", "1回30分×1日最大5回の撮影会・グリーティング。店頭・商業施設の集客と送客を検証", "ステージショー(2回/日)+撮影会のセット。大型集客の瞬発力を検証"],
              ["単位", "1開催=土日2日間", "1開催=土日2日間"],
              [("費用", True, RED), ("100万円/開催(IP許諾・監修・スーツ・アクター・輸送込み)", True, INK), ("150万円/開催(ショー演出・音源・追加スタッフ込み)", True, INK)],
              ["貴社側ご負担", "会場手配・集客告知・当日運営(誘導スタッフ)", "同左+ステージ設営"],
              ["想定規模", ("5施設程度(各1開催)を想定 → 総額500万円", True, INK), "うち1〜2施設でのオプション実施"],
              ["地域", "全国対応可(関東に限定しません)", "同左"],
          ],
          col_widths=[0.6, 1.7, 1.7], row_h=0.56, font_size=11)
textbox(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.3), [
    [("想定プラン:", 12.5, RED, True),
     ("エリア・施設タイプの異なる5施設で実施し、集客・送客効果を比較検証(S1×5施設=総額500万円・IP使用料不要)。", 12.5, INK, False)],
    [("継続パック(S1):", 12.5, RED, True),
     ("隔週(月2開催)190万円/月、毎週(月4開催)360万円/月(単発比の優遇設定)。", 12.5, INK, False)],
    [("実施条件:", 12.5, RED, True),
     ("当社の既存ライセンス条件の調整を経て、実施時期・会場を確定します(調整は当社にて速やかに進めます)。測定項目:来場数・撮影参加組数・ブース送客数。", 12, INK, False)],
    ("※通常イベントキットを用いるテストはキット製作を伴うため来年度からのご提案。※本テスト費用は、年間ご契約成立時に翌年度IP使用料へ充当するご相談が可能です。", 10.5, SUB, False),
], gap=3)

out = "規模別お見積り3パターンと今期スポットテスト.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
