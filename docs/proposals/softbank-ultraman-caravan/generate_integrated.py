# -*- coding: utf-8 -*-
"""SoftBank様ご依頼事項(2026/7/16)への見積回答資料を生成する(円谷フォーマット)。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
TSUB_NAVY = RGBColor(0x1F, 0x33, 0x8C)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")
COVER_IMG = os.path.join(ASSETS, "band_cover.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=12,
              header_size=12, row_h=0.45):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, header_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = text
            set_font(r, font_size, color, bold)
    return tbl


# =============================================================================
# 1. 表紙
# =============================================================================
s = add_slide()
s.shapes.add_picture(COVER_IMG, 0, 0, width=Inches(4.09), height=SH)
textbox(s, Inches(4.5), Inches(1.3), Inches(8.4), Inches(0.5),
        [("ソフトバンク株式会社 御中", 14, SUB, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(4.5), Inches(2.3), Inches(8.4), Inches(2.0), [
    ("ウルトラマンIP 導入のご提案", 30, INK, True),
    ("〜 テストから年間展開・複数年まで 〜", 22, INK, True),
], align=PP_ALIGN.CENTER, gap=10)
textbox(s, Inches(4.5), Inches(4.25), Inches(8.4), Inches(0.9), [
    ("お見積り(年間ご契約条件)・テスト実施パターン・年間展開への接続まで、", 13.5, INK, False),
    ("段階導入の全体像を1冊に統合したご提案です", 13.5, INK, False),
], align=PP_ALIGN.CENTER, gap=4)
textbox(s, Inches(4.5), Inches(5.7), Inches(8.4), Inches(1.0), [
    ("2026年7月", 15, TSUB_NAVY, True),
    ("(株)円谷プロダクション", 15, TSUB_NAVY, True),
    ("金額はすべて税別・概算", 10.5, SILVER, False),
], align=PP_ALIGN.CENTER, gap=3)
textbox(s, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
        [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(10.25), Inches(7.1), Inches(2.5), Inches(0.3),
        [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)

# =============================================================================
# 0b. 全体像:テストから年間展開までの一本道
# =============================================================================
s = add_slide()
header(s, "全体像:小さく始めて、段階的に全国へ",
       "各段階の意思決定は小さく、投資は次の段階に充当・接続されます。後戻りも可能な設計です")
journey = [
    ("STEP1|2026年9月〜", "先行開発テスト(無償):デザイン開発・監修・仕組み検証", PALE),
    ("STEP2|2026年11月〜2027年1月", "協賛型イベント検証(300〜500万円):集客・送客の実地データ取得", PALE),
    ("STEP3|2027年4月〜6月", "期首パイロット(短期IP使用料300万円+着ぐるみ実費):店頭運用を実測 ※年間契約移行時は300万円を全額充当", PALE),
    ("STEP4|2027年7月〜", "年間契約発効:成功基準達成で自動接続 → 夏商戦100開催 → 年末までに全国300〜600開催ペース", PALE_RED),
    ("STEP5|2028年4月〜", "FY28フル年間展開+複数年契約のご協議", None),
]
y = Inches(1.65)
for i, (period, desc, fill) in enumerate(journey):
    is_last = i == len(journey) - 1
    bx = rect(s, Inches(0.9), y, Inches(11.5), Inches(0.78), fill=(NAVY if is_last else fill), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx.adjustments[0] = 0.11
    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = period + "  "
    set_font(r1, 12.5, WHITE if is_last else RED, True)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, 12.5, WHITE if is_last else INK, is_last)
    y += Inches(0.92)
    if not is_last:
        ar = rect(s, Inches(6.5), y - Inches(0.17), Inches(0.22), Inches(0.14), fill=RED, shape=MSO_SHAPE.DOWN_ARROW)
textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.55), [
    [("設計思想:", 12.5, RED, True),
     ("STEP1〜3の投資は年間契約に充当・再利用され、無駄になりません。成功基準は貴社と共同設定し、達成時の年間移行を事前合意(LOI)いただきます。", 12.5, INK, False)],
])

# =============================================================================
# 2. ご依頼事項の整理
# =============================================================================
s = add_slide()
header(s, "ご依頼事項の整理", "貴社資料(2026/7/16)のご依頼・ご確認事項を以下の通り理解しております")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["区分", "内容", "貴社想定規模"],
          [
              ["① 着ぐるみイベント", "店頭写真撮影会(グリーティング型)。イベントキットと併用", "600開催/年(全国)"],
              ["② 通常イベントキット", "展示キット+ノベルティ(作成主体は貴社)", "20,000開催/年(全国)"],
              ["③ イベント外", "店頭訴求ツールへのIP活用(ポスター・POP等)", "各種ツールで展開"],
          ],
          col_widths=[1.1, 2.4, 1.1], row_h=0.55, font_size=12.5)
textbox(s, Inches(0.85), Inches(3.7), Inches(11.8), Inches(0.9), [
    [("ご依頼内容:", 14, RED, True),
     ("年間費用および内訳のご提示(IP利用/契約料、着ぐるみ費用、作成許諾費用〔ロイヤリティ〕)", 14, INK, True)],
], gap=4)
b = rect(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(1.75), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.06
textbox(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(1.45), [
    ("あわせてご回答する確認ポイント(見積もり以外)", 13, RED, True),
    ("1. FY26中のテスト実施可否(単体契約 or 短期契約 等)", 12.5, INK, False),
    ("2. ノベルティ作成時の費用(ロイヤリティ)発生有無", 12.5, INK, False),
    ("3. イベント実施会場に関する制限事項有無(更衣室の準備など)", 12.5, INK, False),
], gap=6)

# =============================================================================
# 3. 年間ご契約費用の全体構造
# =============================================================================
s = add_slide()
header(s, "年間ご契約費用の全体構造", "固定費は年間IP使用料のみ。着ぐるみ費用・ロイヤリティは実施数・作成数に連動する従量型です")
boxes = [
    ("① 年間IP使用料(固定)", "2,000万円/年", "イベント・店頭訴求ツール・Web/SNSでのIP利用と監修対応を包括"),
    ("② 着ぐるみイベント費用(従量)", "20万円/開催(土日2日間)", "撮影会・グリーティング型。実施数に応じたお支払い"),
    ("③ 作成許諾費用(従量)", "作成費用の20%", "イベントキット・ノベルティ・訴求ツール等のロイヤリティ"),
]
x0 = Inches(0.7)
for i, (t, price, desc) in enumerate(boxes):
    y = Inches(1.75 + i * 1.5)
    bx = rect(s, x0, y, Inches(11.9), Inches(1.3), fill=PALE if i != 0 else PALE_RED,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx.adjustments[0] = 0.08
    textbox(s, Inches(1.0), y + Inches(0.14), Inches(4.6), Inches(1.0), [
        (t, 14, RED, True),
        (desc, 11, SUB, False),
    ], gap=4)
    textbox(s, Inches(7.6), y, Inches(4.7), Inches(1.3),
            [(price, 22, INK, True)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
textbox(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5), [
    ("※ポケモン様の現行運用と同一の費用構造(年間契約+着ぐるみ都度+ロイヤリティ)に揃えており、貴社の運用フローを変えずに導入いただけます。", 11, SUB, False),
])

# =============================================================================
# 4. ①年間IP使用料
# =============================================================================
s = add_slide()
header(s, "①年間IP使用料:2,000万円(税別)", "60周年イヤーのプレミアムと、通信業種での優先的なお取り組みを含む年間包括のご提供条件です")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "条件"],
          [
              ["契約期間", "年度単位(12か月)"],
              ["使用業種", "携帯通信、ブロードバンド、でんき(コンシューマ向けプロモーション)"],
              ["使用媒体", "イベント(着ぐるみ・キット)、店頭訴求ツール、LP/Web、SNS、メール"],
              ["含まれるもの", "上記範囲でのIP利用許諾、作成物の監修対応、着ぐるみイベント年600開催の開催権(稼働実費は②)"],
              ["通常キットイベント", "年20,000開催規模での展開可(キット作成費は貴社、ロイヤリティは③)"],
              ["ステージショー(オプション)", "年10回程度まで対応可。1回30万円〜(内容により個別見積)"],
              ["範囲外", "TVCM、交通広告、大型OOH、販売用グッズ、法人向け、グループ他社サービス"],
          ],
          col_widths=[1.15, 2.6], row_h=0.5, font_size=12)
textbox(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.6), [
    ("※60周年イヤーの記念施策・キャンペーン素材のご提供など、初年度限定の特典もあわせてご相談させてください。", 11, SUB, False),
])

# =============================================================================
# 5. ②着ぐるみイベント費用
# =============================================================================
s = add_slide()
header(s, "②着ぐるみイベント費用(撮影会・グリーティング型)",
       "ポケモン様と同一の運用形式。専属スーツアクターによる造形・演技品質でご提供し、貴社店舗・量販店・商業施設での週末稼働に対応します")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "内容"],
          [
              ["形式", "写真撮影会・グリーティング(1回30分×1日最大5回)。イベントキットと併用"],
              [("費用", True, RED), ("20万円/開催(土日2日間の概算)。1日のみの場合は10万円", True, INK)],
              ["内訳", "ヒーロースーツ費用+専属スタッフ人件費(ディレクター+スーツアクターの2名体制)。ポージング・アクション等の演技対応込み"],
              ["交通・宿泊", "首都圏は費用内。首都圏外は交通・宿泊実費を別途"],
              ["体制", "円谷プロ指定アクター制(貴社・代理店社員による着用は不可)。全国600開催に対応する供給体制を整備"],
              ["会場条件", "外部から見えない更衣室(控室)のご用意/他キャラクターとの同時共存は不可/非公式商品の展示・装飾利用は不可"],
          ],
          col_widths=[1, 3.2], row_h=0.55, font_size=12)
textbox(s, Inches(0.85), Inches(5.6), Inches(11.8), Inches(0.9), [
    [("参考:", 12.5, RED, True),
     ("年600開催(貴社想定)の場合、着ぐるみ費用は年間約1億2,000万円(20万円×600開催、実施数に連動)。", 12.5, INK, True)],
    ("開催数が変動した場合も費用は実施分のみ。固定のコミットメントはございません。", 11.5, SUB, False),
], gap=5)

# =============================================================================
# 6. ③作成許諾費用(ロイヤリティ)
# =============================================================================
s = add_slide()
header(s, "③作成許諾費用:作成費用の20%",
       "イベントキット・ノベルティ・店頭訴求ツールの作成主体は貴社。作成費用(税別)の20%をロイヤリティとしてお支払いいただきます")
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
          ["項目", "内容"],
          [
              [("料率", True, RED), ("作成費用(税別)の20%", True, INK)],
              ["対象", "イベントキット(ルーレット・POP・パネル等)、ノベルティ、店頭訴求ツール(ポスター・POP・什器シート等)"],
              ["作成数", "上限なし。数量に応じたロイヤリティのお支払いのみ"],
              ["監修", "作成物ごとに事前監修(デザインデータ確認)+完成サンプルのご提出"],
              ["素材提供", "キービジュアル・キャラクター素材は当社より提供(既存素材は使用料内。新規描き起こしは個別見積)"],
              ["お支払い", "請求書発行月の翌月末までにお振込(作成実績ベース)"],
          ],
          col_widths=[1, 3.2], row_h=0.52, font_size=12)
textbox(s, Inches(0.85), Inches(5.5), Inches(11.8), Inches(0.7), [
    [("ご確認ポイント2への回答:", 12.5, RED, True),
     ("ノベルティ作成時のロイヤリティは発生します(作成費用の20%)。", 12.5, INK, True)],
])

# =============================================================================
# 6b. ご参考:作成物の想定数量とロイヤリティ概算
# =============================================================================
s = add_slide()
header(s, "ご参考:作成物の想定数量とロイヤリティ概算",
       "キット数は貴社ご提示の800キット、その他は貴社資料の前提に基づく試算です")
add_table(s, Inches(0.7), Inches(1.65), Inches(11.9),
          ["作成物", "想定数量(年間)", "単価目安(貴社資料より)", "作成費概算", "ロイヤリティ(20%)"],
          [
              ["イベントキット(ゲームキット一式)", "800キット(貴社ご提示の数量)", "約15.5万円/キット", "約1億2,400万円", "約2,480万円"],
              ["ノベルティ", "4,000セット(1セット=A〜D賞+参加賞327個、計約130万個)", "約9,900円/セット", "約4,000万円", "約800万円"],
              ["店頭訴求ツール(ポスター・POP・什器シート等)", "約3,000店×年2回更新", "約2,500円/店・回", "約1,500万円", "約300万円"],
              [("合計", True), "—", "—", ("約1.8億円", True), ("約3,570万円", True, RED)],
          ],
          col_widths=[1.5, 1.7, 1.1, 0.8, 0.9], row_h=0.62, font_size=11)
textbox(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.8), [
    ("※数量・単価は貴社資料(現行IP実績のキット構成・単価、ノベルティ4,000セット時の単価前提)に基づく参考値です。実際の作成計画に応じて変動し、お支払いは作成実績ベースの従量となります。", 11, SUB, False),
    ("※キット数は貴社ご提示の800キットを反映しています。ノベルティ・訴求ツールは貴社資料の前提に基づく想定値です。", 11, SUB, False),
], gap=5)

# =============================================================================
# 7. 年間費用の概算シミュレーション
# =============================================================================
s = add_slide()
header(s, "年間費用の概算シミュレーション",
       "貴社想定の開催数(着ぐるみ600開催/キット20,000開催)をそのまま当てはめた場合の年間概算です")
add_table(s, Inches(0.7), Inches(1.65), Inches(11.9),
          ["費用項目", "単価・料率", "貴社想定規模", "年間概算"],
          [
              ["① 年間IP使用料(固定)", "—", "—", ("2,000万円", True, RED)],
              ["② 着ぐるみイベント", "20万円/開催(土日)", "600開催/年", "約1億2,000万円(従量)"],
              ["  ステージショー(任意)", "30万円〜/回", "年10回程度まで", "〜約300万円(任意)"],
              ["③ 作成許諾ロイヤリティ", "作成費用の20%", "キット800・ノベルティ4,000セット・訴求ツール約3,000店×2回", "約3,570万円(従量)"],
              [("参考:年間合計(上記想定時)", True), "—", "—", ("約1.76億円", True, RED)],
          ],
          col_widths=[1.5, 1.1, 1.3, 1.4], row_h=0.55, font_size=12)
b = rect(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.1), fill=PALE_RED)
textbox(s, Inches(1.0), Inches(5.32), Inches(11.3), Inches(0.85), [
    ("固定でお支払いいただくのは①2,000万円のみ。②③は実施・作成した分だけの従量型のため、", 13.5, INK, False),
    ("展開規模を貴社側でコントロールしながらご活用いただけます。", 14, RED, True),
], gap=4)
textbox(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6), [
    ("※金額はすべて税別・概算です。②は実施数、③は作成実績により変動します。訴求ツール作成のロイヤリティはイベント作成物と同水準(20%)です。", 10.5, SUB, False),
])

# =============================================================================
# 1. テスト実施のご提案パターン
# =============================================================================
s = add_slide()
header(s, "テスト実施のご提案パターン(確認ポイント1へのご回答)",
       "既存ライセンス条件を踏まえ、2026年度内から段階的に検証を始められる3つのパターンをご用意しました")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "パターン1:先行開発テスト", "パターン2:イベント協賛型テスト", "パターン3:期首パイロット"],
          [
              ["時期", "2026年9月〜(即時開始可)", "2026年秋〜冬", "2027年4月〜6月"],
              ["内容", "キット・訴求ツールのデザイン開発と監修、お客様調査、デジタル施策(QR等)の仕組み検証 ※IP実働なし",
               "当社主催(または商業施設主催)の60周年イベントに貴社がご協賛・ブース出展。着ぐるみは当社イベントとして稼働",
               "首都圏5〜10店舗・着ぐるみ10〜20開催の単体テスト(全国展開と同一の運用形式)"],
              ["ご契約形態", "秘密保持契約+監修合意", "イベント協賛契約", "短期単体契約(3か月)"],
              ["費用目安", "原則無償(制作実費のみ)", "ご協賛金 300〜500万円", "短期IP使用料300万円+着ぐるみ20万円/開催+作成物ロイヤリティ"],
              ["検証できること", "クリエイティブ品質・監修フロー・運用設計", "集客力・ブース送客・相談転換の実地データ", "全国展開前の店頭オペレーション・KPI実測"],
          ],
          col_widths=[0.55, 1.3, 1.45, 1.3], row_h=0.72, font_size=10.5)
textbox(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.7), [
    [("※パターン3の短期IP使用料は、年間ご契約へ移行された場合、初年度IP使用料に全額充当します。", 11.5, RED, True)],
    ("※2026年度内の貴社主催・貴社店頭でのイベント実施は、既存ライセンス条件の調整が必要なため個別協議となります。", 10.5, SUB, False),
], gap=4)

# =============================================================================
# 2. 推奨ロードマップ
# =============================================================================
s = add_slide()
header(s, "推奨ロードマップ:2026年度内に手応えを作り、2027年4月に全国へ",
       "パターン1→2→3を段階的に重ね、リスクなく検証を積み上げて年間契約へ接続します")
steps = [
    ("2026年8月", "条件協議・テスト計画の合意", PALE),
    ("2026年9月〜", "パターン1:デザイン開発・監修・仕組み検証を開始(無償)", PALE),
    ("2026年11月〜2027年1月", "パターン2:協賛型イベントで集客・送客データを実地検証", PALE),
    ("2027年2月〜3月", "検証結果の評価 → 年間契約の締結", PALE_RED),
    ("2027年4月〜", "年間契約開始。期首3か月はパターン3(パイロット)として重点運用 → 7月から全国拡大", None),
]
y = Inches(1.7)
for i, (period, desc, fill) in enumerate(steps):
    is_last = i == len(steps) - 1
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72))
    bx.adjustments[0] = 0.12
    bx.shadow.inherit = False
    bx.fill.solid()
    bx.fill.fore_color.rgb = NAVY if is_last else fill
    bx.line.fill.background()
    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = period + "  "
    set_font(r1, 13, WHITE if is_last else RED, True)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, 13, WHITE if is_last else INK, is_last)
    y += Inches(0.86)
    if not is_last:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y - Inches(0.16), Inches(0.22), Inches(0.14))
        ar.shadow.inherit = False
        ar.fill.solid(); ar.fill.fore_color.rgb = RED
        ar.line.fill.background()
textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.6), [
    [("ポイント:", 12.5, RED, True),
     ("2026年度は「作る・確かめる」に集中し、実働は2027年4月から空白なくスタート。パターン2の協賛実績が年間契約の社内稟議の裏付けになります。", 12.5, INK, False)],
])

# =============================================================================
# 3. パイロットの成功基準と年間契約への接続
# =============================================================================
s = add_slide()
header(s, "パイロットの成功基準と年間契約への接続(ご提案)",
       "パイロットを「やって終わり」にしないため、成功基準と移行条件を事前にご合意いただく設計です")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["成功基準(貴社と共同設定)", "測定方法", "目安"],
          [
              ["店頭集客(1開催あたり来場)", "会場カウント+キットQR読取数", "現行IP施策の同等水準(比80%以上)を基準に共同設定"],
              ["撮影会参加組数", "参加受付数(1回30分×5回/日)", "同上"],
              ["ブース接触→相談への転換", "QRアンケート+店頭相談記録", "パイロットで実測しベンチマーク化"],
              ["現場運用品質", "店舗・代理店ヒアリング、事故・クレーム件数", "重大事故・クレームゼロ"],
          ],
          col_widths=[1.5, 1.4, 1.7], row_h=0.55, font_size=11.5)
b = rect(s, Inches(0.7), Inches(4.45), Inches(11.9), Inches(1.85), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(4.62), Inches(11.3), Inches(1.6), [
    ("年間契約への接続条件(基本合意書で事前にご合意)", 13.5, RED, True),
    ("① 成功基準を達成した場合、年間契約へ移行することをパイロット開始前に基本合意(LOI)", 12.5, INK, False),
    ("② 短期IP使用料300万円は、年間契約の初年度IP使用料に全額充当(=パイロット費用は実質ゼロに)", 12.5, INK, False),
    ("③ パイロット期間中のご締結で、60周年記念素材・限定ビジュアルを初年度特典としてご提供", 12.5, INK, False),
    ("④ 計測データは月次レポートで貴社ご報告用フォーマットに整えてご提出(社内ご説明にそのままご利用いただけます)", 12.5, INK, False),
], gap=5)

# =============================================================================
# 4. 年間展開へのランプアップ計画
# =============================================================================
s = add_slide()
header(s, "年間展開へのランプアップ計画",
       "パイロットと並行して供給体制を整備し、判定後は待ち時間ゼロで全国拡大へ移行します")
ramp = [
    ("2027年4月〜6月", "パイロット:首都圏5〜10店舗・10〜20開催。月次レビューで基準進捗を確認", PALE),
    ("2027年7月", "成功基準の判定 → 年間契約へ移行(基本合意済みのためスムーズに発効)", PALE_RED),
    ("2027年7月〜9月", "夏休み商戦:約100開催規模へ拡大(関東+主要都市)", PALE),
    ("2027年10月〜12月", "年末商戦:全国展開へ拡大(年率300〜600開催ペース)", PALE),
    ("2028年1月〜3月", "春商戦:フル稼働。FY28のフル年間展開+複数年契約のご協議", None),
]
y = Inches(1.7)
for i, (period, desc, fill) in enumerate(ramp):
    is_last = i == len(ramp) - 1
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72))
    bx.adjustments[0] = 0.12
    bx.shadow.inherit = False
    bx.fill.solid()
    bx.fill.fore_color.rgb = NAVY if is_last else fill
    bx.line.fill.background()
    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = period + "  "
    set_font(r1, 13, WHITE if is_last else RED, True)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, 13, WHITE if is_last else INK, is_last)
    y += Inches(0.86)
    if not is_last:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y - Inches(0.16), Inches(0.22), Inches(0.14))
        ar.shadow.inherit = False
        ar.fill.solid(); ar.fill.fore_color.rgb = RED
        ar.line.fill.background()
textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.6), [
    [("ポイント:", 12.5, RED, True),
     ("スーツ増産・指定アクターの研修はパイロット期間中に並行整備するため、判定後すぐに拡大できます。パイロットの計測データが、そのまま拡大時のKPI・出店計画の根拠になります。", 12.5, INK, False)],
])

# =============================================================================
# 8. ご確認ポイントへのご回答
# =============================================================================
s = add_slide()
header(s, "ご確認ポイントへのご回答", "貴社資料の確認事項3点への回答です")
add_table(s, Inches(0.7), Inches(1.65), Inches(11.9),
          ["No.", "ご確認内容", "ご回答"],
          [
              ["1", "FY26中のテスト実施可否(単体契約 or 短期契約)",
               "短期・単体契約での対応は可能です。ただし当社の既存ライセンス条件との調整が必要なため、FY26中の実施は形態・会場・時期を個別協議とさせてください。本格開始は2027年4月(FY27期首)を推奨します"],
              ["2", "ノベルティ作成時の費用(ロイヤリティ)発生有無",
               "発生します。作成費用(税別)の20%をロイヤリティとしてお支払いいただきます(作成数の上限なし)"],
              ["3", "イベント実施会場に関する制限事項有無",
               "①外部から見えない更衣室(控室)のご用意 ②他キャラクターとの同時共存は不可 ③非公式商品の店頭ツール・装飾利用は不可。詳細は運用ガイドラインをご提示します"],
          ],
          col_widths=[0.35, 1.6, 2.9], row_h=0.85, font_size=11.5)
textbox(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.6), [
    ("※着ぐるみの稼働ルール(1回30分×1日最大5回、専属スタッフ2名体制)は、ポケモン様の現行運用と同一の水準で設計しています。", 11, SUB, False),
])

# =============================================================================
# 9. 次アクション
# =============================================================================
s = add_slide()
header(s, "本ご提案のまとめと次アクション", "段階導入で意思決定リスクを抑えながら、60周年イヤーのウルトラマンを貴社の店頭資産にしていただくご提案です")
b1 = rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(2.3), fill=PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b1.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.0), [
    ("本ご提案のポイント", 14, RED, True),
    ("・年間条件は3階建て(IP使用料2,000万円+着ぐるみ20万円/開催+ロイヤリティ20%)。固定は①のみの従量型", 13, INK, False),
    ("・テスト(無償)→協賛検証→期首パイロット→年間契約の段階導入。各段階の投資は次に充当され無駄になりません", 13, INK, False),
    ("・成功基準を共同設定し、達成時の年間移行を事前合意(LOI)。パイロットのデータが貴社の社内ご説明資料になります", 13, INK, False),
], gap=7)
add_table(s, Inches(0.7), Inches(4.3), Inches(11.9),
          ["時期", "アクション"],
          [
              ["2026年8月", "条件協議・テスト計画の基本合意(LOI)"],
              ["2026年9月〜2027年3月", "STEP1〜2:先行開発・協賛型イベント検証+供給体制整備"],
              [("2027年4月〜", True, RED), ("STEP3〜4:期首パイロット → 7月 年間契約発効 → 全国拡大", True, RED)],
          ],
          col_widths=[1, 2.8], row_h=0.5, font_size=12.5)

out = "ウルトラマンIP導入のご提案_統合版.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
