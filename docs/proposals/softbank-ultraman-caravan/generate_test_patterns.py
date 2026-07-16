# -*- coding: utf-8 -*-
"""テスト実施のご提案パターン(2枚・円谷フォーマット)を生成する。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
BAND_IMG = os.path.join(ASSETS, "band_header.png")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
page_no = [0]


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def add_slide():
    page_no[0] += 1
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def footer(slide):
    textbox(slide, Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
            [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
            [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
            [(str(page_no[0]), 10, SILVER, False)], align=PP_ALIGN.RIGHT)


def header(slide, title, sub=None):
    slide.shapes.add_picture(BAND_IMG, 0, 0, width=SW, height=Inches(0.86))
    textbox(slide, Inches(0.5), Inches(0.08), Inches(11.2), Inches(0.7),
            [(title, 21, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        textbox(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.4),
                [(sub, 12.5, SUB, False)])
    footer(slide)


def add_table(slide, x, y, w, headers, rows, col_widths=None, font_size=11.5, row_h=0.5):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(headers), x, y, w, Inches(row_h * n_rows)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.08)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        set_font(r, font_size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            text, bold, color = val, False, INK
            if isinstance(val, tuple):
                text, bold = val[0], val[1]
                color = val[2] if len(val) > 2 else INK
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = text
            set_font(r, font_size, color, bold)
    return tbl


# =============================================================================
# 1. テスト実施のご提案パターン
# =============================================================================
s = add_slide()
header(s, "テスト実施のご提案パターン(確認ポイント1へのご回答)",
       "既存ライセンス条件を踏まえ、2026年度内から段階的に検証を始められる3つのパターンをご用意しました")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["", "パターン1:先行開発テスト", "パターン2:イベント協賛型テスト", "パターン3:期首パイロット"],
          [
              ["時期", "2026年9月〜(即時開始可)", "2026年秋〜冬", "2027年4月〜6月"],
              ["内容", "キット・訴求ツールのデザイン開発と監修、お客様調査、デジタル施策(QR等)の仕組み検証 ※IP実働なし",
               "当社主催(または商業施設主催)の60周年イベントに貴社がご協賛・ブース出展。着ぐるみは当社イベントとして稼働",
               "首都圏5〜10店舗・着ぐるみ10〜20開催の単体テスト(全国展開と同一の運用形式)"],
              ["ご契約形態", "秘密保持契約+監修合意", "イベント協賛契約", "短期単体契約(3か月)"],
              ["費用目安", "原則無償(制作実費のみ)", "ご協賛金 300〜500万円", "短期IP使用料300万円+着ぐるみ20万円/開催+作成物ロイヤリティ"],
              ["検証できること", "クリエイティブ品質・監修フロー・運用設計", "集客力・ブース送客・相談転換の実地データ", "全国展開前の店頭オペレーション・KPI実測"],
          ],
          col_widths=[0.55, 1.3, 1.45, 1.3], row_h=0.72, font_size=10.5)
textbox(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.7), [
    [("※パターン3の短期IP使用料は、年間ご契約へ移行された場合、初年度IP使用料に全額充当します。", 11.5, RED, True)],
    ("※2026年度内の貴社主催・貴社店頭でのイベント実施は、既存ライセンス条件の調整が必要なため個別協議となります。", 10.5, SUB, False),
], gap=4)

# =============================================================================
# 2. 推奨ロードマップ
# =============================================================================
s = add_slide()
header(s, "推奨ロードマップ:2026年度内に手応えを作り、2027年4月に全国へ",
       "パターン1→2→3を段階的に重ね、リスクなく検証を積み上げて年間契約へ接続します")
steps = [
    ("2026年8月", "条件協議・テスト計画の合意", PALE),
    ("2026年9月〜", "パターン1:デザイン開発・監修・仕組み検証を開始(無償)", PALE),
    ("2026年11月〜2027年1月", "パターン2:協賛型イベントで集客・送客データを実地検証", PALE),
    ("2027年2月〜3月", "検証結果の評価 → 年間契約の締結", PALE_RED),
    ("2027年4月〜", "年間契約開始。期首3か月はパターン3(パイロット)として重点運用 → 7月から全国拡大", None),
]
y = Inches(1.7)
for i, (period, desc, fill) in enumerate(steps):
    is_last = i == len(steps) - 1
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72))
    bx.adjustments[0] = 0.12
    bx.shadow.inherit = False
    bx.fill.solid()
    bx.fill.fore_color.rgb = NAVY if is_last else fill
    bx.line.fill.background()
    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = period + "  "
    set_font(r1, 13, WHITE if is_last else RED, True)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, 13, WHITE if is_last else INK, is_last)
    y += Inches(0.86)
    if not is_last:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y - Inches(0.16), Inches(0.22), Inches(0.14))
        ar.shadow.inherit = False
        ar.fill.solid(); ar.fill.fore_color.rgb = RED
        ar.line.fill.background()
textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.6), [
    [("ポイント:", 12.5, RED, True),
     ("2026年度は「作る・確かめる」に集中し、実働は2027年4月から空白なくスタート。パターン2の協賛実績が年間契約の社内稟議の裏付けになります。", 12.5, INK, False)],
])

# =============================================================================
# 3. パイロットの成功基準と年間契約への接続
# =============================================================================
s = add_slide()
header(s, "パイロットの成功基準と年間契約への接続(ご提案)",
       "パイロットを「やって終わり」にしないため、成功基準と移行条件を事前にご合意いただく設計です")
add_table(s, Inches(0.7), Inches(1.55), Inches(11.9),
          ["成功基準(貴社と共同設定)", "測定方法", "目安"],
          [
              ["店頭集客(1開催あたり来場)", "会場カウント+キットQR読取数", "現行IP施策の同等水準(比80%以上)を基準に共同設定"],
              ["撮影会参加組数", "参加受付数(1回30分×5回/日)", "同上"],
              ["ブース接触→相談への転換", "QRアンケート+店頭相談記録", "パイロットで実測しベンチマーク化"],
              ["現場運用品質", "店舗・代理店ヒアリング、事故・クレーム件数", "重大事故・クレームゼロ"],
          ],
          col_widths=[1.5, 1.4, 1.7], row_h=0.55, font_size=11.5)
b = rect(s, Inches(0.7), Inches(4.45), Inches(11.9), Inches(1.85), fill=PALE_RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b.adjustments[0] = 0.05
textbox(s, Inches(1.0), Inches(4.62), Inches(11.3), Inches(1.6), [
    ("年間契約への接続条件(基本合意書で事前にご合意)", 13.5, RED, True),
    ("① 成功基準を達成した場合、年間契約へ移行することをパイロット開始前に基本合意(LOI)", 12.5, INK, False),
    ("② 短期IP使用料300万円は、年間契約の初年度IP使用料に全額充当(=パイロット費用は実質ゼロに)", 12.5, INK, False),
    ("③ パイロット期間中のご締結で、60周年記念素材・限定ビジュアルを初年度特典としてご提供", 12.5, INK, False),
    ("④ 計測データは月次レポートで貴社ご報告用フォーマットに整えてご提出(社内ご説明にそのままご利用いただけます)", 12.5, INK, False),
], gap=5)

# =============================================================================
# 4. 年間展開へのランプアップ計画
# =============================================================================
s = add_slide()
header(s, "年間展開へのランプアップ計画",
       "パイロットと並行して供給体制を整備し、判定後は待ち時間ゼロで全国拡大へ移行します")
ramp = [
    ("2027年4月〜6月", "パイロット:首都圏5〜10店舗・10〜20開催。月次レビューで基準進捗を確認", PALE),
    ("2027年7月", "成功基準の判定 → 年間契約へ移行(基本合意済みのためスムーズに発効)", PALE_RED),
    ("2027年7月〜9月", "夏休み商戦:約100開催規模へ拡大(関東+主要都市)", PALE),
    ("2027年10月〜12月", "年末商戦:全国展開へ拡大(年率300〜600開催ペース)", PALE),
    ("2028年1月〜3月", "春商戦:フル稼働。FY28のフル年間展開+複数年契約のご協議", None),
]
y = Inches(1.7)
for i, (period, desc, fill) in enumerate(ramp):
    is_last = i == len(ramp) - 1
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72))
    bx.adjustments[0] = 0.12
    bx.shadow.inherit = False
    bx.fill.solid()
    bx.fill.fore_color.rgb = NAVY if is_last else fill
    bx.line.fill.background()
    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = period + "  "
    set_font(r1, 13, WHITE if is_last else RED, True)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, 13, WHITE if is_last else INK, is_last)
    y += Inches(0.86)
    if not is_last:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y - Inches(0.16), Inches(0.22), Inches(0.14))
        ar.shadow.inherit = False
        ar.fill.solid(); ar.fill.fore_color.rgb = RED
        ar.line.fill.background()
textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.6), [
    [("ポイント:", 12.5, RED, True),
     ("スーツ増産・指定アクターの研修はパイロット期間中に並行整備するため、判定後すぐに拡大できます。パイロットの計測データが、そのまま拡大時のKPI・出店計画の根拠になります。", 12.5, INK, False)],
])

out = "テスト実施のご提案パターン.pptx"
prs.save(out)
print(f"saved: {out} / slides: {page_no[0]}")
