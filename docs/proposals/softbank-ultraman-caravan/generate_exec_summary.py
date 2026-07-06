# -*- coding: utf-8 -*-
"""役員回覧用エグゼクティブ1枚サマリー(円谷フォーマット)を生成する。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x5A, 0x62, 0x6B)
RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x23, 0x22, 0x78)
SILVER = RGBColor(0x8E, 0x94, 0x9B)
PALE = RGBColor(0xF2, 0xF4, 0xF6)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
s = prs.slides.add_slide(prs.slide_layouts[6])


def set_font(run, size, color=INK, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', FONT)


def tb(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, gap=3):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        for (t, sz, c, b) in (item if isinstance(item, list) else [item]):
            r = p.add_run()
            r.text = t
            set_font(r, sz, c, b)
    return box


def panel(x, y, w, h, title, lines, fill=PALE):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.adjustments[0] = 0.045
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    tb(x + Inches(0.18), y + Inches(0.1), w - Inches(0.36), Inches(0.3),
       [(title, 12.5, RED, True)])
    tb(x + Inches(0.18), y + Inches(0.42), w - Inches(0.36), h - Inches(0.5), lines, gap=4)


# ヘッダー帯
s.shapes.add_picture(os.path.join(ASSETS, "band_header.png"), 0, 0, width=SW, height=Inches(0.8))
tb(Inches(0.45), Inches(0.06), Inches(11.3), Inches(0.68),
   [("SoftBank × ウルトラマン 未来防衛隊キャラバン|エグゼクティブサマリー", 18, WHITE, True)],
   anchor=MSO_ANCHOR.MIDDLE)

# 結論帯
sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.95), Inches(12.43), Inches(0.62))
sp.shadow.inherit = False
sp.fill.solid(); sp.fill.fore_color.rgb = PALE_RED; sp.line.fill.background()
tb(Inches(0.65), Inches(1.0), Inches(12.0), Inches(0.5), [
    [("結論:", 13, RED, True),
     ("ウルトラマンを商業施設集客の核とし、親子・ファミリーをSoftBankショップ・BB・でんき相談へ送客。推奨はFY27年間1.5億円(B案)。", 13, INK, True)],
], anchor=MSO_ANCHOR.MIDDLE)

# 4パネル(上段)
panel(Inches(0.45), Inches(1.75), Inches(6.05), Inches(1.78), "課題(FY27以降)", [
    ("・現行キャラクターIPは、FY27以降イベントでの着ぐるみ活用が終了予定", 11, INK, False),
    ("・商業施設イベントの「会いに行く理由」=リアル集客機能に空白", 11, INK, False),
    ("・週末ファミリー集客の目的来館装置が必要", 11, INK, False),
])
panel(Inches(6.83), Inches(1.75), Inches(6.05), Inches(1.78), "解決(ウルトラマン60周年)", [
    ("・ヒーロー撮影会・ステージ・館内ミッションラリー・デジタル隊員証を1パッケージ化", 11, INK, False),
    ("・現行IPの置き換えではなく「リアルイベント集客機能」を補完(他IPと併用可)", 11, INK, False),
    ("・自社アクター体制で年160回規模の稼働を安定供給", 11, INK, False),
])
# 4パネル(下段)
panel(Inches(0.45), Inches(3.65), Inches(6.05), Inches(1.9), "SoftBank様への効果", [
    ("・最終KGI:新規契約(モバイル・BB・でんき)の獲得", 11, INK, True),
    ("・店舗送客3,000〜8,000人/相談1,000〜3,000件(B案・年間)", 11, INK, False),
    ("・デジタル隊員証で送客・属性を施設別に実測、月次レポート", 11, INK, False),
    ("・ファミリー集客企画は商業施設の場所代交渉の材料に(全国で年間数億円規模のコスト構造に波及)", 11, INK, False),
])
panel(Inches(6.83), Inches(3.65), Inches(6.05), Inches(1.9), "プランと金額(税別・概算)", [
    ("・A案 3,500万円:FY27期首(2027年4月〜)3〜4か月の小規模検証", 11, INK, False),
    [("・B案 1億5,000万円:FY27年間・首都圏10〜20施設 ", 11, INK, True), ("【推奨】", 11, RED, True)],
    ("・C案 2億5,000万円:全国大型展開(B案実施後の拡張)", 11, INK, False),
    ("・年間IP使用料+ヒーロー稼働+イベント運営+制作+計測まで含む実行費一式の総額", 10.5, SUB, False),
])

# スケジュール+次アクション
tb(Inches(0.45), Inches(5.72), Inches(12.4), Inches(0.9), [
    [("スケジュール:", 11.5, RED, True),
     ("FY26下期=施設選定・制作準備(実働なし)→ 2027年4月 FY27契約開始・実働(期首に検証→7月以降拡大)→ 2028年4月〜 複数年・全国展開を協議", 11.5, INK, False)],
    [("次アクション:", 11.5, RED, True),
     ("方向性ご合意のうえ、条件打診・概算精査(2026年9〜10月)、施設選定・設計(11月〜)に着手。60周年イヤーの稼働枠確保のため、8月末までの内示を推奨。", 11.5, INK, False)],
], gap=6)

# フッター
tb(Inches(5.17), Inches(7.16), Inches(3.0), Inches(0.28),
   [("© TSUBURAYA PRODUCTIONS", 8, SILVER, False)], align=PP_ALIGN.CENTER)
tb(Inches(10.25), Inches(7.1), Inches(2.2), Inches(0.3),
   [("Strictly Confidential", 10, RED, True)], align=PP_ALIGN.RIGHT)
tb(Inches(0.45), Inches(7.16), Inches(4.0), Inches(0.28),
   [("2026年7月|(株)円谷プロダクション", 8, SILVER, False)])

prs.save("エグゼクティブサマリー_1枚.pptx")
print("saved")
