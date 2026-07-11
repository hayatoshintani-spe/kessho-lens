#!/usr/bin/env python3
"""S6とS7〜S13の改修(7/11指示)。

- S6: 円谷中心の放射図 → 円谷×パートナーの対等な共創図へ
- S7〜S12: 各企画に「例|」の具体イメージ行を追加(ライセンシーが自社事業で想像できる粒度に)
- S13: コラボ4文脈の商品化機会に具体例を追記

    python3 rework_s6_categories.py  # v16を読み、v17を出力
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

BG    = RGBColor(0x00, 0x07, 0x14)
WHITE = RGBColor(0xF4, 0xF7, 0xFD)
SUB   = RGBColor(0xA9, 0xB6, 0xD2)
DIM   = RGBColor(0x6A, 0x77, 0x94)
RED   = RGBColor(0xE8, 0x38, 0x2F)
REDDK = RGBColor(0xA8, 0x1C, 0x18)
PANEL = RGBColor(0x0C, 0x17, 0x30)
PANEL2= RGBColor(0x11, 0x1F, 0x3E)
LINE  = RGBColor(0x25, 0x36, 0x5E)
BLUE  = RGBColor(0x7F, 0xA8, 0xE0)
JP_FONT = "Yu Gothic"
ML = Inches(0.55)
CW = Inches(12.23)


def style_run(run, size, color=SUB, bold=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = JP_FONT
    f.language_id = MSO_LANGUAGE_ID.JAPANESE
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP_FONT)


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for spec in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = spec if isinstance(spec, list) else [spec]
        opts = {}
        for rs in runs:
            text, size, color, bold = rs[0], rs[1], rs[2], rs[3]
            opts = rs[4] if len(rs) > 4 else opts
            run = p.add_run(); run.text = text
            style_run(run, size, color, bold)
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(opts.get("before", 0))
        p.space_after = Pt(opts.get("after", 3))
        p.line_spacing = opts.get("line", 1.25)
    return box


def rect(slide, x, y, w, h, fill, line_color=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


# ============================================================ S6:共創図
def rebuild_s6(slide):
    # 既存のカテゴリチップ・中央箱・下部コピーを撤去(ヘッダーは残す)
    for sh in list(slide.shapes):
        keep = False
        if sh.has_text_frame:
            t = sh.text_frame.text
            if t.startswith(('06|', '本日ご参加', '7つのカテゴリー')):
                keep = True
        if Emu(sh.top).inches < 1.6 and Emu(sh.height).inches < 0.1:  # 赤ルール
            keep = True
        if not keep:
            sh._element.getparent().remove(sh._element)

    # 左:円谷が持ち寄るもの
    LX, PW, PY, PH = ML, Inches(5.40), Inches(2.20), Inches(2.10)
    rect(slide, LX, PY, PW, PH, PANEL, LINE, Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(slide, LX, PY, PW, Pt(3), RED)
    add_text(slide, LX + Inches(0.28), PY + Inches(0.18), PW - Inches(0.56), PH - Inches(0.36),
             [[("円谷プロダクション", 13, WHITE, True), ("  が持ち寄るもの", 10.5, BLUE, True)],
              ("・60年のIPと歴代ヒーロー・怪獣", 10.5, SUB, False, {"before": 6, "line": 1.45}),
              ("・スタイルガイドと監修の伴走", 10.5, SUB, False, {"line": 1.45}),
              ("・公式SNS・ULTRA MART・イベントの発信力", 10.5, SUB, False, {"line": 1.45})])

    # 中央:×
    cx = Inches(6.28)
    rect(slide, cx, Inches(2.95), Inches(0.78), Inches(0.78), None, RED, Pt(1.6), shape=MSO_SHAPE.OVAL)
    add_text(slide, cx, Inches(2.95), Inches(0.78), Inches(0.78),
             [("×", 20, RED, True, {"align": PP_ALIGN.CENTER, "after": 0, "line": 1.0})],
             anchor=MSO_ANCHOR.MIDDLE)

    # 右:パートナーが持ち寄るもの
    RX = Inches(7.38)
    rect(slide, RX, PY, PW, PH, PANEL, LINE, Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(slide, RX, PY, PW, Pt(3), BLUE)
    add_text(slide, RX + Inches(0.28), PY + Inches(0.18), PW - Inches(0.56), PH - Inches(0.36),
             [[("パートナー企業の皆さま", 13, WHITE, True), ("  が持ち寄るもの", 10.5, BLUE, True)],
              ("・商品開発力・製造・品質", 10.5, SUB, False, {"before": 6, "line": 1.45}),
              ("・売場・販路・EC・顧客接点", 10.5, SUB, False, {"line": 1.45}),
              ("・顧客理解とプロモーションの知見", 10.5, SUB, False, {"line": 1.45})])

    # ▼ → 共創の帯
    add_text(slide, ML, Inches(4.40), CW, Inches(0.22),
             [("▼", 11, RED, True, {"align": PP_ALIGN.CENTER, "after": 0, "line": 1.0})])
    rect(slide, ML, Inches(4.68), CW, Inches(0.60), REDDK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    add_text(slide, ML, Inches(4.68), CW, Inches(0.60),
             [[("2027年の商機を、一緒につくる", 15, WHITE, True),
               ("  — 主導するのは、皆さまの事業です", 11, WHITE, False)]],
             anchor=MSO_ANCHOR.MIDDLE)

    # カテゴリチップ列
    add_text(slide, ML, Inches(5.44), CW, Inches(0.26),
             [("募集カテゴリー(この後、1枚ずつご提案します)", 10, BLUE, True)])
    cats = ["MD・グッズ", "食品・外食", "流通・EC", "広告・SP", "デジタル", "イベント・観光", "IPコラボ"]
    gap = Inches(0.14)
    cw_ = Emu(int((CW - gap * 6) / 7))
    for i, c in enumerate(cats):
        x = ML + Emu(int((cw_ + gap) * i))
        rect(slide, x, Inches(5.76), cw_, Inches(0.50), PANEL2, LINE, Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.20)
        add_text(slide, x, Inches(5.76), cw_, Inches(0.50),
                 [(c, 9.5, WHITE, True, {"align": PP_ALIGN.CENTER, "after": 0})],
                 anchor=MSO_ANCHOR.MIDDLE)

    # 下部コピー
    rect(slide, ML, Inches(6.42), CW, Inches(0.56), PANEL)
    rect(slide, ML, Inches(6.42), Pt(4), Inches(0.56), RED)
    add_text(slide, ML + Inches(0.24), Inches(6.42), CW - Inches(0.48), Inches(0.56),
             [("どのカテゴリーでも、主役は皆さまの事業。円谷はIPと発信力で伴走します。", 14, WHITE, True)],
             anchor=MSO_ANCHOR.MIDDLE)


# ============================================================ S7〜S12:具体例つきリスト
PROPOSALS = {
    7: (2.32, 0.78, 11.5, 9.5, [
        ("新作連動キッズMD", "2027年新シリーズの変身なりきり玩具を、放送開始週に量販店の店頭へ"),
        ("映画・ゼロ関連商品", "劇場公開に合わせたフィギュア・前売券セット・入場者特典の連動設計"),
        ("歴代ヒーロー/怪獣コレクション", "60年の歴代からブラインドトイ・ガチャ・カードを通年シリーズ化"),
        ("大人ファン向けプレミアム商品", "受注生産の高価格帯フィギュア・アパレルを限定箱+ECで展開"),
        ("女性ライト層向け雑貨", "モフサンド文脈のポーチ・文具をバラエティショップの平台へ"),
    ]),
    8: (2.26, 0.655, 11, 9, [
        ("限定パッケージ", "全6種のヒーローパッケージで「選ぶ楽しさ」を棚に作る(菓子・飲料)"),
        ("購入特典カード", "1個につきカード1枚。集めたくなる仕掛けでリピート購買を作る"),
        ("映画半券キャンペーン", "2027夏の映画半券×対象商品レシートで限定グッズ応募"),
        ("親子ペア応募", "バーコード応募でウルサマ親子ペア招待・限定ショー観覧"),
        ("外食コラボメニュー", "「ウルトラ創業祭」型のコラボメニュー+店内装飾+ノベルティ"),
        ("ウルサマ来場者特典", "会場配布のクーポンで店舗へ送客。イベント→店の導線を作る"),
    ]),
    9: (2.32, 0.78, 11.5, 9.5, [
        ("映画公開記念棚/ウルサマ連動棚", "2027夏、映画×イベントの山に合わせて全国エンド棚を確保"),
        ("ゼロ特集棚/怪獣特集棚", "キャラクター軸の編集棚で、既存商品も一緒に動かす"),
        ("IPコラボ棚", "モフサンド棚のように、女性客が立ち止まる棚を日用品売場に"),
        ("POPUP/EC限定販売", "有楽町マルイ型の女性向けPOPUP。EC先行→店舗拡大の二段展開"),
        ("購入特典", "チェーン限定の特典カード・ノベルティで「この店で買う理由」を作る"),
    ]),
    10: (2.32, 0.78, 11.5, 9.5, [
        ("親子集客/夏休み販促", "SC・量販の夏休み集客にヒーローショー+限定ノベルティを一括提案"),
        ("防災/交通安全", "自治体・インフラ企業の啓発企画に、世代を超えた信頼感を乗せる"),
        ("地域創生", "ご当地怪獣×観光キャンペーン。地域イベント登場+限定グッズ"),
        ("未来・テクノロジー/CSR", "企業ブランディング広告にヒーロー文脈を活用(掲出実績あり)"),
        ("商業施設回遊", "スタンプラリー+デジタルスタンプで館内回遊・滞在時間を伸ばす"),
    ]),
    11: (2.32, 0.78, 11.5, 9.5, [
        ("ゲーム内コラボ/限定アバター", "人気タイトルに期間限定のウルトラマンスキン・コラボイベント"),
        ("デジタルスタンプ", "位置情報スタンプで店舗・施設を回遊させ、コンプリート特典へ"),
        ("SNSキャンペーン", "変身ポーズ投稿・フィルター企画でUGCを量産する参加型設計"),
        ("配信連動/EC限定デジタル特典", "配信視聴×EC購入でデジタル壁紙・限定コンテンツを付与"),
        ("ファンコミュニティ施策", "会員限定の先行情報・投票企画でコアファンの熱量を可視化"),
    ]),
    12: (2.32, 0.78, 11.5, 9.5, [
        ("ヒーローショー/グリーティング", "ウルサマ型のショー+握手会を、施設の集客の山に合わせて開催"),
        ("撮影会", "人気ヒーローとの撮影会を有料企画化。SNS拡散も同時に生む"),
        ("館内スタンプラリー/地域回遊", "台紙+スタンプ台で親子の全館回遊・滞在時間を延ばす"),
        ("商業施設イベント/観光地コラボ", "吹き抜けイベント+期間限定売場のセットで館全体を活性化"),
        ("物販連動", "会場限定商品と物販列で「来た人が必ず買う」導線を設計"),
    ]),
}


def rework_list(slide, y0, step, tsize, esize, items):
    # 既存の番号チップ(0.30幅)とタイトル行(x=1.00, w=5.75)を撤去
    for sh in list(slide.shapes):
        xi, wi = Emu(sh.left).inches, Emu(sh.width).inches
        if abs(xi - 0.55) < 0.03 and abs(wi - 0.30) < 0.03:
            sh._element.getparent().remove(sh._element)
        elif abs(xi - 1.00) < 0.03 and abs(wi - 5.75) < 0.03:
            sh._element.getparent().remove(sh._element)
    for i, (title, ex) in enumerate(items):
        y = y0 + step * i
        rect(slide, ML, Inches(y + 0.01), Inches(0.28), Inches(0.28), RED, shape=MSO_SHAPE.OVAL)
        add_text(slide, ML, Inches(y + 0.01), Inches(0.28), Inches(0.28),
                 [(str(i + 1), 10.5, WHITE, True, {"align": PP_ALIGN.CENTER, "after": 0, "line": 1.0})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(0.98), Inches(y - 0.02), Inches(5.95), Inches(0.30),
                 [(title, tsize, WHITE, True, {"after": 0})])
        add_text(slide, Inches(0.98), Inches(y + 0.275), Inches(5.95), Inches(0.42),
                 [[("例|", esize - 0.5, BLUE, True), (ex, esize, SUB, False, {"line": 1.15, "after": 0})]])


# ============================================================ S13:テーブルに具体例
S13_CELLS = {
    "雑貨・文具": "雑貨・文具(ぬいポーチ・マステ・カフェコラボ)",
    "玩具・カード": "玩具・カード(ブラインドトイ・ゲーム内コラボ)",
    "流通棚・雑貨": "流通棚・雑貨(ファミリーギフト・季節催事)",
    "プレミアム・アパレル": "プレミアム・アパレル(ハイエンド・限定箱)",
}


def rework_s13(slide):
    for sh in slide.shapes:
        if not getattr(sh, 'has_table', False) or not sh.has_table:
            continue
        for row in sh.table.rows:
            for cell in row.cells:
                t = cell.text_frame.text.strip()
                if t in S13_CELLS:
                    p = cell.text_frame.paragraphs[0]
                    for r in list(p.runs):
                        r._r.getparent().remove(r._r)
                    run = p.add_run(); run.text = S13_CELLS[t]
                    style_run(run, 9.5, SUB, False)


if __name__ == "__main__":
    prs = Presentation("ULTRAMAN_BUSINESS_PARTNERS_MEETING_2026_v16.pptx")
    rebuild_s6(prs.slides[5])
    for sidx, (y0, step, ts, es, items) in PROPOSALS.items():
        rework_list(prs.slides[sidx - 1], y0, step, ts, es, items)
    rework_s13(prs.slides[12])
    prs.save("ULTRAMAN_BUSINESS_PARTNERS_MEETING_2026_v17.pptx")
    print("saved v17 (S6 + S7-13 reworked)")
