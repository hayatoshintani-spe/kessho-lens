#!/usr/bin/env python3
"""ダミー画像(カンプ)を「テキスト+ボックス」のネイティブ図形に置換する後処理.

Googleスライド変換後も編集できるよう、画像プレースホルダーをすべて
PowerPoint図形(角丸四角+ラベル)に差し替える。ユーザーのテキスト修正は保持。

    python3 replace_comps_with_boxes.py <入力.pptx> <出力.pptx>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

WHITE = RGBColor(0xF4, 0xF7, 0xFD)
SUB   = RGBColor(0xA9, 0xB6, 0xD2)
DIM   = RGBColor(0x8A, 0x96, 0xB2)
RED   = RGBColor(0xE8, 0x38, 0x2F)
REDDK = RGBColor(0xA8, 0x1C, 0x18)
PANEL = RGBColor(0x0C, 0x17, 0x30)
PANEL2= RGBColor(0x11, 0x1F, 0x3E)
BORDER= RGBColor(0x4A, 0x5F, 0x8F)
JP_FONT = "Yu Gothic"


def style_run(run, size, color, bold=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = JP_FONT
    f.language_id = MSO_LANGUAGE_ID.JAPANESE
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP_FONT)


def ph_box(slide, x, y, w, h, label, sub="写真差し替え", label_size=None,
           fill=PANEL, dash=True, border=BORDER, sub_color=DIM):
    """ラベル付きプレースホルダーボックス."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = min(0.12, 0.08 * Emu(h).inches)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(1.1)
    if dash:
        shp.line.dash_style = 4
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    wi, hi = Emu(w).inches, Emu(h).inches
    ls = label_size or (12 if wi > 2.4 else 10 if wi > 1.5 else 8.5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.15
    r = p.add_run(); r.text = label
    style_run(r, ls, WHITE, bold=True)
    if sub and hi > 0.9:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        style_run(r2, max(ls - 3, 7), sub_color)
    return shp


def remove_pic(slide, pic):
    """画像シェイプと参照リレーションを削除."""
    rIds = set()
    for blip in pic._element.iter(qn("a:blip")):
        rId = blip.get(qn("r:embed"))
        if rId:
            rIds.add(rId)
    pic._element.getparent().remove(pic._element)
    for rId in rIds:
        try:
            slide.part.drop_rel(rId)
        except KeyError:
            pass


# スライド番号(1始まり) → 画像位置順(上→下、左→右)のラベル
LABELS = {
    1: ["__KV__"],                                  # 全面背景(特別処理)
    3: ["会場に来る", "写真を撮る", "商品を手に取る", "誰かに共有する", "グリーティング"],
    4: ["ウルサマ客席", "ヒーローショー", "商品化", "IPコラボ",
        "ウルトラマート", "流通展開", "広告タイアップ", "海外展開"],
    7: ["アパレルライン(モック)", "アクスタ・カード", "プレミアムフィギュア", "売場イメージ"],
    8: ["限定パッケージ+コラボメニュー(モック)", "購入特典カード", "半券キャンペーン", "親子購入シーン"],
    9: ["映画公開記念棚+EC特集(モック)", "POPUP売場", "スタンプラリー台紙"],
    10: ["交通広告 駅貼りB0(モック)", "SNS投稿画面", "グリーティング集客"],
    11: ["ゲーム内コラボ画面(モック)", "ガチャ・アバター", "YouTubeサムネ", "ライブ配信画面"],
    12: ["ヒーローショー+観客", "グリーティング", "スタンプラリー", "商業施設イベント"],
    13: ["かわいい文脈", "ホビー文脈", "キャラクター文脈", "ヒーロー文脈", "コラボ棚"],
    14: ["マート外観", "店内の賑わい", "限定コーナー", "POPUP"],
    15: ["ロゴ使用例+カラーパレット", "キャラ配置例+OK/NG", "商品展開例"],
    16: ["__FLOW__"],                               # フローストリップ(特別処理)
    17: ["__QR__"],
    18: ["__CONTENT__"],                            # コンテンツパノラマ(特別処理)
    19: ["TVシリーズ", "映画・大型映像", "配信・YouTube"],
    21: ["__QR__", "クロージングビジュアル(ヒーロー集合)"],
}

FLOW_STEPS = [
    ("① 個別相談", "QR登録から1週間以内にご連絡"),
    ("② 素材・ガイド共有", "NDA後 3営業日以内【仮】"),
    ("③ 企画提出", "パートナー様側リードタイム"),
    ("④ 監修・修正", "初稿10営業日/修正 各5営業日【仮】"),
    ("⑤ 商品化・販促展開", "量産前見本の確認 5営業日以内【仮】"),
]

CONTENT_ZONES = [
    ("TVシリーズ", "継続接点"), ("映画・大型映像", "大型話題化"),
    ("配信・YouTube", "コアファン接点"), ("イベント", "熱量づくり"),
]


def main(src, dst):
    prs = Presentation(src)
    replaced = 0
    for idx, slide in enumerate(prs.slides, 1):
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]
        if not pics:
            continue
        # スライド1の右上60thロゴ(小さい実ロゴ)は残す
        targets = [p for p in pics
                   if not (idx == 1 and Emu(p.width).inches < 4 and Emu(p.top).inches < 1)]
        targets.sort(key=lambda p: (round(Emu(p.top).inches, 1), Emu(p.left).inches))
        labels = LABELS.get(idx, [])
        for i, pic in enumerate(targets):
            x, y, w, h = pic.left, pic.top, pic.width, pic.height
            label = labels[i] if i < len(labels) else "ビジュアル"
            remove_pic(slide, pic)
            replaced += 1
            if label == "__KV__":
                # 全面背景 → 右側にKVプレースホルダーのみ配置
                ph_box(slide, Inches(9.55), Inches(1.25), Inches(3.2), Inches(4.6),
                       "キービジュアル", "ヒーロー集合/写真差し替え", label_size=13)
            elif label == "__QR__":
                ph_box(slide, x, y, w, h, "QRコード",
                       "本番URL発行後に差し替え(7/18)", label_size=14,
                       border=RED, sub_color=SUB)
            elif label == "__FLOW__":
                n = len(FLOW_STEPS)
                gap = Inches(0.12)
                bw = int((w - gap * (n - 1)) / n)
                for j, (t, term) in enumerate(FLOW_STEPS):
                    shp = ph_box(slide, x + (bw + gap) * j, y, bw, h, t, term,
                                 label_size=10.5,
                                 fill=REDDK if j == n - 1 else PANEL2,
                                 dash=False, border=BORDER, sub_color=SUB)
            elif label == "__CONTENT__":
                n = len(CONTENT_ZONES)
                gap = Inches(0.16)
                bw = int((w - gap * (n - 1)) / n)
                for j, (t, role) in enumerate(CONTENT_ZONES):
                    ph_box(slide, x + (bw + gap) * j, y, bw, h, t,
                           role + "/開示可能ビジュアル差し替え", label_size=14)
            else:
                sub = "モック差し替え" if "モック" in label else "写真差し替え"
                ph_box(slide, x, y, w, h, label.replace("(モック)", ""), sub)
    prs.save(dst)
    print(f"replaced {replaced} pictures -> {dst}")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
