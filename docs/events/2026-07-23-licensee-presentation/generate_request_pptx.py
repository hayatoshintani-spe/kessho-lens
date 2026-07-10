#!/usr/bin/env python3
"""素材依頼資料(ビジュアルショットリスト)PowerPoint生成.

依頼画像素材リスト(全49点+確認事項)を、各素材の「欲しい絵」見本
(制作済みカンプ)のサムネイル付きで1冊のpptxにする。社内回覧・
説明会・各部への依頼にそのまま使える。公式テンプレートベース。

    python3 generate_request_pptx.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

BG    = RGBColor(0x00, 0x07, 0x14)
WHITE = RGBColor(0xF4, 0xF7, 0xFD)
SUB   = RGBColor(0xA9, 0xB6, 0xD2)
DIM   = RGBColor(0x6A, 0x77, 0x94)
RED   = RGBColor(0xE8, 0x38, 0x2F)
REDDK = RGBColor(0xA8, 0x1C, 0x18)
PANEL = RGBColor(0x0C, 0x17, 0x30)
PANEL2= RGBColor(0x11, 0x1F, 0x3E)
LINE  = RGBColor(0x25, 0x36, 0x5E)
BLUE  = RGBColor(0x7F, 0xA8, 0xE0)
GOLD  = RGBColor(0xF4, 0xC0, 0x4E)

JP_FONT = "Yu Gothic"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.55)
CW = Inches(12.23)
TITLE_W = Inches(9.55)
DRIVE_URL = "drive.google.com/drive/folders/1Nk6C5s9mTxFCz0cYNK0Qap2HEdQQLSNW"

prs = Presentation("template_PARTNERS_MEETING2026.pptx")
BLANK = prs.slide_layouts[6]
for sldId in list(prs.slides._sldIdLst):
    prs.part.drop_rel(sldId.rId)
    prs.slides._sldIdLst.remove(sldId)


def style_run(run, size, color=SUB, bold=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = JP_FONT
    f.language_id = MSO_LANGUAGE_ID.JAPANESE
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP_FONT)
    if spacing is not None:
        rPr.set("spc", str(spacing))


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for spec in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = spec if isinstance(spec, list) else [spec]
        opts = {}
        for rs in runs:
            text, size, color, bold = rs[0], rs[1], rs[2], rs[3]
            opts = rs[4] if len(rs) > 4 else opts
            run = p.add_run(); run.text = text
            style_run(run, size, color, bold, spacing=opts.get("spc"))
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opts.get("after", 3))
        p.line_spacing = opts.get("line", 1.25)
    return box


def rect(slide, x, y, w, h, fill, line_color=None, line_w=None, dash=False,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
        if dash:
            shp.line.dash_style = 4
    shp.shadow.inherit = False
    return shp


def new_slide(eyebrow_txt, title_txt, lead_txt=None, notes=""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    if eyebrow_txt:
        add_text(s, ML, Inches(0.32), TITLE_W, Inches(0.3),
                 [(eyebrow_txt, 11, RED, True, {"spc": 80})])
    if title_txt:
        add_text(s, ML, Inches(0.62), TITLE_W, Inches(0.95),
                 [(title_txt, 22, WHITE, True, {"line": 1.18})])
        rect(s, ML, Inches(1.42), Inches(1.2), Pt(3.2), RED)
    if lead_txt:
        add_text(s, ML, Inches(1.6), CW, Inches(0.4),
                 [(lead_txt, 13, BLUE, True)])
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


ASPECT = {}
def aspect(path):
    if path not in ASPECT:
        with Image.open(os.path.join("assets", path)) as im:
            ASPECT[path] = im.width / im.height
    return ASPECT[path]


def item_grid(slide, items, y0, n_cols, cell_h, img_h, label_size=9.5, gap=0.18):
    """items: [(番号, ラベル, カンプパス or None, 注記 or None)] のカードグリッド."""
    cw = (12.23 - gap * (n_cols - 1)) / n_cols
    for i, (num, label, img, note) in enumerate(items):
        col, row = i % n_cols, i // n_cols
        x = ML + Inches(col * (cw + gap))
        y = y0 + row * (cell_h + gap)
        rect(slide, x, Inches(y), Inches(cw), Inches(cell_h), PANEL, LINE, Pt(0.75),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        if img:
            a = aspect(img)
            iw = min(img_h * a, cw - 0.16)
            ih = iw / a
            slide.shapes.add_picture(os.path.join("assets", img),
                                     x + Inches((cw - iw) / 2), Inches(y + 0.08),
                                     Inches(iw), Inches(ih))
        else:
            rect(slide, x + Inches(0.08), Inches(y + 0.08), Inches(cw - 0.16),
                 Inches(img_h), PANEL2, LINE, Pt(0.75), dash=True)
            add_text(slide, x + Inches(0.08), Inches(y + 0.08), Inches(cw - 0.16),
                     Inches(img_h), [((note or "見本なし(内容は下記)"), 8.5, DIM, False,
                                     {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
        d = 0.26
        rect(slide, x + Inches(0.08), Inches(y + img_h + 0.14), Inches(d), Inches(d),
             RED, shape=MSO_SHAPE.OVAL)
        add_text(slide, x + Inches(0.08), Inches(y + img_h + 0.14), Inches(d), Inches(d),
                 [(str(num), 9, WHITE, True, {"align": PP_ALIGN.CENTER})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.42), Inches(y + img_h + 0.1), Inches(cw - 0.5),
                 Inches(cell_h - img_h - 0.14),
                 [(label, label_size, WHITE, True, {"line": 1.1})])


def footer_deadline(slide, text):
    band_h = Inches(0.5)
    rect(slide, ML, Inches(6.55), CW, band_h, PANEL)
    rect(slide, ML, Inches(6.55), Pt(4), band_h, RED)
    add_text(slide, ML + Inches(0.22), Inches(6.55), CW - Inches(0.44), band_h,
             [(text, 12.5, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)


# ============ 1. タイトル ============
s = new_slide(None, None)
rect(s, 0, 0, Inches(0.28), SLIDE_H, RED)
add_text(s, Inches(0.9), Inches(1.5), Inches(11.3), Inches(0.4),
         [("素材ご提供のお願い", 14, BLUE, True, {"spc": 200})])
add_text(s, Inches(0.9), Inches(2.0), Inches(11.6), Inches(1.6),
         [("依頼画像素材リスト", 40, WHITE, True, {"line": 1.1, "after": 0}),
          ("全49点+確認事項", 40, WHITE, True, {"line": 1.1})])
add_text(s, Inches(0.9), Inches(3.85), Inches(11.6), Inches(0.5),
         [("ULTRAMAN BUSINESS PARTNERS MEETING 2026(7/23・約450名)投影資料用", 16, RED, True)])
add_text(s, Inches(0.9), Inches(5.1), Inches(11.6), Inches(1.2),
         [("各素材に「欲しい絵」の見本(仮ビジュアル)を付けています。", 13, SUB, False, {"after": 0}),
          ("見本と同じ構図の実物写真・実案がありましたら、ご提供をお願いします。", 13, SUB, False)])
footer_deadline(s, "写真系:7/11(土)/モック類:7/12(日)/ガイド類:7/15(水)|納品先:" + DRIVE_URL)

# ============ 2. 共通ルール ============
s = new_slide("RULES", "共通ルール", "全素材に共通するお願い")
rules = [
    ("サイズ", "長辺1600px以上(最低1200px)。全面背景用のみ 3840×2160"),
    ("形式", "写真=JPEG高品質/ロゴ・切り抜き=PNG背景透過"),
    ("点数", "各項目につき候補3〜5点(採用はこちらで選定します)"),
    ("NG", "動画からの静止画切り出しは不可(解像度不足)"),
    ("選び方", "「商品単体」より「売場・人・熱量」が写っている絵を優先"),
    ("開示区分", "[1]投影可 [2]配布可 [3]NDA後開示 [4]開示不可 [5]権利確認中 を各素材に付記"),
    ("回答方法", "番号に○×と区分を付けるだけでOK(例:1○[1]、2×)"),
]
y = 2.0
for k, v in rules:
    rect(s, ML, Inches(y), Inches(1.7), Inches(0.52), REDDK,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
    add_text(s, ML, Inches(y), Inches(1.7), Inches(0.52),
             [(k, 12, WHITE, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ML + Inches(1.95), Inches(y), Inches(10.2), Inches(0.52),
             [(v, 12.5, SUB, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.63
footer_deadline(s, "納品先フォルダ:" + DRIVE_URL + " (02_photo/03_mock/04_kv)")

# ============ 3. カテゴリ1a 行動カット ============
s = new_slide("CATEGORY 1|イベント・ファンの写真(1/2)",
              "ファンの「行動」が分かるカット|5点",
              "行動の連鎖(来る→撮る→手に取る→共有する)がそのまま資料になります")
item_grid(s, [
    (1, "親子がヒーローを見上げている", "photo/s03_01_arrive_M.png", None),
    (2, "スマホで撮影している手元", "photo/s03_02_shoot_M.png", None),
    (3, "売場で子どもが商品に手を伸ばす", "photo/s03_03_pickup_M.png", None),
    (4, "SNS投稿画面(許諾済み)", "photo/s03_04_share_M.png", None),
    (5, "グリーティングで握手", "photo/s03_05_greeting_M.png", None),
], 2.1, 5, 3.15, 2.2, gap=0.16)
footer_deadline(s, "締切 7/11(土)|納品先:02_photo フォルダ")

# ============ 4. カテゴリ1b 熱量+施設 ============
s = new_slide("CATEGORY 1|イベント・ファンの写真(2/2)",
              "会場の「熱量」+施設・集客の実績カット|8点",
              "客席の埋まり・行列・賑わい——「人が動いている」ことが分かる絵を")
item_grid(s, [
    (6, "ウルサマ客席の引き画", "photo/s04_01_ulsama_M.png", None),
    (7, "ショーの決め場面+観客", "photo/s12_01_heroshow_M.png", None),
    (8, "物販の行列", None, "見本なし:列の長さが分かる引き画"),
    (9, "会場限定商品+購入者", None, "見本なし:手渡しの瞬間・手元中心"),
    (10, "撮影会の様子", None, "見本なし"),
    (11, "館内スタンプラリー", "photo/s12_03_stamprally_M.png", None),
    (12, "商業施設・観光地の賑わい", "photo/s12_04_mall-event_M.png", None),
    (13, "地域イベント(フリーショー可)", "photo/s04_02_heroshow_M.png", None),
], 2.1, 4, 2.15, 1.35, gap=0.16)
footer_deadline(s, "締切 7/11(土)|8・9はウルサマ会場での撮り下ろしリストを別途お渡しします")

# ============ 5. カテゴリ2 実績写真 ============
s = new_slide("CATEGORY 2|実績写真", "商品化・コラボ・広告・海外の実績|8点",
              "「すでに動いている」ことの証明に使います")
item_grid(s, [
    (14, "モフサンドコラボ商品", "photo/s13_01_cute-collab_S.png", None),
    (15, "ベイブレードコラボ商品(×回答済)", "photo/s13_02_hobby-collab_S.png", None),
    (16, "コラボ・一般商品の売場棚", "photo/s13_05_collab-shelf_S.png", None),
    (17, "商品化事例:アパレル", "mock/s07_01_apparel-main_M.png", None),
    (18, "商品化事例:ホビー・カード", "mock/s07_02_acsta-cards_S.png", None),
    (19, "商品化事例:プレミアム", "mock/s07_03_premium-figure_S.png", None),
    (20, "広告タイアップ掲出", "photo/s04_07_ad_M.png", None),
    (21, "海外イベント・展開", "photo/s04_08_global_M.png", None),
], 2.1, 4, 2.15, 1.35, gap=0.16)
footer_deadline(s, "締切 7/11(土)|16はIPコラボ限定ではなく一般のウルトラマン売場も歓迎です")

# ============ 6. カテゴリ3 マート ============
s = new_slide("CATEGORY 3|ウルトラマート・売場写真", "売場の「出口」を見せるカット|6点",
              "「作っても売る場所があるのか」への答えになります(全点○回答済・納品待ち)")
item_grid(s, [
    (22, "外観(看板が分かるもの)", "photo/s14_01_storefront_M.png", None),
    (23, "店内の賑わい(客入り)", "photo/s14_02_interior_M.png", None),
    (24, "商品棚クローズアップ", "photo/s04_06_retail_M.png", None),
    (25, "限定商品コーナー", "photo/s14_03_limited-corner_M.png", None),
    (26, "POPUP出店時の売場", "photo/s14_04_popup_M.png", None),
    (27, "EC特集ページのスクショ", "mock/s09_01_shelf-ec-main_M.png", None),
], 2.1, 3, 2.15, 1.35, gap=0.18)
footer_deadline(s, "締切 7/11(土)|納品先:02_photo フォルダ")

# ============ 7. カテゴリ4 スタイルガイド ============
s = new_slide("CATEGORY 4|スタイルガイド・監修資料", "「作れる」安心感を出す資料|5点",
              "当日投影可能な範囲の抜粋で結構です(31・32はドラフト作成済み・監修のみお願いします)")
item_grid(s, [
    (28, "ロゴ使用例+カラーパレット", "mock/s15_01_guide-logo_M.png", None),
    (29, "キャラ配置例+OK/NG例", "mock/s15_02_guide-okng_M.png", None),
    (30, "商品展開例(Tシャツ・PKG・POP)", "mock/s15_03_guide-products_M.png", None),
], 2.1, 3, 2.3, 1.5, gap=0.18)
add_text(s, ML, Inches(4.7), CW, Inches(0.4),
         [("31 監修フロー図/32 標準期間・チェックリスト → 当方ドラフトを送付済み。数字と要件の赤入れのみお願いします",
           12, GOLD, True)])
footer_deadline(s, "締切 7/15(水)|納品先:03_mock フォルダ")

# ============ 8. カテゴリ5 コンテンツビジュアル ============
s = new_slide("CATEGORY 5|コンテンツビジュアル", "2027年度以降を見せるビジュアル|5点",
              "可否のご回答を特に急ぎます(黒澤さん確認中の項目あり)")
rows5 = [
    ("33", "ヒーロー集合ビジュアル", "60周年KVから流用でOK。元データ(理想3840×2160)をご共有ください"),
    ("34", "2027新作関連", "ゼロ映画の場面写真+キャラデザイラストの方針で進行(黒澤さん確認中)"),
    ("35", "ゼロ・映画関連", "ポスター・KVなし→場面写真で構成(黒澤さん確認中)"),
    ("36", "配信・YouTube展開", "配信作品のKVでOKです"),
    ("37", "歴代シリーズ・歴代ヒーロー", "TVシリーズ紹介用(○回答済・納品待ち)"),
]
y = 2.05
for num, t, note in rows5:
    rect(s, ML, Inches(y), CW, Inches(0.78), PANEL, LINE, Pt(0.75),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    d = 0.3
    rect(s, ML + Inches(0.14), Inches(y + 0.24), Inches(d), Inches(d), RED,
         shape=MSO_SHAPE.OVAL)
    add_text(s, ML + Inches(0.14), Inches(y + 0.24), Inches(d), Inches(d),
             [(num, 9.5, WHITE, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ML + Inches(0.6), Inches(y + 0.09), Inches(4.1), Inches(0.6),
             [(t, 13, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ML + Inches(4.9), Inches(y + 0.09), Inches(7.1), Inches(0.6),
             [(note, 11, SUB, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.89
footer_deadline(s, "可否回答 7/11(土)まで|納品先:04_kv フォルダ")

# ============ 9. カテゴリ6a 商品モック ============
s = new_slide("CATEGORY 6|モック類(1/2)", "商品モック|5点",
              "仮ビジュアル組込済み——実物・実案・過去実績があるものだけで結構です")
item_grid(s, [
    (38, "Tシャツ・スウェット・キャップ", "mock/s07_01_apparel-main_M.png", None),
    (39, "アクスタ・カード・フィギュア", "mock/s07_02_acsta-cards_S.png", None),
    (40, "プレミアムフィギュア", "mock/s07_03_premium-figure_S.png", None),
    (41, "限定パッケージ(コラボPKG可)", "mock/s08_01_package-main_M.png", None),
    (42, "購入特典カード(×回答済)", "mock/s08_02_bonus-cards_S.png", None),
], 2.1, 5, 2.5, 1.55, gap=0.16)
footer_deadline(s, "締切 7/12(日)|納品先:03_mock フォルダ")

# ============ 10. カテゴリ6b 売場・広告・デジタル ============
s = new_slide("CATEGORY 6|モック類(2/2)", "売場・広告・デジタル|7点",
              "過去事例の写真・画面キャプチャで結構です")
item_grid(s, [
    (43, "映画公開記念棚(シン・ウルトラマン'22確認中)", "mock/s09_01_shelf-ec-main_M.png", None),
    (44, "POPUP売場(過去事例可)", "mock/s09_02_popup_S.png", None),
    (45, "スタンプラリー台紙(過去事例可)", "mock/s09_03_stamp-rally_S.png", None),
    (46, "交通広告(過去事例可)", "mock/s10_01_station-ad-main_M.png", None),
    (47, "SNSキャンペーン投稿(過去事例可)", "mock/s10_02_sns_S.png", None),
    (48, "ゲーム内コラボ・ガチャ(過去事例可)", "mock/s11_01_game-main_M.png", None),
    (49, "YouTubeサムネ・配信画面", "mock/s11_03_youtube_S.png", None),
], 2.1, 4, 2.15, 1.35, gap=0.16)
footer_deadline(s, "締切 7/12(日)|納品先:03_mock フォルダ")

# ============ 11. 画像以外の確認事項+納品先 ============
s = new_slide("OTHERS", "画像以外の確認事項+納品先", None)
others = [
    ("a", "来場データ(満席率・稼働率・物販購買率など)", "まず提供可否のみ。数字の強さで当日の表現を決めます"),
    ("b", "2026後半〜2027の商機カレンダー情報", "時期×施策の一覧"),
    ("c", "登壇者が話してよい表現/NG表現", "制作部・IP戦略さま"),
    ("d", "コラボ先への掲出許諾確認", "モフサンド・ベイブレード・サンリオ等。外部確認のため本日着手をお願いします"),
    ("e", "来場者が写る写真の肖像権確認", "顔が判別できる写真の掲出許諾"),
]
y = 1.85
for num, t, note in others:
    rect(s, ML, Inches(y), CW, Inches(0.66), PANEL, LINE, Pt(0.75),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    rect(s, ML + Inches(0.14), Inches(y + 0.18), Inches(0.3), Inches(0.3), REDDK,
         shape=MSO_SHAPE.OVAL)
    add_text(s, ML + Inches(0.14), Inches(y + 0.18), Inches(0.3), Inches(0.3),
             [(num, 10, WHITE, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ML + Inches(0.6), Inches(y + 0.06), Inches(4.6), Inches(0.54),
             [(t, 12, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ML + Inches(5.4), Inches(y + 0.06), Inches(6.6), Inches(0.54),
             [(note, 10.5, SUB, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.77
rect(s, ML, Inches(5.9), CW, Inches(0.55), PANEL2, GOLD, Pt(1))
add_text(s, ML + Inches(0.2), Inches(5.9), CW - Inches(0.4), Inches(0.55),
         [("納品先:https://" + DRIVE_URL + "(02_photo/03_mock/04_kv)|ファイル名の頭に素材番号と用途を記載",
           11.5, GOLD, True)], anchor=MSO_ANCHOR.MIDDLE)
footer_deadline(s, "ご回答は番号+○×+開示区分でメール返信ください。あるものから順次で構いません!")

OUT = "素材依頼_ビジュアルショットリスト_UBPM2026.pptx"
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides")
