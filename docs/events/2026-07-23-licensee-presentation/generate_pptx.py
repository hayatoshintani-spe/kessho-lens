#!/usr/bin/env python3
"""ULTRAMAN BUSINESS PARTNERS MEETING 2026 — デッキ生成スクリプト v5.

公式テンプレート(template_PARTNERS_MEETING2026.pptx:濃紺黒背景・
右上ULTRAMAN 60thロゴ・右下ページ番号)を土台に、21枚構成で生成する。
設計書は 11-presentation-spec.md。コンセプト:商談起点イベント。
1スライド1メッセージ/ビジュアル中心/全スライド下部に強いコピー。

    python3 generate_pptx.py
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

# ---- design tokens(公式テンプレートのダーク基調)----
BG    = RGBColor(0x00, 0x07, 0x14)   # テンプレート背景色
WHITE = RGBColor(0xF4, 0xF7, 0xFD)   # タイトル・強調
SUB   = RGBColor(0xA9, 0xB6, 0xD2)   # 本文
DIM   = RGBColor(0x6A, 0x77, 0x94)   # 注記
RED   = RGBColor(0xE8, 0x38, 0x2F)   # アクセント(ウルトラマンレッド)
REDDK = RGBColor(0xA8, 0x1C, 0x18)   # 濃赤(表ヘッダー・強カード)
PANEL = RGBColor(0x0C, 0x17, 0x30)   # パネル
PANEL2= RGBColor(0x11, 0x1F, 0x3E)   # パネル(明)
LINE  = RGBColor(0x25, 0x36, 0x5E)   # 罫線
BLUE  = RGBColor(0x7F, 0xA8, 0xE0)   # 淡青(リード・サブ強調)

JP_FONT = "Yu Gothic"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.55)
CW = Inches(12.23)
TITLE_W = Inches(9.55)   # 右上ロゴを避けるタイトル幅

TEMPLATE = "template_PARTNERS_MEETING2026.pptx"
prs = Presentation(TEMPLATE)
BLANK = prs.slide_layouts[6]   # 白紙(ロゴはマスターから継承)

# テンプレート同梱のサンプル2枚を削除
for sldId in list(prs.slides._sldIdLst):
    prs.part.drop_rel(sldId.rId)
    prs.slides._sldIdLst.remove(sldId)


def style_run(run, size, color=SUB, bold=False, spacing=None):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = JP_FONT
    f.language_id = MSO_LANGUAGE_ID.JAPANESE
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", JP_FONT)
    if spacing is not None:
        rPr.set("spc", str(spacing))


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for spec in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = spec if isinstance(spec, list) else [spec]
        opts = {}
        for rs in runs:
            text, size, color, bold = rs[0], rs[1], rs[2], rs[3]
            opts = rs[4] if len(rs) > 4 else opts
            run = p.add_run()
            run.text = text
            style_run(run, size, color, bold, spacing=opts.get("spc"))
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(opts.get("before", 0))
        p.space_after = Pt(opts.get("after", 3))
        p.line_spacing = opts.get("line", 1.25)
    return box


def rect(slide, x, y, w, h, fill, line_color=None, line_w=None, dash=False,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
        if dash:
            shp.line.dash_style = 4
    shp.shadow.inherit = False
    return shp


def new_slide(eyebrow_txt, notes=""):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    if eyebrow_txt:
        add_text(slide, ML, Inches(0.32), TITLE_W, Inches(0.3),
                 [(eyebrow_txt, 11, RED, True, {"spc": 80})])
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def title(slide, text, size=24):
    add_text(slide, ML, Inches(0.62), TITLE_W, Inches(0.95),
             [(text, size, WHITE, True, {"line": 1.18})])


def rule(slide, y=1.5):
    rect(slide, ML, Inches(y), Inches(1.2), Pt(3.2), RED)


def lead(slide, text, y=1.7, size=14):
    add_text(slide, ML, Inches(y), CW, Inches(0.45),
             [(text, size, BLUE, True)])


def bottom_copy(slide, text, y=6.42):
    """全スライド共通:下部の強いコピー帯."""
    band_h = Inches(0.56)
    rect(slide, ML, Inches(y), CW, band_h, PANEL)
    rect(slide, ML, Inches(y), Pt(4), band_h, RED)
    box = slide.shapes.add_textbox(ML + Inches(0.24), Inches(y), CW - Inches(0.48), band_h)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    style_run(r, 14, WHITE, bold=True)


def photo_zone(slide, x, y, w, h, label, items, size=9.5):
    """写真・モック差し込みゾーン(点線枠+必要素材リスト)."""
    rect(slide, x, y, w, h, PANEL, LINE, Pt(1.2), dash=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.16),
                                   w - Inches(0.44), h - Inches(0.32))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.space_after = Pt(5)
    r = p.add_run(); r.text = "📷 " + label
    style_run(r, 10.5, RED, bold=True)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.35
    r2 = p2.add_run(); r2.text = items
    style_run(r2, size, DIM)


def cards(slide, items, y, h=1.5, n_cols=None, title_size=13.5, body_size=10,
          gap_in=0.22, x=ML, total_w=None):
    n = n_cols or len(items)
    gap = Inches(gap_in)
    tw = total_w or CW
    w = Emu(int((tw - gap * (n - 1)) / n))
    for i, (t, b, variant) in enumerate(items):
        col = i % n
        row = i // n
        xx = x + Emu(int((w + gap) * col))
        yy = Inches(y) + Emu(int((Inches(h) + gap) * row))
        if variant == "red":
            fill, border, tc, bc = REDDK, None, WHITE, RGBColor(0xF3, 0xC9, 0xC5)
        else:
            fill, border, tc, bc = PANEL2, LINE, WHITE, SUB
        shp = rect(slide, xx, yy, w, Inches(h), fill, border,
                   Pt(1) if border else None,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(3)
        r = p.add_run(); r.text = t
        style_run(r, title_size, tc, bold=True)
        if b:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.2
            r2 = p2.add_run(); r2.text = b
            style_run(r2, body_size, bc)


def table_block(slide, headers, rows, y, col_w, size=11, row_h=0.44):
    n_rows, n_cols = len(rows) + 1, len(headers)
    total_w = Emu(sum(int(c) for c in col_w))
    gfx = slide.shapes.add_table(n_rows, n_cols, ML, Inches(y), total_w,
                                 Inches(row_h * n_rows))
    tbl = gfx.table
    tbl.first_row = False; tbl.horz_banding = False
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Emu(int(w))
    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = REDDK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        style_run(r, size, WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri % 2 else PANEL2
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.1
            r = p.add_run(); r.text = val
            style_run(r, size, WHITE if c == 0 else SUB, bold=(c == 0))


def plan_rows(slide, items, y, row_h=0.5, gap=0.12, size=12.5, w=Inches(6.2)):
    """カテゴリースライド左側:企画リスト(赤丸番号+白文字)."""
    for i, text in enumerate(items):
        yy = Inches(y) + Emu(int((Inches(row_h) + Inches(gap)) * i))
        d = 0.3
        circ = rect(slide, ML, yy + Inches((row_h - d) / 2), Inches(d), Inches(d),
                    RED, shape=MSO_SHAPE.OVAL)
        tfc = circ.text_frame
        tfc.margin_left = 0; tfc.margin_right = 0
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run(); rc.text = str(i + 1)
        style_run(rc, 10.5, WHITE, bold=True)
        box = slide.shapes.add_textbox(ML + Inches(0.45), yy, w - Inches(0.45), Inches(row_h))
        tf = box.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0; tf.margin_right = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        style_run(r, size, WHITE, bold=True)


def category_slide(no, audience, package, main_copy, plans, materials,
                   bottom, notes):
    s = new_slide(f"{no:02d}|FOR {audience}", notes)
    title(s, f"{audience}向け|{package}", size=22)
    rule(s)
    lead(s, main_copy)
    plan_rows(s, plans, 2.4)
    photo_zone(s, Inches(7.1), Inches(2.3), Inches(5.68), Inches(3.85),
               "ビジュアルゾーン(モック・写真)", materials)
    bottom_copy(s, bottom)
    return s


NOTE_QR = "(QRは以降のスライドでも画面隅に残す)"

# =====================================================================
# Slide 1|タイトル
s = new_slide(None,
              "皆さま、本日はお集まりいただきありがとうございます。最初に申し上げます。本日は、作品の"
              "発表会ではありません。2027年度以降のウルトラマンの商機を、皆さまと一緒につくり始める"
              "1時間です。")
rect(s, 0, 0, Inches(0.28), SLIDE_H, RED)
add_text(s, Inches(0.9), Inches(1.5), Inches(11.3), Inches(0.4),
         [("TSUBURAYA PRODUCTIONS", 13, BLUE, True, {"spc": 300})])
add_text(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(1.7),
         [("ULTRAMAN BUSINESS", 42, WHITE, True, {"line": 1.08, "after": 0}),
          ("PARTNERS MEETING 2026", 42, WHITE, True, {"line": 1.08})])
add_text(s, Inches(0.9), Inches(3.85), Inches(11.6), Inches(0.6),
         [("2027年度以降の商機を、いま一緒につくる", 22, RED, True)])
add_text(s, Inches(0.9), Inches(5.35), Inches(11.3), Inches(0.35),
         [("2026年7月23日(木)16:00-17:00|ウルサマ イベントステージ", 12, SUB, False)])
add_text(s, Inches(0.9), Inches(5.7), Inches(11.3), Inches(0.35),
         [("円谷プロダクション", 12, SUB, True)])
photo_zone(s, Inches(9.0), Inches(4.9), Inches(3.75), Inches(1.9),
           "背景コラージュ",
           "キービジュアル/ヒーロー集合/60周年ロゴ/イベント・物販・売場・映画館・商品・親子ファン")
bottom_copy(s, "2027年度以降の売場は、今日ここから始まる。")

# Slide 2|本日の目的
s = new_slide("02|PURPOSE",
              "本日の目的は3つです。2027年度以降の展開を共有すること。皆さまの商品化・販促・売場展開の"
              "可能性をご提示すること。そして、本日から個別商談と企画相談を開始することです。"
              "ぜひ「自社なら何ができるか」という目線でお聞きください。")
title(s, "本日は、2027年度以降の商機を共有する場です")
rule(s)
lead(s, "「作品を聞く」のではなく、「自社の商談機会」として")
cards(s, [
    ("① 展開を共有する", "2027年度以降のウルトラマン展開", "panel"),
    ("② 可能性を提示する", "商品化・販促・売場展開の具体案", "panel"),
    ("③ 商談を開始する", "本日から個別商談・企画相談を受付", "red"),
], y=2.6, h=1.9, title_size=15, body_size=11)
bottom_copy(s, "作品を聞く1時間ではなく、商談が始まる1時間です。")

# Slide 3|人を動かすIP
s = new_slide("03|WHY ULTRAMAN",
              "いまご覧いただいた熱量が、ウルトラマンの原点です。ファンは会場に来て、写真を撮り、商品を"
              "手に取り、誰かに共有します。この「人が動く力」を、皆さまの商品・売場・キャンペーンに接続"
              "するのが、本日のテーマです。")
title(s, "ウルトラマンは、“人を動かすIP”")
rule(s)
lead(s, "ファンの行動が、そのまま売上への動線になる")
chain = ["会場に来る", "写真を撮る", "商品を手に取る", "誰かに共有する"]
for i, t in enumerate(chain):
    x = ML + Inches(i * 3.12)
    shp = rect(s, x, Inches(2.5), Inches(2.7), Inches(1.05), PANEL2, LINE, Pt(1),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    style_run(r, 15, WHITE, bold=True)
    if i < 3:
        add_text(s, x + Inches(2.72), Inches(2.62), Inches(0.4), Inches(0.7),
                 [("→", 20, RED, True, {"align": PP_ALIGN.CENTER})])
photo_zone(s, ML, Inches(3.95), CW, Inches(2.2),
           "熱量ビジュアル(横並び4-5点)",
           "親子がヒーローを見ている写真/グリーティング/物販列/商品購入シーン/SNS投稿イメージ")
bottom_copy(s, "この熱量を、商品・売場・キャンペーンへ。")

# Slide 4|60周年は、もう動き出している
s = new_slide("04|MOMENTUM",
              "60周年は、すでに動き出しています。イベント、商品化、コラボ、流通、広告、海外。ご覧いただき"
              "たいのは商品そのものではなく、そこに集まっている「人」です。この勢いは実績として証明されて"
              "います。")
title(s, "60周年は、もう動き出している")
rule(s)
lead(s, "勢いは、証明済み。")
moves = ["ウルサマ", "ヒーローショー", "商品化", "IPコラボ",
         "ウルトラマート", "流通展開", "広告タイアップ", "海外展開"]
cards(s, [(m, "", "panel") for m in moves], y=2.35, h=0.72, n_cols=4,
      title_size=13, gap_in=0.18)
photo_zone(s, ML, Inches(4.25), CW, Inches(1.95),
           "熱量ビジュアル(人・売場中心)",
           "人が集まっている写真/売場の熱量がある写真を使用(商品単体写真は使わない)— ウルサマ会場・物販列・コラボ棚・海外イベント")
bottom_copy(s, "この勢いは、2027年度以降につながっていく。")

# Slide 5|商機カレンダー
s = new_slide("05|BUSINESS CALENDAR",
              "ライセンシーの皆さまが一番知りたいのは「いつ仕込めばいいか」だと思います。答えはこの"
              "カレンダーです。2027年夏の映画・イベント・大型キャンペーンが最大の山です。商品開発の"
              "リードタイムから逆算すると——2026年冬には企画と監修が始まっている必要があります。"
              "つまり、商談開始は今です。")
title(s, "2027年の売場は、今日から仕込む")
rule(s)
lead(s, "2026-2027 ULTRAMAN 商機カレンダー")
table_block(s,
            ["時期", "商機"],
            [["2026夏", "ウルサマ・イベント連動"],
             ["2026秋", "IPコラボ・年末商戦"],
             ["2026冬", "2027商品の企画・監修開始"],
             ["2027春", "新生活・文具・雑貨・アパレル"],
             ["2027夏", "映画・イベント・大型キャンペーン"],
             ["2027秋冬", "ホビー・EC・プレミアム・年末商戦"]],
            2.3, [Inches(2.2), Inches(10.0)], size=12, row_h=0.52)
bottom_copy(s, "“発売タイミング”から逆算すると、商談開始は今です。")

# Slide 6|7カテゴリーのハブ図
s = new_slide("06|FOR ALL PARTNERS",
              "ここからは、本日お越しの皆さま一人ひとりに関わる話です。7つのカテゴリーで、具体的な企画を"
              "ご提案します。ご自身のカテゴリーが出てきたら、「この枠にうちが入れるか」という目線で"
              "ご覧ください。")
title(s, "本日ご参加の皆さまと、作りたいこと")
rule(s)
lead(s, "7つのカテゴリーで、企画を募集します")
center = rect(s, Inches(5.42), Inches(3.35), Inches(2.5), Inches(1.35), REDDK,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
tfc = center.text_frame; tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
rc = pc.add_run(); rc.text = "ULTRAMAN"
style_run(rc, 17, WHITE, bold=True)
pc2 = tfc.add_paragraph(); pc2.alignment = PP_ALIGN.CENTER
rc2 = pc2.add_run(); rc2.text = "円谷プロダクション"
style_run(rc2, 10.5, RGBColor(0xF3, 0xC9, 0xC5), bold=False)
hub = [
    ("MD・グッズ・玩具・アパレル", Inches(0.7),  Inches(2.35)),
    ("食品・飲料・外食",           Inches(5.02), Inches(2.2)),
    ("流通・EC・商業施設",         Inches(9.35), Inches(2.35)),
    ("広告代理店・SP",            Inches(0.7),  Inches(4.0)),
    ("デジタル・ゲーム",           Inches(9.35), Inches(4.0)),
    ("イベント・施設・観光",        Inches(2.6),  Inches(5.35)),
    ("IPコラボ",                 Inches(7.45), Inches(5.35)),
]
for t, x, y in hub:
    shp = rect(s, x, y, Inches(3.3), Inches(0.8), PANEL2, LINE, Pt(1),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    style_run(r, 12.5, WHITE, bold=True)
bottom_copy(s, "この中に、貴社の席があります。")

# Slide 7-13|カテゴリー別
category_slide(7, "MD・グッズ・玩具・アパレル", "2027新展開 商品化パッケージ",
               "2027年の初速を、一緒に作る5つの企画",
               ["新作連動キッズMD", "映画・ゼロ関連商品", "歴代ヒーロー/怪獣コレクション",
                "大人ファン向けプレミアム商品", "女性ライト層向け雑貨"],
               "Tシャツ/スウェット/キャップ/アクスタ/フィギュア/カード/文具/雑貨/プレミアムフィギュア/怪獣柄パターン/ゼロ関連素材/スタイルガイド抜粋/売場イメージ",
               "2027年の初速商品は、今日から仕込みが始まります。",
               "まず、最大のカテゴリーであるMD・グッズの皆さまへ。2027年の新展開に向けて、5つの商品化"
               "パッケージをご用意しました。キッズMDから、大人向けプレミアム、女性ライト層向けまで。"
               "新作の初速商品は、発売日から逆算すると今日から仕込みが始まります。")
category_slide(8, "食品・飲料・外食", "親子キャンペーンパッケージ",
               "親子需要 × 映画・イベントの山",
               ["限定パッケージ", "購入特典カード", "映画半券キャンペーン",
                "親子ペア応募", "外食コラボメニュー", "ウルサマ来場者特典"],
               "菓子・飲料・アイス・ゼリー・パンのパッケージモック/キッズプレート/コラボドリンク/ノベルティ/親子写真",
               "親子の熱量を、購入・来店・応募につなげる。",
               "食品・飲料・外食の皆さまには、親子キャンペーンパッケージをご提案します。映画やイベントの"
               "山に合わせて、限定パッケージ、半券キャンペーン、コラボメニュー。ウルトラマンの親子熱量を、"
               "購入・来店・応募という行動に変える企画です。")
category_slide(9, "流通・EC・商業施設", "ULTRA MART / 売場連動パッケージ",
               "「商品を置く」ではなく「売場を作る」",
               ["映画公開記念棚/ウルサマ連動棚", "ゼロ特集棚/怪獣特集棚",
                "IPコラボ棚", "POPUP/EC限定販売", "購入特典"],
               "棚モック/POPUPイメージ/ECページモック/店頭POP/棚帯/購入特典告知/館内装飾/スタンプラリー",
               "商品を作るだけでなく、売場まで一緒に設計する。",
               "流通・EC・商業施設の皆さまには、「商品を置く」のではなく「売場を作る」提案です。映画棚、"
               "ウルサマ棚、怪獣棚、IPコラボ棚。期間限定の売場ジャックからPOPUP、EC限定まで、売場その"
               "ものを一緒に設計させてください。")
category_slide(10, "広告代理店・SP", "企業キャンペーン活用パッケージ",
               "クライアント提案に、そのまま使える",
               ["親子集客/夏休み販促", "防災/交通安全", "地域創生",
                "未来・テクノロジー/CSR", "商業施設回遊"],
               "店頭キャンペーン事例/交通広告モック/SNS投稿施策/ヒーローグリーティング写真/スタンプラリー台紙/企業タイアップ動画イメージ",
               "ウルトラマンは、企業課題を解決するキャンペーン装置になる。",
               "広告代理店・SPの皆さまには、クライアント提案にそのまま使えるパッケージをご用意します。"
               "親子集客、防災、交通安全、地域創生。ウルトラマンは「正義」「勇気」「未来」という文脈を"
               "持っているので、企業課題と接続しやすいIPです。")
category_slide(11, "デジタル・ゲーム", "デジタル・ゲーム連動パッケージ",
               "リアルIPの熱量を、デジタルへ",
               ["ゲーム内コラボ/限定アバター", "デジタルスタンプ", "SNSキャンペーン",
                "配信連動/EC限定デジタル特典", "ファンコミュニティ施策"],
               "ゲーム画面風モック/ガチャ画面/限定キャラクター/デジタルアイテム/SNS投稿テンプレート/YouTubeサムネイル風/ライブ配信画面",
               "リアル商品だけでなく、デジタル接点でもファンを広げる。",
               "デジタル・ゲームの皆さまへ。ウルトラマンはリアルだけのIPではありません。ゲーム内コラボ、"
               "限定アバター、配信連動。歴代ヒーローと怪獣というキャラクター資産は、デジタル施策と非常に"
               "相性が良い領域です。")
category_slide(12, "イベント・施設・観光", "会いに来るIP 集客パッケージ",
               "ウルトラマンは、来場理由になる",
               ["ヒーローショー/グリーティング", "撮影会", "館内スタンプラリー/地域回遊",
                "商業施設イベント/観光地コラボ", "物販連動"],
               "ヒーローショー写真/親子観覧写真/決めポーズ写真/グリーティング写真/館内マップ/スタンプラリー台紙/地域イベント写真",
               "来場理由を作り、購買理由につなげる。",
               "イベント・施設・観光の皆さまへ。ウルトラマンの最大の特長は「会いに来てもらえる」ことです。"
               "ヒーローショー、グリーティング、スタンプラリー。来場理由を作り、そのまま館内の回遊と購買に"
               "つなげる集客パッケージをご提案します。")

# Slide 13|IPコラボ
s = new_slide("13|FOR IP COLLABORATION",
              "IPコラボは、ウルトラマンの入口を広げる装置です。かわいい文脈で女性・ライト層へ。ホビー文脈で"
              "コレクターへ。すでにモフサンド、ベイブレードで実証されています。コラボを起点に、新しい売場と"
              "新しいファン層を一緒に開拓しませんか。")
title(s, "IPコラボで、新しいファン層へ")
rule(s)
lead(s, "ウルトラマンは、男性・キッズ向けだけではない")
table_block(s,
            ["方向性", "狙う層", "商品化機会"],
            [["かわいい文脈", "女性・ライト層", "雑貨・文具"],
             ["ホビー文脈", "キッズ・コレクター", "玩具・カード"],
             ["キャラクター文脈", "ファミリー・ギフト", "流通棚・雑貨"],
             ["ヒーロー文脈", "コアファン", "プレミアム・アパレル"]],
            2.3, [Inches(3.2), Inches(3.6), Inches(5.4)], size=11.5, row_h=0.5)
photo_zone(s, ML, Inches(5.0), CW, Inches(1.2),
           "コラボ実績・イメージ",
           "モフサンド/ベイブレード/サンリオ(開示可能なら)/コラボ商品写真/コラボ棚モック/新規ファン層マップ")
bottom_copy(s, "ウルトラマンの入口を、もっと広げる。")

# Slide 14|ULTRA MART
s = new_slide("14|ULTRA MART",
              "「作った商品はどこで売れるのか」。この不安にお応えするのがULTRA MARTです。売る場所であり、"
              "見せる場所であり、話題化する場所。皆さまの商品の出口を、円谷側でも用意していきます。")
title(s, "ULTRA MART|売場の出口")
rule(s)
lead(s, "ULTRA MARTの、3つの役割")
cards(s, [
    ("① 売る場所", "商品の販売拠点・EC連動", "panel"),
    ("② 見せる場所", "限定コーナー・POPUP展開", "panel"),
    ("③ 話題化する場所", "映画・イベント連動の情報発信", "red"),
], y=2.45, h=1.6, title_size=15, body_size=11)
photo_zone(s, ML, Inches(4.45), CW, Inches(1.75),
           "ULTRA MART ビジュアル",
           "ウルトラマート外観/店内写真/商品棚/限定商品コーナー/POPUP展開イメージ/EC連動ページ/イベント物販との連動図")
bottom_copy(s, "商品を、ファンと出会う場所へ。")

# Slide 15|スタイルガイド戦略
s = new_slide("15|STYLE GUIDE",
              "「企画したいが、素材はあるのか」。あります。しかも、ただ渡すのではなく、カテゴリー別・"
              "ターゲット別に商品化しやすい形に整備していきます。キッズ、大人、女性、怪獣、映画連動、流通、"
              "広告。貴社の企画に合うガイドを用意します。")
title(s, "スタイルガイド戦略")
rule(s)
lead(s, "カテゴリー別・ターゲット別に、「使える形」で")
guides = ["キッズ向け", "大人ファン向け", "女性・ライト層向け", "怪獣・ヴィラン向け",
          "映画連動向け", "流通販促向け", "広告キャンペーン向け"]
cards(s, [(g, "", "panel") for g in guides[:4]], y=2.4, h=0.75, n_cols=4, title_size=12.5, gap_in=0.18)
cards(s, [(g, "", "panel") for g in guides[4:]], y=3.35, h=0.75, n_cols=3, title_size=12.5, gap_in=0.18)
photo_zone(s, ML, Inches(4.4), CW, Inches(1.8),
           "スタイルガイド抜粋(見開き風)",
           "ロゴ使用例/キャラクター配置例/カラーパレット/背景パターン/OK・NG例/Tシャツ展開例/パッケージ展開例/POP展開例")
bottom_copy(s, "素材を渡すだけでなく、商品化しやすい形にする。")

# Slide 16|監修・素材提供フロー
s = new_slide("16|PROCESS",
              "進め方はシンプルです。個別相談から始まり、素材共有、企画提出、監修、商品化。各ステップの"
              "標準期間も明示します。「作りたい」で終わらせず、「作れる」状態まで私たちが伴走します。")
title(s, "監修・素材提供フロー")
rule(s)
lead(s, "相談から商品化まで、5ステップ")
steps = [("① 個別相談", ""), ("② 素材・ガイド共有", ""), ("③ 企画提出", ""),
         ("④ 監修・修正", ""), ("⑤ 商品化・販促展開", "")]
cards(s, [(t, b, "red" if i == 4 else "panel") for i, (t, b) in enumerate(steps)],
      y=2.55, h=1.15, n_cols=5, title_size=12, gap_in=0.16)
add_text(s, ML, Inches(4.05), CW, Inches(0.4),
         [("各ステップの標準期間・提出物チェックリスト・開示区分(当日/配布/NDA後)を個別商談時にご提示します",
           11, DIM, False)])
photo_zone(s, ML, Inches(4.6), CW, Inches(1.6),
           "フロー補足素材",
           "5ステップのフロー図/提出物チェックリスト/開示区分表/監修期間目安/商品化スケジュール例")
bottom_copy(s, "作りたい、で終わらせず、作れる状態まで伴走する。")

# Slide 17|本日から受付開始(QR)
s = new_slide("17|START TODAY",
              "ここで、本日一番大事なご案内です。いまこの瞬間から、個別商談の受付を開始します。お手元の"
              "スマートフォンで、このQRからご登録ください。相談カテゴリを選ぶだけ。登録いただいた企業様には、"
              "担当者から順にご連絡し、個別商談、素材共有、企画検討へ進みます。この後のコンテンツ紹介の間も、"
              "QRは画面の隅に残しておきます。【運営】QRは最低60秒表示。以降のスライドでも画面隅に残す。")
title(s, "本日から受付開始")
rule(s)
lead(s, "個別商談・企画相談、いま登録できます")
qr = rect(s, ML, Inches(2.3), Inches(3.7), Inches(3.7), PANEL2, RED, Pt(1.6),
          dash=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
tfq = qr.text_frame; tfq.word_wrap = True
tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
for qi, qline in enumerate(["【QRコード】", "個別商談・企画相談", "登録は約3分"]):
    pq = tfq.paragraphs[0] if qi == 0 else tfq.add_paragraph()
    pq.alignment = PP_ALIGN.CENTER
    rq = pq.add_run(); rq.text = qline
    style_run(rq, 16 if qi != 2 else 11, WHITE if qi != 2 else BLUE, bold=(qi != 2))
add_text(s, Inches(4.6), Inches(2.3), Inches(8.2), Inches(0.35),
         [("相談カテゴリ(複数選択可)", 12, RED, True)])
cats = ["商品化", "食品・外食", "流通・売場", "広告タイアップ", "イベント・施設",
        "IPコラボ", "デジタル", "スタイルガイド請求", "NDA後説明会"]
cards(s, [(c, "", "panel") for c in cats], y=2.75, h=0.55, n_cols=3,
      title_size=11, gap_in=0.14, x=Inches(4.6), total_w=Inches(8.18))
add_text(s, Inches(4.6), Inches(4.95), Inches(8.2), Inches(0.35),
         [("登録後の流れ", 12, RED, True)])
add_text(s, Inches(4.6), Inches(5.35), Inches(8.2), Inches(0.6),
         [("QR登録 → 担当者連絡(1週間以内)→ 個別商談 → 素材共有 → 企画検討",
           12.5, WHITE, True)])
bottom_copy(s, "早期にご相談いただいた企業様から、具体的な企画検討を進めます。")

# Slide 18|商機を支えるコンテンツ展開
s = new_slide("18|CONTENT",
              "最後に5分だけ、この商機を支えるコンテンツの話をします。ここまでの提案が「続けられる」理由です。")
title(s, "最後に、商機を支えるコンテンツ展開")
rule(s)
lead(s, "だから、商機は続く")
photo_zone(s, ML, Inches(2.35), CW, Inches(3.85),
           "コンテンツ接点ビジュアル(大)",
           "TV・映画・配信・イベントのアイコン/2027年度以降のコンテンツ接点図/ヒーロー集合ビジュアル/"
           "作品情報の開示可能素材/新作ビジュアル/ゼロ関連/配信系/映画館イメージ 【開示区分は7/14確認会で確定】")
bottom_copy(s, "コンテンツが、商品化・販促・売場展開の理由を作る。")

# Slide 19|3つのコンテンツ接点
s = new_slide("19|CONTENT",
              "2027年度以降の接点は3つ。TVシリーズが継続接点を作り、映画が大型の話題化を作り、配信・"
              "YouTubeがコアファン接点を広げます。この3つが年間を通じて連動するので、商品もキャンペーンも"
              "単発で終わりません。")
title(s, "2027年度以降の、3つのコンテンツ接点")
rule(s)
lead(s, "継続 × 話題化 × コアファン")
cards(s, [
    ("① TVシリーズ", "継続接点\nファンが追い続けたくなる世界観・キャラクター展開", "panel"),
    ("② 映画・大型映像", "大型話題化\n商品化・販促・流通施策を集中させる最大の山", "panel"),
    ("③ 配信・YouTube", "コアファン接点\n歴代ヒーロー・怪獣を横断活用できる領域", "panel"),
], y=2.5, h=2.1, title_size=15, body_size=11)
photo_zone(s, ML, Inches(4.95), CW, Inches(1.25),
           "各接点のビジュアル",
           "TVシリーズ素材/映画館・ポスター・ゼロ関連/配信サムネイル風/歴代ヒーロー/怪獣/コアファン向けビジュアル")
bottom_copy(s, "接点が続くから、商品もキャンペーンも続けられる。")

# Slide 20|コンテンツ×商機マトリクス
s = new_slide("20|CONTENT × BUSINESS",
              "この表が、本日のまとめです。コンテンツごとに、皆さまのカテゴリーの商品化機会が対応しています。"
              "作品の話をしているようで、実はずっと、皆さまの事業機会の話をしてきました。")
title(s, "コンテンツが、商品化の理由を作る")
rule(s)
lead(s, "コンテンツ × 商機マトリクス")
table_block(s,
            ["コンテンツ", "役割", "商品化機会"],
            [["TVシリーズ", "継続接点", "キッズ、文具、玩具"],
             ["映画", "大型話題化", "食品、流通、広告、ホビー"],
             ["配信・YouTube", "コアファン接点", "コレクション、EC、プレミアム"],
             ["イベント", "熱量化", "物販、限定商品、施設集客"],
             ["流通", "購買接点", "売場、POPUP、EC"]],
            2.3, [Inches(3.0), Inches(3.2), Inches(6.0)], size=11.5, row_h=0.5)
bottom_copy(s, "作品の話ではなく、皆さまの事業機会の話です。")

# Slide 21|クロージング
s = new_slide("21|CLOSING",
              "2027年度以降のウルトラマンの売場は、今日ここから始まります。映像の話題、イベントの熱量、"
              "流通の売場——そこに皆さまの商品とキャンペーンがつながることで、次の市場が生まれます。"
              "個別相談は、本日より受付開始です。会場を出る前に、ぜひQRのご登録を。皆さまと一緒に、"
              "次のウルトラマン市場をつくれることを楽しみにしています。")
title(s, "2027年度以降の売場は、今日ここから始まる", size=26)
rule(s)
add_text(s, ML, Inches(1.95), Inches(8.3), Inches(1.2),
         [("次のウルトラマン市場を、", 24, WHITE, True, {"line": 1.3, "after": 0}),
          ("皆さまと一緒につくりたい。", 24, WHITE, True, {"line": 1.3})])
qr2 = rect(s, Inches(9.3), Inches(1.9), Inches(3.45), Inches(3.45), PANEL2, RED,
           Pt(1.6), dash=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
tfq = qr2.text_frame; tfq.word_wrap = True
tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
for qi, qline in enumerate(["【QRコード再掲】", "個別商談・企画相談", "会場を出る前にご登録を"]):
    pq = tfq.paragraphs[0] if qi == 0 else tfq.add_paragraph()
    pq.alignment = PP_ALIGN.CENTER
    rq = pq.add_run(); rq.text = qline
    style_run(rq, 14 if qi != 2 else 10.5, WHITE if qi != 2 else BLUE, bold=(qi != 2))
add_text(s, ML, Inches(3.35), Inches(8.3), Inches(0.4),
         [("相談カテゴリ:商品化/食品・外食/流通・売場/広告/イベント・施設/IPコラボ/デジタル",
           11, SUB, False)])
photo_zone(s, ML, Inches(3.95), Inches(8.3), Inches(2.2),
           "クロージングビジュアル",
           "ヒーロー集合ビジュアル/商品・売場・イベント・広告のコラージュ/未来感のある背景")
bottom_copy(s, "個別相談は、本日より受付開始。")

OUT = "ULTRAMAN_BUSINESS_PARTNERS_MEETING_2026_draft_v5.pptx"
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides")
