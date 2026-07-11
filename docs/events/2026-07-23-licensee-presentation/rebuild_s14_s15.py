#!/usr/bin/env python3
"""S14・S15を改稿版に差し替える(7/11指示)。

S14: ULTRA MART → 「スタイルガイドを、企業ごとの提案プラットフォームへ」(商品化の質と提案力)
S15: スタイルガイド戦略 → 「商品化した後も、円谷が一緒に市場へ届ける」(販売・話題化まで伴走する支援力)

    python3 rebuild_s14_s15.py  # v9を読み、v10を出力
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

BG    = RGBColor(0x00, 0x07, 0x14)
WHITE = RGBColor(0xF4, 0xF7, 0xFD)
SUB   = RGBColor(0xA9, 0xB6, 0xD2)
DIM   = RGBColor(0x6A, 0x77, 0x94)
RED   = RGBColor(0xE8, 0x38, 0x2F)
REDDK = RGBColor(0xA8, 0x1C, 0x18)
PANEL = RGBColor(0x0C, 0x17, 0x30)
PANEL2= RGBColor(0x11, 0x1F, 0x3E)
LINE  = RGBColor(0x25, 0x36, 0x5E)
BLUE  = RGBColor(0x7F, 0xA8, 0xE0)
JP_FONT = "Yu Gothic"
ML = Inches(0.55)
CW = Inches(12.23)
TITLE_W = Inches(9.55)


def style_run(run, size, color=SUB, bold=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = JP_FONT
    f.language_id = MSO_LANGUAGE_ID.JAPANESE
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP_FONT)
    if spacing is not None:
        rPr.set("spc", str(spacing))


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for spec in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = spec if isinstance(spec, list) else [spec]
        opts = {}
        for rs in runs:
            text, size, color, bold = rs[0], rs[1], rs[2], rs[3]
            opts = rs[4] if len(rs) > 4 else opts
            run = p.add_run(); run.text = text
            style_run(run, size, color, bold, spacing=opts.get("spc"))
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(opts.get("before", 0))
        p.space_after = Pt(opts.get("after", 3))
        p.line_spacing = opts.get("line", 1.25)
    return box


def rect(slide, x, y, w, h, fill, line_color=None, line_w=None, dash=False,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = line_w or Pt(1)
        if dash:
            shp.line.dash_style = 4
    shp.shadow.inherit = False
    return shp


def header(slide, eyebrow, title_txt, lead_txt):
    add_text(slide, ML, Inches(0.32), TITLE_W, Inches(0.3),
             [(eyebrow, 11, RED, True, {"spc": 80})])
    add_text(slide, ML, Inches(0.62), TITLE_W, Inches(0.95),
             [(title_txt, 24, WHITE, True, {"line": 1.18})])
    rect(slide, ML, Inches(1.5), Inches(1.2), Pt(3.2), RED)
    add_text(slide, ML, Inches(1.7), CW, Inches(0.45),
             [(lead_txt, 14, BLUE, True)])


def bottom_copy(slide, text, y=6.42):
    band_h = Inches(0.56)
    rect(slide, ML, Inches(y), CW, band_h, PANEL)
    rect(slide, ML, Inches(y), Pt(4), band_h, RED)
    box = slide.shapes.add_textbox(ML + Inches(0.24), Inches(y), CW - Inches(0.48), band_h)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0
    r = tf.paragraphs[0].add_run(); r.text = text
    style_run(r, 14, WHITE, bold=True)


def clear_slide(slide):
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def chip_row(slide, x, y, w, label, items, accent=RED):
    """組み合わせ図の1行:左にカテゴリチップ、右に要素パネル."""
    ch, cw_ = Inches(0.50), Inches(1.55)
    rect(slide, x, y, cw_, ch, None, accent, Pt(1.4),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    add_text(slide, x, y, cw_, ch, [(label, 11, accent, True, {"align": PP_ALIGN.CENTER})],
             anchor=MSO_ANCHOR.MIDDLE)
    px = x + cw_ + Inches(0.18)
    rect(slide, px, y, w - cw_ - Inches(0.18), ch, PANEL2, LINE, Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    add_text(slide, px + Inches(0.2), y, w - cw_ - Inches(0.58), ch,
             [(items, 11.5, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================ S14
def build_s14(slide):
    clear_slide(slide)
    header(slide, "14|STYLE GUIDE PLATFORM",
           "スタイルガイドを、企業ごとの提案プラットフォームへ",
           "素材を選ぶのではなく、売れる表現を一緒に組み立てる。")

    LX, LW = ML, Inches(7.35)
    rows = [
        ("TARGET", "キッズ・親子・女性・若者・大人ファン・コレクター"),
        ("CATEGORY", "食品・文具・アパレル・雑貨・ホビー・流通・広告"),
        ("STYLE PARTS", "ヒーロー・怪獣・ポーズ・背景・パターン・ロゴ・カラー・コピー"),
    ]
    y = 2.30
    for i, (label, items) in enumerate(rows):
        chip_row(slide, LX, Inches(y), LW, label, items)
        if i < 2:
            add_text(slide, LX, Inches(y + 0.485), LW, Inches(0.20),
                     [("×", 11, RED, True, {"align": PP_ALIGN.CENTER, "after": 0, "line": 1.0})])
        y += 0.72
    add_text(slide, LX, Inches(y - 0.04), LW, Inches(0.20),
             [("▼", 11, RED, True, {"align": PP_ALIGN.CENTER, "after": 0, "line": 1.0})])
    ry = y + 0.22
    rect(slide, LX, Inches(ry), LW, Inches(0.70), REDDK,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_text(slide, LX + Inches(0.25), Inches(ry), LW - Inches(0.5), Inches(0.70),
             [[("企業別提案パッケージ", 14, WHITE, True),
               ("  商品モック / 売場 / キャンペーン / デジタル施策", 11, WHITE, False)]],
             anchor=MSO_ANCHOR.MIDDLE)

    # 現在 → 今後
    ny = ry + 0.92
    rect(slide, LX, Inches(ny), Inches(0.72), Inches(0.34), None, DIM, Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    add_text(slide, LX, Inches(ny), Inches(0.72), Inches(0.34),
             [("現在", 10, DIM, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, LX + Inches(0.88), Inches(ny), LW - Inches(0.88), Inches(0.34),
             [("完成アートを見せ、合うものを選んでもらう", 10.5, DIM, False)],
             anchor=MSO_ANCHOR.MIDDLE)
    ny += 0.44
    rect(slide, LX, Inches(ny), Inches(0.72), Inches(0.34), None, RED, Pt(1.2),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    add_text(slide, LX, Inches(ny), Inches(0.72), Inches(0.34),
             [("今後", 10, RED, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, LX + Inches(0.88), Inches(ny), LW - Inches(0.88), Inches(0.34),
             [("市場・顧客・価格・販路から逆算し、最適なアートパーツを組み合わせる", 10.5, WHITE, True)],
             anchor=MSO_ANCHOR.MIDDLE)

    # 右:同じボトルでも表現が変わる
    RX, RW = Inches(8.15), Inches(4.63)
    add_text(slide, RX, Inches(2.28), RW, Inches(0.34),
             [[("例|", 12, RED, True), ("同じ「ボトル・水筒」でも、ここまで変わる", 12, WHITE, True)]])
    ex = [
        ("幼児", "明るい・親しみ・デフォルメ", "水筒・ランチボックス・園児向け売場"),
        ("小学生", "アクション・強さ・メカ", "水筒・リュック・映画連動棚"),
        ("大人", "ミニマル・上質・シルエット", "ボトル・革ケース・限定箱・EC"),
    ]
    ey = 2.72
    for tgt, expr, prop in ex:
        rect(slide, RX, Inches(ey), RW, Inches(1.06), PANEL, LINE, Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(slide, RX + Inches(0.18), Inches(ey + 0.18), Inches(0.86), Inches(0.70), PANEL2,
             RED, Pt(1.2), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
        add_text(slide, RX + Inches(0.18), Inches(ey + 0.18), Inches(0.86), Inches(0.70),
                 [(tgt, 12, WHITE, True, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, RX + Inches(1.24), Inches(ey + 0.14), RW - Inches(1.44), Inches(0.82),
                 [[("表現  ", 9.5, BLUE, True), (expr, 10.5, WHITE, False)],
                  [("提案  ", 9.5, BLUE, True), (prop, 10.5, SUB, False)]])
        ey += 1.18
    bottom_copy(slide, "同じウルトラマンでも、顧客・商品・売場によって、最適な表現は変わる。")
    slide.notes_slide.notes_text_frame.text = (
        "現在も、多様なスタイルガイド、アートガイドをご用意しています。"
        "今後はさらに一歩進め、完成されたビジュアルを選んでいただくだけではなく、"
        "キャラクター、ポーズ、背景、パターン、カラー、ロゴなどを組み合わせ、"
        "企業ごとの顧客、商品、価格帯、販路に適した提案を一緒に作れる形へ進化させていきます。\n"
        "同じ水筒でも、幼児向け、小学生向け、大人向けでは、必要なウルトラマン表現は全く異なります。"
        "皆さまの顧客に最適なウルトラマンを、一緒に設計していきます。")


# ============================================================ S15
def build_s15(slide):
    clear_slide(slide)
    header(slide, "15|LAUNCH SUPPORT",
           "商品化した後も、円谷が一緒に市場へ届ける",
           "商品を発売するだけでなく、“売れる瞬間”を一緒に作る。")

    stages = [
        ("① 発売前", "期待を作る",
         ["ティザー投稿・先行告知", "商品開発ストーリー", "インフルエンサー先行体験",
          "予約・抽選受付", "ファン投票"]),
        ("② 発売時", "話題を最大化する",
         ["円谷公式SNS × ULTRA MART", "店頭大型モニター", "インフルエンサー投稿・ライブ配信",
          "店頭イベント・ヒーロー登場", "限定・先行販売"]),
        ("③ 発売後", "拡散と購買を伸ばす",
         ["購入者投稿・UGC活用", "EC誘導・再入荷告知", "使用シーン投稿",
          "二次展開", "既存流通への波及"]),
    ]
    SX, SW, GAP = ML, Inches(2.56), Inches(0.17)
    for i, (num, goal, items) in enumerate(stages):
        x = SX + Emu(int((SW + GAP) * i))
        rect(slide, x, Inches(2.30), SW, Inches(3.72), PANEL, LINE, Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
        rect(slide, x, Inches(2.30), SW, Pt(3), RED)
        add_text(slide, x + Inches(0.22), Inches(2.52), SW - Inches(0.44), Inches(0.80),
                 [(num, 14, WHITE, True, {"after": 1}),
                  (goal, 11.5, RED, True)])
        add_text(slide, x + Inches(0.22), Inches(3.42), SW - Inches(0.44), Inches(2.50),
                 [("・" + it, 10, SUB, False, {"line": 1.3, "after": 4}) for it in items])

    # 右:円谷の支援資産
    RX, RW = Inches(8.85), Inches(3.93)
    rect(slide, RX, Inches(2.30), RW, Inches(1.72), PANEL2, LINE, Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    add_text(slide, RX + Inches(0.22), Inches(2.46), RW - Inches(0.44), Inches(1.44),
             [("円谷が持つ発信・販売接点", 12, WHITE, True, {"after": 5}),
              ("公式SNS / ULTRA MART / ウルトラマンストア / ウルサマ・イベント / EC / POPUP / インフルエンサー / メディアPR / 店頭施策", 10, SUB, False, {"line": 1.4})])
    rect(slide, RX, Inches(4.18), RW, Inches(1.84), PANEL, RED, Pt(1.2),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    add_text(slide, RX + Inches(0.22), Inches(4.34), RW - Inches(0.44), Inches(1.56),
             [[("ULTRA MART = ", 12, WHITE, True), ("発信・実験拠点", 12, RED, True)],
              ("① 女性・若年・ライト層の開拓", 10, SUB, False, {"line": 1.35, "before": 4}),
              ("② 限定・先行販売による話題化", 10, SUB, False, {"line": 1.35}),
              ("③ 商品・デザインのテストマーケ", 10, SUB, False, {"line": 1.35}),
              ("④ 量販店・専門店・ECへの波及", 10, SUB, False, {"line": 1.35})])

    bottom_copy(slide, "原宿で話題を作り、全国の売場へ広げる。")
    slide.notes_slide.notes_text_frame.text = (
        "今後、円谷が強化していくもう一つの価値が、商品化後のマーケティング支援です。\n"
        "商品を作っていただき、監修して終わりではありません。"
        "円谷公式SNS、ULTRA MART、イベント、インフルエンサー、EC、店頭施策を組み合わせながら、"
        "発売前の期待形成、発売時の話題化、発売後の購買促進まで一緒に設計していきます。\n"
        "ULTRA MARTも、すべての商品を大量に販売するための場所ではなく、"
        "新しい商品、新しいターゲット、新しいデザインを試し、原宿で話題を作り、"
        "全国の売場へ広げていく発信拠点として活用していきます。\n"
        "商品化だけでなく、ファンと出会い、話題になり、売れるところまで、皆さまと一緒に作っていきます。")


if __name__ == "__main__":
    prs = Presentation("ULTRAMAN_BUSINESS_PARTNERS_MEETING_2026_v9.pptx")
    build_s14(prs.slides[13])
    build_s15(prs.slides[14])
    prs.save("ULTRAMAN_BUSINESS_PARTNERS_MEETING_2026_v10.pptx")
    print("saved v10 (S14/S15 rebuilt)")
