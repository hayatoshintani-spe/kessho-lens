# -*- coding: utf-8 -*-
# 経営会議_40億定点資料 — J.P.モルガン(HD向け資料)デザインルール準拠版
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def _set_dash(line_format, val="dash"):
    ln = line_format._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {})
        ln.append(d)
    d.set('val', val)

FONT = "ＭＳ Ｐゴシック"

CHARCOAL = RGBColor(0x3F, 0x3F, 0x3F)
DARKINK  = RGBColor(0x26, 0x26, 0x26)
GRAY     = RGBColor(0x7F, 0x7F, 0x7F)
LINEGR   = RGBColor(0xBF, 0xBF, 0xBF)
BROWN    = RGBColor(0x7B, 0x62, 0x44)
NAVY     = RGBColor(0x1F, 0x38, 0x64)
TEAL     = RGBColor(0x31, 0x85, 0x9B)
ORANGE   = RGBColor(0xC5, 0x5A, 0x11)
RED      = RGBColor(0xC0, 0x00, 0x00)
BEIGE    = RGBColor(0xE9, 0xE1, 0xD3)
BEIGE_L  = RGBColor(0xF2, 0xEC, 0xE1)
LTEAL    = RGBColor(0xDA, 0xEE, 0xF3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH


def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, bold=False, color=CHARCOAL, first=False, align=None,
         space_after=None, spc=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)
    if spc:
        rPr.set('spc', str(spc))
    return p


def bullet(tf, text, size, bcolor, color=CHARCOAL, first=False, space_after=3, indent_text=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    for t, c, b in [("●  ", bcolor, False), (text, color, False)]:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size if t != "●  " else size - 2.5)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = FONT
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
        ea.set('typeface', FONT)
    return p


def rect(s, x, y, w, h, fill, line=None, line_w=0.75):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def dashbox(s, x, y, w, h, color=RED, weight=1.4):
    sh = rect(s, x, y, w, h, None, line=color, line_w=weight)
    _set_dash(sh.line)
    return sh


def rule(s, x, y, w, color=LINEGR, weight=0.75):
    rect(s, x, y, w, Pt(weight), color)


def chart_title(s, x, y, w, text):
    """JPM流：小さめ太字ラベル＋直下に細罫線"""
    _, tf = box(s, x, y, w, Inches(0.26))
    para(tf, text, 10.5, bold=True, color=DARKINK, first=True)
    rule(s, x, y + Inches(0.28), w, LINEGR, 1.0)


def header(s, sec_no, sec_title, sec_color, title, message, chip_w=2.3):
    # セクションチップ
    rect(s, Inches(0.55), Inches(0.30), Inches(chip_w), Inches(0.32), sec_color)
    _, tf = box(s, Inches(0.55), Inches(0.345), Inches(chip_w), Inches(0.26))
    para(tf, f"{sec_no}. {sec_title}", 10.5, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
    # 社外秘（右上）
    _, tf = box(s, Inches(11.3), Inches(0.34), Inches(1.5), Inches(0.24))
    para(tf, "社 外 秘", 8, color=CHARCOAL, first=True, align=PP_ALIGN.RIGHT, spc=300)
    # ブラウンのアクセントバー
    rect(s, Inches(0.85), Inches(0.80), Inches(1.35), Inches(0.055), BROWN)
    # 表題
    _, tf = box(s, Inches(0.83), Inches(0.96), Inches(11.7), Inches(0.5))
    para(tf, title, 21, bold=False, color=DARKINK, first=True)
    # メッセージライン
    _, tf = box(s, Inches(0.85), Inches(1.52), Inches(11.9), Inches(0.55))
    para(tf, message, 12.5, bold=False, color=CHARCOAL, first=True)


PAGE_NOTE = "出所：部門別通期見込サマリー（2026年8月13日時点・財務検証済み）"


def footer(s, n, note=PAGE_NOTE):
    _, tf = box(s, Inches(0.55), Inches(7.06), Inches(2.6), Inches(0.3))
    para(tf, "円谷プロダクション STO", 9.5, bold=True, color=DARKINK, first=True)
    _, tf = box(s, Inches(3.25), Inches(7.10), Inches(5.9), Inches(0.3))
    para(tf, note, 7, color=GRAY, first=True)
    _, tf = box(s, Inches(9.2), Inches(7.09), Inches(3.6), Inches(0.3))
    para(tf, f"TSUBURAYA PRODUCTIONS   {n}", 8.5, color=CHARCOAL, first=True,
         align=PP_ALIGN.RIGHT, spc=150)


# ============================================================
# 表紙
# ============================================================
s = slide()
_, tf = box(s, Inches(0.85), Inches(0.50), Inches(4.5), Inches(0.6))
para(tf, "営業改革推進室（STO）", 13, bold=True, color=CHARCOAL, first=True)
_, tf = box(s, Inches(8.3), Inches(0.52), Inches(4.2), Inches(0.6))
para(tf, "TSUBURAYA PRODUCTIONS", 11, bold=True, color=GRAY, first=True,
     align=PP_ALIGN.RIGHT, spc=200)

_, tf = box(s, Inches(0.85), Inches(2.55), Inches(11.5), Inches(0.8))
para(tf, "業績定点報告 — 営業利益40億の達成管理", 32, bold=False, color=DARKINK, first=True)
_, tf = box(s, Inches(0.88), Inches(3.55), Inches(11.5), Inches(0.5))
para(tf, "円谷プロダクション 経営会議 | 2026年8月（8月13日正本ベース）", 15, color=CHARCOAL, first=True)

rect(s, Inches(0), Inches(6.90), SW, Inches(0.085), BROWN)
_, tf = box(s, Inches(0.85), Inches(7.08), Inches(3.0), Inches(0.3))
para(tf, "対 外 厳 秘", 8, color=CHARCOAL, first=True, spc=300)

# ============================================================
# 1. 今月のサマリー
# ============================================================
s = slide()
header(s, 1, "今月の全体像", NAVY, "営業利益40億 達成管理 — 今月のサマリー",
       "毎月同じ7枚のみを更新して報告する。見る数字は3つ：通期見込／40億までの残り／確定（黒）の額", chip_w=2.1)

tiles = [
    (0.55, "通期見込（PL・販管費込み）", "19.8億", "前期5.7億から＋14.2億。部門目標23.7億比▲3.9億", DARKINK, False),
    (4.70, "40億まで", "あと20.2億", "挽回3.9＋アップサイド15＋コスト1.3で到達", RED, True),
    (8.85, "アップサイドの確定（黒）", "0.3億", "カード9弾（8/17受注確定）。青・赤の転換が毎月の焦点", DARKINK, False),
]
for x, k, v, note, vc, emph in tiles:
    rect(s, Inches(x), Inches(2.15), Inches(3.94), Inches(1.55), WHITE, line=LINEGR)
    rect(s, Inches(x), Inches(2.15), Inches(3.94), Inches(0.36), BEIGE)
    _, tf = box(s, Inches(x + 0.18), Inches(2.21), Inches(3.6), Inches(0.28))
    para(tf, k, 10.5, bold=True, color=CHARCOAL, first=True)
    _, tf = box(s, Inches(x + 0.20), Inches(2.62), Inches(3.6), Inches(0.62))
    para(tf, v, 27, bold=True, color=vc, first=True)
    _, tf = box(s, Inches(x + 0.20), Inches(3.28), Inches(3.6), Inches(0.4))
    para(tf, note, 8.5, color=GRAY, first=True)
    if emph:
        dashbox(s, Inches(x - 0.07), Inches(2.08), Inches(4.08), Inches(1.69))

chart_title(s, Inches(0.55), Inches(4.02), Inches(5.5), "今月のポイント")
_, tf = box(s, Inches(0.62), Inches(4.42), Inches(12.1), Inches(1.3))
pts = [
    ("全社共通のPL正本が初めて確定（8/13・財務検証済み）。以後この数字だけで議論する", NAVY),
    ("目標比▲3.9億は3カ所に集中：ライブ▲1.8億／ライセンス（海外ディス・国内MD）▲2.2億／管理販管費＋1.5億。挽回オーナーを本日確定する", NAVY),
    ("構造リスク：中国単体の利益22.4億が全社19.8億を上回り、中国以外は▲2.5億。中国の契約防衛を常設の確認事項とする", RED),
]
for i, (p_, bc) in enumerate(pts):
    bullet(tf, p_, 11, bc, first=(i == 0), space_after=6)

rect(s, Inches(0.55), Inches(5.80), Inches(12.23), Inches(1.02), BEIGE_L)
rect(s, Inches(0.55), Inches(5.80), Inches(0.06), Inches(1.02), BROWN)
_, tf = box(s, Inches(0.85), Inches(5.93), Inches(11.6), Inches(0.8))
para(tf, "本日決めていただきたいこと（3件）", 11.5, bold=True, color=DARKINK, first=True, space_after=4)
para(tf, "①挽回3本（ライブ・海外ディス／国内MD・管理販管費）のオーナーと期限　②コスト削減の目標額　③トップ外交のターゲット",
     11, color=CHARCOAL, space_after=0)
footer(s, 2)

# ============================================================
# 2. ブリッジ
# ============================================================
s = slide()
header(s, 1, "今月の全体像", NAVY, "40億まで、あと20.2億 — 現在地からの道筋",
       "通期見込19.8億から4つのブロックで積み上げる。このページの形は毎月変えない", chip_w=2.1)

base_y = 6.15
scale = 0.095
bar_w = 1.62
xs = [0.80, 3.05, 5.30, 7.55, 9.80]

def wbar(x, bottom_val, height_val, fill, label_top, label_in=None, tcol=DARKINK):
    h = height_val * scale
    yy = base_y - (bottom_val + height_val) * scale
    rect(s, Inches(x), Inches(yy), Inches(bar_w), Inches(h), fill)
    _, tf = box(s, Inches(x - 0.25), Inches(yy - 0.31), Inches(bar_w + 0.5), Inches(0.3))
    para(tf, label_top, 13, bold=True, color=tcol, first=True, align=PP_ALIGN.CENTER)
    if label_in and h > 0.3:
        _, tf = box(s, Inches(x - 0.06), Inches(yy + h / 2 - 0.14), Inches(bar_w + 0.12), Inches(0.3))
        para(tf, label_in, 10.5, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)

wbar(xs[0], 0, 19.8, NAVY, "19.8億", "現在地")
wbar(xs[1], 19.8, 3.9, RED, "＋3.9億", "ベース挽回", tcol=RED)
wbar(xs[2], 23.7, 15.0, TEAL, "＋15億", "アップサイド", tcol=TEAL)
wbar(xs[3], 38.7, 1.3, ORANGE, "＋1.3億", tcol=ORANGE)
wbar(xs[4], 0, 40.0, NAVY, "40億", "今期目標")

cums = [19.8, 23.7, 38.7, 40.0]
for i in range(4):
    ln = rect(s, Inches(xs[i] + bar_w), Inches(base_y - cums[i] * scale),
              Inches(xs[i + 1] - xs[i] - bar_w), Pt(1.0), None, line=GRAY, line_w=0.9)
    _set_dash(ln.line)
rule(s, Inches(0.65), Inches(base_y), Inches(10.9), CHARCOAL, 1.2)

notes = [
    (xs[1], "毀損の復元", "ライブ・海外ディス・国内MD（5頁）"),
    (xs[2], "確度で管理", "黒0.3／青／赤（7頁）"),
    (xs[3], "コスト削減（第3のレバー）", "販管費超過＋1.0億の解消を含む"),
]
for x, t1, t2 in notes:
    _, tf = box(s, Inches(x - 0.35), Inches(6.30), Inches(2.35), Inches(0.6))
    para(tf, t1, 9.5, bold=True, color=NAVY, first=True, align=PP_ALIGN.CENTER, space_after=1)
    para(tf, t2, 7.5, color=GRAY, align=PP_ALIGN.CENTER, space_after=0)

chart_title(s, Inches(0.85), Inches(2.42), Inches(3.9), "現在地の評価")
_, tf = box(s, Inches(0.90), Inches(2.82), Inches(3.9), Inches(1.2))
para(tf, "前期5.7億 → 見込19.8億（＋14.2億の改善）。", 10, bold=True, color=CHARCOAL, first=True, space_after=2)
para(tf, "ただし部門目標23.7億比では▲3.9億", 10, bold=True, color=RED, space_after=4)
para(tf, "※カード9弾等は見込に未反映。毎月「あといくら」を更新する", 8.5, color=GRAY, space_after=0)
footer(s, 3)

# ============================================================
# 3. 利益の構造
# ============================================================
s = slide()
header(s, 2, "利益の構造と毀損", ORANGE, "利益の構造 — どこで稼ぎ、どこで使っているか",
       "単位：億円（営業利益・共通費配賦後）。この構造理解が全アクションの前提となる", chip_w=2.35)

rows_ = [
    ("営業事業本部", "42.1", "38.4", "▲3.7", "うち中国22.4億・国内ライセンス（粗利率86%）が利益の源泉", False),
    ("　うち中国", "22.7", "22.4", "▲0.3", "全社利益19.8億を単体で上回る。契約・取り分率の防衛が生命線", True),
    ("商品開発事業本部", "1.0", "0.7", "▲0.3", "カード9弾（8/17受注）が未反映。反転見込み", False),
    ("映像事業本部", "▲4.0", "▲3.8", "＋0.2", "構造赤字。製作投資としての位置づけを別途整理", False),
    ("管理・共通ほか", "▲15.5", "▲15.5", "▲0.1", "販管費超過＋1.0億の主因（総務・法務/人事/財務経理）", False),
    ("全社合計", "23.7", "19.8", "▲3.9", "粗利率56%・販管費26.0億（売上比32%）", False),
]
W = [2.9, 1.35, 1.35, 1.35, 5.28]
xpos = [0.55]
for w in W[:-1]:
    xpos.append(xpos[-1] + w)
heads = ["部門", "目標", "見込", "差", "読み方"]
cy = 2.15
x = 0.55
for h, w in zip(heads, W):
    rect(s, Inches(x), Inches(cy), Inches(w - 0.03), Inches(0.38), NAVY)
    _, tf = box(s, Inches(x), Inches(cy + 0.07), Inches(w - 0.03), Inches(0.28))
    para(tf, h, 10.5, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
    x += w
cy += 0.42
row_h = 0.58
for name, t, m, d, note, china in rows_:
    total = name == "全社合計"
    rect(s, Inches(0.55), Inches(cy), Inches(W[0] - 0.03), Inches(row_h - 0.04),
         BEIGE if not total else BEIGE)
    if total:
        rule(s, Inches(0.55), Inches(cy - 0.02), Inches(sum(W)), CHARCOAL, 1.4)
    _, tf = box(s, Inches(0.67), Inches(cy + 0.14), Inches(W[0] - 0.3), Inches(0.3))
    para(tf, name, 10.5, bold=(total or china), color=DARKINK, first=True)
    for (val, w, xx, bold, col, al) in [
        (t, W[1], xpos[1], total, CHARCOAL, PP_ALIGN.CENTER),
        (m, W[2], xpos[2], True, CHARCOAL, PP_ALIGN.CENTER),
        (d, W[3], xpos[3], True, (RED if "▲" in d else TEAL), PP_ALIGN.CENTER),
        (note, W[4], xpos[4], china, (NAVY if china else GRAY), None),
    ]:
        _, tf = box(s, Inches(xx + 0.06), Inches(cy + (0.14 if w != W[4] else 0.09)),
                    Inches(w - 0.16), Inches(0.44))
        para(tf, val, 10.5 if w != W[4] else 8.5, bold=bold, color=col, first=True, align=al)
    rule(s, Inches(0.55), Inches(cy + row_h - 0.02), Inches(sum(W)), LINEGR, 0.75)
    if china:
        dashbox(s, Inches(0.50), Inches(cy - 0.03), Inches(sum(W) + 0.10), Inches(row_h + 0.02))
    cy += row_h

rect(s, Inches(0.55), Inches(6.16), Inches(12.23), Inches(0.80), BEIGE_L)
rect(s, Inches(0.55), Inches(6.16), Inches(0.06), Inches(0.80), BROWN)
_, tf = box(s, Inches(0.85), Inches(6.27), Inches(11.6), Inches(0.62))
para(tf, "中国を除くと全社は▲2.5億。40億の議論は「中国が計画どおり続く」前提の上にある",
     12, bold=True, color=RED, first=True, space_after=2)
para(tf, "→ 実効取り分率・座組の変化を毎月この場で確認する（変化の兆候があれば月次を待たず報告）", 10, color=CHARCOAL, space_after=0)
footer(s, 4)

# ============================================================
# 4. ▲3.9億の中身
# ============================================================
s = slide()
header(s, 2, "利益の構造と毀損", ORANGE, "どこで漏れているか — 目標比▲3.9億の中身",
       "悪化は3カ所に集中。それぞれ挽回オーナーと期限を確定する（単位：百万円）", chip_w=2.35)

chart_title(s, Inches(0.55), Inches(2.10), Inches(6.5), "部門別の目標差（営業利益）")
_, tf = box(s, Inches(11.55), Inches(2.10), Inches(1.25), Inches(0.26))
para(tf, "オーナー", 9.5, bold=True, color=GRAY, first=True)

items = [
    ("ライセンス営業部", -220, "海外ディス▲90（売上見込がマイナス値・正体の特定が必要）／国内MD▲62／中国▲33", "南谷"),
    ("ライブイベント部", -183, "粗利▲185。原価付替＋有料ライブの後ズレ。時期差か実力減かの切り分けを", "南谷・菅野"),
    ("管理部門の販管費", -150, "総務＋101・法務/人事/財務経理＋106の超過。コスト削減チームの第1案件", "高橋"),
    ("商品開発（カード）", -29, "9弾受注（8/17確定）が未反映。次回更新で反転見込み", "山田"),
    ("プラス側", 60, "新規開拓＋23（海外＋33）・映像＋22・流通＋15。全部門が悪いわけではない", "－"),
]
cy = 2.58
maxv = 220
for name, v, note, owner in items:
    _, tf = box(s, Inches(0.55), Inches(cy + 0.03), Inches(2.5), Inches(0.3))
    para(tf, name, 11, bold=True, color=DARKINK, first=True)
    w = abs(v) / maxv * 3.0
    col = RED if v < 0 else TEAL
    rect(s, Inches(3.15), Inches(cy + 0.04), Inches(max(w, 0.15)), Inches(0.26), col)
    _, tf = box(s, Inches(3.15 + max(w, 0.15) + 0.1), Inches(cy + 0.02), Inches(1.0), Inches(0.3))
    para(tf, f"{v:+d}", 11, bold=True, color=(RED if v < 0 else TEAL), first=True)
    _, tf = box(s, Inches(7.35), Inches(cy + 0.02), Inches(4.15), Inches(0.56))
    para(tf, note, 8.5, color=CHARCOAL, first=True)
    _, tf = box(s, Inches(11.60), Inches(cy + 0.03), Inches(1.2), Inches(0.3))
    para(tf, owner, 9.5, bold=True, color=NAVY, first=True)
    rule(s, Inches(0.55), Inches(cy + 0.68), Inches(12.23), LINEGR, 0.5)
    cy += 0.76

rect(s, Inches(0.55), Inches(6.48), Inches(12.23), Inches(0.52), BEIGE_L)
rect(s, Inches(0.55), Inches(6.48), Inches(0.06), Inches(0.52), BROWN)
_, tf = box(s, Inches(0.85), Inches(6.59), Inches(11.6), Inches(0.36))
para(tf, "販管費超過＋1.0億の最大要因は事業部門ではなく管理部門。前回会長ご指摘の「コストの緩み」と符合する",
     11, bold=True, color=DARKINK, first=True)
footer(s, 5)

# ============================================================
# 5. アクションと効果額
# ============================================================
s = slide()
header(s, 3, "アクションと確度", DARKINK, "アクションと数字での効果 — 誰が・何を・いくら",
       "効果額は現時点の仮置き。本会議で各オーナーがコミット値に確定し、来月から実績で答え合わせする", chip_w=2.3)

acts = [
    ("挽回1", "ライブの後ズレ復元", "計上時期の確定と通期での戻り額の確定。戻らない分は代替積み上げ", "南谷・菅野", "＋1.0〜1.5億", RED),
    ("挽回2", "海外ディス・国内MDの立て直し", "売上マイナス値の正体特定（返品・解約・計上）→復元プラン", "南谷", "＋0.8〜1.2億", RED),
    ("挽回3", "管理販管費の削減", "総務・法務/人事/財務経理の超過分解→削減第1弾（コストチーム）", "高橋", "＋0.7〜1.0億", RED),
    ("転換1", "カード9弾の見込反映", "8/17受注確定分を月次更新に反映", "山田・柳沼", "＋0.3億（確定）", TEAL),
    ("転換2", "アップサイドの赤→青転換", "トップ外交・上位22社の四半期MTG・MG改定・UM/EC", "南谷・新谷", "15億の確度向上", TEAL),
]
cy = 2.12
for tag, title, note, owner, eff, tcol in acts:
    rect(s, Inches(0.55), Inches(cy), Inches(12.23), Inches(0.80), WHITE, line=LINEGR)
    rect(s, Inches(0.78), Inches(cy + 0.22), Inches(0.92), Inches(0.36), tcol)
    _, tf = box(s, Inches(0.78), Inches(cy + 0.27), Inches(0.92), Inches(0.28))
    para(tf, tag, 10, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
    _, tf = box(s, Inches(1.92), Inches(cy + 0.08), Inches(5.4), Inches(0.32))
    para(tf, title, 12, bold=True, color=DARKINK, first=True)
    _, tf = box(s, Inches(1.92), Inches(cy + 0.44), Inches(5.6), Inches(0.32))
    para(tf, note, 8.5, color=GRAY, first=True)
    _, tf = box(s, Inches(7.72), Inches(cy + 0.24), Inches(1.6), Inches(0.32))
    para(tf, owner, 10, bold=True, color=NAVY, first=True)
    _, tf = box(s, Inches(9.35), Inches(cy + 0.20), Inches(3.25), Inches(0.36))
    para(tf, eff, 13, bold=True, color=(RED if tcol == RED else TEAL), first=True, align=PP_ALIGN.RIGHT)
    cy += 0.88

rect(s, Inches(0.55), Inches(6.58), Inches(12.23), Inches(0.44), BEIGE_L)
rect(s, Inches(0.55), Inches(6.58), Inches(0.06), Inches(0.44), BROWN)
_, tf = box(s, Inches(0.85), Inches(6.66), Inches(11.6), Inches(0.32))
para(tf, "挽回・削減の合計 ＋2.8〜4.0億 ＝ ベース復元（＋3.9億）のめど。アップサイド15億は次頁の確度管理で積み上げる",
     11.5, bold=True, color=DARKINK, first=True)
footer(s, 6, note="注：効果額は仮置き。本会議で各オーナーがコミット値に確定")
# 「あと20.2億」構成の破線強調（効果額列）
dashbox(s, Inches(9.10), Inches(2.04), Inches(3.60), Inches(4.42), color=RED, weight=1.1)

# ============================================================
# 6. 確度パイプライン
# ============================================================
s = slide()
header(s, 3, "アクションと確度", DARKINK, "アップサイド15億の中身 — 確度別パイプライン",
       "黒（確定）・青（根拠あり）・赤（構想）で管理し、毎月「赤が青に、青が黒に変わったか」だけを確認する", chip_w=2.3)

rows = [
    ("ライセンス ＋5億（社内7億）", [
        ("青", "大型広告の新規・継続（4-7月実績で新規4社94百万の実勢）"),
        ("青", "上位22社（売上の3分の2）のアップセル・四半期戦略MTG"),
        ("赤", "トップ外交による大型案件（本日ターゲットをご相談）"),
    ]),
    ("商品 ＋5億（社内6億）", [
        ("黒", "カード9弾受注確定分（8/17確定・単弾で営業利益＋0.3億水準）"),
        ("青", "ウルトラマート（日販100万円ライン）・EC受注前年比2倍"),
        ("赤", "mofusand型の60周年横展開・フワヌイ追加発注"),
    ]),
    ("中国 ＋5億（社内7億）", [
        ("青", "MG改定・Blokees（該当部門より詳述）"),
        ("赤", "越境EC拡大（月商1,000万→3,000万の段階）・新規座組"),
        ("赤", "※実効取り分率の変動リスクを毎月監視（過去に原価率悪化の実例）"),
    ]),
]
cy = 2.12
cmap = {"黒": DARKINK, "青": NAVY, "赤": RED}
for title, items2 in rows:
    rect(s, Inches(0.55), Inches(cy), Inches(12.23), Inches(1.44), WHITE, line=LINEGR)
    rect(s, Inches(0.55), Inches(cy), Inches(12.23), Inches(0.36), BEIGE)
    _, tf = box(s, Inches(0.80), Inches(cy + 0.05), Inches(11.5), Inches(0.3))
    para(tf, title, 12, bold=True, color=DARKINK, first=True)
    iy = cy + 0.46
    for tag, txt in items2:
        rect(s, Inches(0.88), Inches(iy + 0.02), Inches(0.34), Inches(0.24), cmap[tag])
        _, tf = box(s, Inches(0.88), Inches(iy + 0.045), Inches(0.34), Inches(0.2))
        para(tf, tag, 8.5, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
        _, tf = box(s, Inches(1.36), Inches(iy + 0.015), Inches(11.2), Inches(0.28))
        para(tf, txt, 10, color=CHARCOAL, first=True)
        iy += 0.31
    cy += 1.58

_, tf = box(s, Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.26))
para(tf, "※ 見込に未反映：カード9弾／家賃▲26.4・退職▲18.4／出向＋40／海外ディス▲58（要確認）。毎月消し込む", 9, color=GRAY, first=True)
footer(s, 7)

# ============================================================
# 7. 月次トラッキング
# ============================================================
s = slide()
header(s, 4, "月次トラッキング", TEAL, "月次トラッキング — この表を毎月1列ずつ埋める",
       "単位：億円。「見込が上がる・黒が増える・販管費が超えない」の3つが揃えば40億に着地する", chip_w=2.35)

heads = ["", "8月\n(今回)", "9月", "10月", "11月", "12月\n(到達判定)", "1月", "2月", "3月\n(着地)"]
rows_t = [
    ("通期見込（営業利益）", ["19.8", "", "", "", "", "", "", ""]),
    ("40億まで（あといくら）", ["20.2", "", "", "", "", "", "", ""]),
    ("アップサイド 黒（確定）", ["0.3", "", "", "", "", "", "", ""]),
    ("アップサイド 青（根拠あり）", ["集計中", "", "", "", "", "", "", ""]),
    ("販管費 見込（計画25.0）", ["26.0", "", "", "", "", "", "", ""]),
]
x0 = 0.55
col_w = [3.3] + [1.11] * 8
cy = 2.22
# 8月列の地色
rect(s, Inches(x0 + col_w[0]), Inches(cy), Inches(col_w[1] - 0.03), Inches(0.52 + 5 * 0.54), LTEAL)
x = x0
for i, (h, w) in enumerate(zip(heads, col_w)):
    if i > 0:
        rect(s, Inches(x), Inches(cy), Inches(w - 0.03), Inches(0.52), NAVY if i != 1 else TEAL)
    _, tf = box(s, Inches(x), Inches(cy + 0.05), Inches(w - 0.03), Inches(0.46))
    for j, line in enumerate(h.split("\n")):
        para(tf, line, 9, bold=True, color=(WHITE if i > 0 else CHARCOAL), first=(j == 0),
             align=PP_ALIGN.CENTER, space_after=0)
    x += w
cy += 0.52
for name, vals in rows_t:
    rect(s, Inches(x0), Inches(cy), Inches(col_w[0] - 0.03), Inches(0.50), BEIGE)
    _, tf = box(s, Inches(x0 + 0.12), Inches(cy + 0.11), Inches(3.1), Inches(0.3))
    para(tf, name, 10, bold=True, color=DARKINK, first=True)
    x = x0 + col_w[0]
    for k, v in enumerate(vals):
        _, tf = box(s, Inches(x), Inches(cy + 0.11), Inches(col_w[1] - 0.06), Inches(0.3))
        para(tf, v if v else "－", 10.5, bold=bool(v), color=(DARKINK if v else LINEGR),
             first=True, align=PP_ALIGN.CENTER)
        x += col_w[1]
    rule(s, Inches(x0), Inches(cy + 0.50), Inches(sum(col_w)), LINEGR, 0.75)
    cy += 0.54

chart_title(s, Inches(0.55), Inches(5.70), Inches(4.0), "この資料の運用ルール")
_, tf = box(s, Inches(0.62), Inches(6.10), Inches(12.1), Inches(1.0))
rules = [
    "毎月、同じ7枚の数字だけを更新する（体裁・構成は変えない。作成はSTO・検証は財務）",
    "議論は「前月からの変化」と「コミット値との差異」のみ。差異には理由と対策を添える",
    "12月の点検会議で40億到達を判定。届かない場合は緊急レバー（条件改定前倒し・コスト第2弾）を発動",
]
for i, r_ in enumerate(rules):
    bullet(tf, r_, 10.5, TEAL, first=(i == 0), space_after=4)
footer(s, 8)

out = "/home/user/kessho-lens/docs/sales-reform-toolkit/meetings/経営会議_40億定点資料_202608.pptx"
prs.save(out)
print("OK:", out)
